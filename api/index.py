"""
OpenGov Auto Implementation Backend - Single File FastAPI Application
Complete backend for AI-powered implementation configuration from CSV data
"""

import base64
import csv
import io
import json
import os
from collections import Counter
from datetime import datetime
from typing import Any, Dict, List, Optional
from contextlib import asynccontextmanager

import uuid
from pydantic import BaseModel, Field
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware

# Try to import anthropic, fallback to mock if not available
try:
    from anthropic import Anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False


# ============================================================================
# PYDANTIC MODELS
# ============================================================================

class Condition(BaseModel):
    """Conditional logic rule: if field X has value Y, then do Z"""
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    source_field: str = ""
    operator: str = "equals"
    value: Any = None
    description: str = ""


class FormField(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    name: str
    field_type: str
    required: bool = True
    description: str = ""
    options: Optional[List[str]] = None
    conditions: List[Condition] = []


class WorkflowStep(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    name: str
    order: int
    assigned_role: str = ""
    status_from: Optional[str] = None
    status_to: str
    actions: List[str] = []
    conditions: List[Condition] = []


class Fee(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    name: str
    amount: float
    fee_type: str
    when_applied: str
    conditions: List[Condition] = []
    formula: str = ""


class RequiredDocument(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    name: str
    required: bool = True
    description: str = ""
    stage: str
    conditions: List[Condition] = []


class RecordType(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    name: str
    description: str = ""
    department_id: str = ""
    category: str = ""
    form_fields: List[FormField] = []
    workflow_steps: List[WorkflowStep] = []
    fees: List[Fee] = []
    required_documents: List[RequiredDocument] = []


class Department(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    name: str
    description: str = ""
    record_type_ids: List[str] = []


class UserRole(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    name: str
    description: str = ""
    permissions: List[str] = []
    department_ids: List[str] = []


class UploadedFile(BaseModel):
    filename: str
    size: int
    upload_time: datetime = Field(default_factory=datetime.utcnow)
    rows_count: int = 0
    columns: List[str] = []


class Configuration(BaseModel):
    record_types: List[RecordType] = []
    departments: List[Department] = []
    user_roles: List[UserRole] = []
    generated_at: Optional[datetime] = None
    summary: str = ""


class DataSource(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    source_type: str  # "municipal_code", "existing_form", "fee_schedule", "peer_template"
    name: str
    status: str = "pending"  # "pending", "processing", "completed", "error"
    url: Optional[str] = ""
    filename: Optional[str] = ""
    created_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    extracted_data: Optional[Dict[str, Any]] = None
    raw_text: Optional[str] = ""
    error_message: Optional[str] = ""

class ValidationFinding(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    severity: str  # "critical", "warning", "info", "success"
    category: str  # "completeness", "workflow", "fees", "documents", "best_practice"
    title: str
    description: str
    record_type_id: Optional[str] = ""
    recommendation: str = ""
    auto_fixable: bool = False
    fix_data: Optional[Dict[str, Any]] = None

class ReconciliationItem(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    action: str  # "add", "update", "flag"
    target: str  # "record_type", "fee", "form_field", "document", "workflow_step", "department"
    target_id: Optional[str] = ""
    record_type_name: Optional[str] = ""
    confidence: float = 0.0
    source_ids: List[str] = []
    title: str
    description: str
    suggested_data: Optional[Dict[str, Any]] = None
    status: str = "pending"  # "pending", "accepted", "rejected"


class Project(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4())[:8])
    name: str
    customer_name: str
    product_type: str = "PLC"
    status: str = "setup"
    created_at: datetime = Field(default_factory=datetime.utcnow)
    updated_at: datetime = Field(default_factory=datetime.utcnow)
    uploaded_files: List[UploadedFile] = []
    configuration: Optional[Configuration] = None
    analysis_progress: int = 0
    analysis_stage: str = ""
    community_url: Optional[str] = ""
    community_name: Optional[str] = ""
    community_research: Optional[str] = ""
    data_sources: List[Dict[str, Any]] = []
    validation_findings: List[Dict[str, Any]] = []
    reconciliation_items: List[Dict[str, Any]] = []
    intelligence_report: Optional[str] = ""


# Request models
class CreateProjectRequest(BaseModel):
    name: str
    customer_name: str
    product_type: str = "PLC"
    community_url: Optional[str] = ""
    id: Optional[str] = ""  # Allow re-creating with same ID (for cold start recovery)


class UpdateRecordTypeRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    department_id: Optional[str] = None
    category: Optional[str] = None
    form_fields: Optional[List[FormField]] = None
    workflow_steps: Optional[List[WorkflowStep]] = None
    fees: Optional[List[Fee]] = None
    required_documents: Optional[List[RequiredDocument]] = None


class UpdateDepartmentRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None


class UpdateRoleRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    permissions: Optional[List[str]] = None
    department_ids: Optional[List[str]] = None


# ============================================================================
# AI USAGE TRACKING
# ============================================================================

class AIUsageTracker:
    """Tracks Claude API calls and token usage"""
    def __init__(self):
        self.total_calls = 0
        self.total_tokens = 0
        self.call_history = []
        self.success_count = 0
        self.error_count = 0

    def record_call(self, operation: str, tokens_used: int, success: bool = True):
        self.total_calls += 1
        self.total_tokens += tokens_used
        if success:
            self.success_count += 1
        else:
            self.error_count += 1
        self.call_history.append({
            "operation": operation,
            "tokens": tokens_used,
            "timestamp": datetime.utcnow().isoformat(),
            "success": success
        })

    def get_stats(self) -> dict:
        return {
            "total_calls": self.total_calls,
            "total_tokens": self.total_tokens,
            "success_count": self.success_count,
            "error_count": self.error_count,
            "success_rate": round((self.success_count / self.total_calls * 100), 1) if self.total_calls > 0 else 0
        }

ai_usage_tracker = AIUsageTracker()

# ============================================================================
# PERSISTENT STORE (Vercel KV with /tmp per-project fallback)
# ============================================================================

STORE_PATH = os.path.join("/tmp", "plc_store.json")
PROJECT_DIR = os.path.join("/tmp", "plc_projects")
os.makedirs(PROJECT_DIR, exist_ok=True)

# Try to use Vercel KV (Upstash Redis) for persistent storage
KV_AVAILABLE = False
_redis_client = None
try:
    from upstash_redis import Redis
    _kv_url = os.environ.get("KV_REST_API_URL") or os.environ.get("UPSTASH_REDIS_REST_URL", "")
    _kv_token = os.environ.get("KV_REST_API_TOKEN") or os.environ.get("UPSTASH_REDIS_REST_TOKEN", "")
    if _kv_url and _kv_token:
        _redis_client = Redis(url=_kv_url, token=_kv_token)
        KV_AVAILABLE = True
        print(f"[KV] Upstash Redis connected: {_kv_url[:40]}...")
    else:
        print("[KV] No Upstash credentials found, using /tmp per-project fallback")
except ImportError:
    print("[KV] upstash-redis package not installed, using /tmp per-project fallback")
except Exception as e:
    print(f"[KV] Failed to connect: {e}")


def _kv_get(key):
    """GET from Upstash Redis"""
    if not KV_AVAILABLE or not _redis_client:
        return None
    try:
        result = _redis_client.get(key)
        if result and isinstance(result, str):
            return json.loads(result)
        return result
    except Exception as e:
        print(f"[KV] GET error for {key}: {e}")
        return None


def _kv_set(key, value):
    """SET to Upstash Redis"""
    if not KV_AVAILABLE or not _redis_client:
        return False
    try:
        encoded = json.dumps(value, default=lambda o: o.isoformat() if isinstance(o, datetime) else str(o))
        _redis_client.set(key, encoded)
        return True
    except Exception as e:
        print(f"[KV] SET error for {key}: {e}")
        return False


def _kv_delete(key):
    """DELETE from Upstash Redis"""
    if not KV_AVAILABLE or not _redis_client:
        return False
    try:
        _redis_client.delete(key)
        return True
    except Exception as e:
        print(f"[KV] DEL error for {key}: {e}")
        return False


def _kv_keys(pattern="project:*"):
    """List keys from Upstash Redis"""
    if not KV_AVAILABLE or not _redis_client:
        return []
    try:
        keys = _redis_client.keys(pattern)
        return keys if keys else []
    except Exception as e:
        print(f"[KV] KEYS error for {pattern}: {e}")
        return []


# --- Per-project /tmp file helpers (works even without KV) ---
def _project_file_path(project_id: str) -> str:
    """Return /tmp/plc_projects/{project_id}.json"""
    return os.path.join(PROJECT_DIR, f"{project_id}.json")


def _save_project_file(project_id: str, data: dict):
    """Save a single project to its own /tmp file."""
    try:
        path = _project_file_path(project_id)
        encoded = json.dumps(data, default=lambda o: o.isoformat() if isinstance(o, datetime) else str(o))
        with open(path, "w") as f:
            f.write(encoded)
            f.flush()
            os.fsync(f.fileno())
    except Exception as e:
        print(f"[DISK] Error saving project file {project_id}: {e}")


def _load_project_file(project_id: str) -> dict:
    """Load a single project from its own /tmp file."""
    try:
        path = _project_file_path(project_id)
        if os.path.exists(path):
            with open(path, "r") as f:
                return json.loads(f.read())
    except Exception as e:
        print(f"[DISK] Error loading project file {project_id}: {e}")
    return None


def _list_project_files() -> list:
    """List all project IDs from /tmp files."""
    try:
        files = os.listdir(PROJECT_DIR)
        return [f.replace(".json", "") for f in files if f.endswith(".json")]
    except Exception:
        return []


class InMemoryStore:
    """Store with Vercel KV persistence and per-project /tmp file fallback."""

    def __init__(self):
        self._projects = {}
        self._load_from_disk()

    def _load_from_disk(self):
        """Load from /tmp as fallback cache."""
        # Load from legacy monolithic file
        try:
            if os.path.exists(STORE_PATH):
                with open(STORE_PATH, "r") as f:
                    data = json.load(f)
                self._projects.update(data.get("projects", {}))
        except Exception:
            pass
        # Also load from per-project files
        for pid in _list_project_files():
            if pid not in self._projects:
                data = _load_project_file(pid)
                if data:
                    self._projects[pid] = data

    def _save_to_disk(self):
        """Save to /tmp as local cache (both monolithic and per-project)."""
        try:
            def serialize(obj):
                if isinstance(obj, datetime):
                    return obj.isoformat()
                raise TypeError(f"Type {type(obj)} not serializable")
            with open(STORE_PATH, "w") as f:
                json.dump({"projects": self._projects}, f, default=serialize)
        except Exception:
            pass

    def _persist_project(self, project_id):
        """Save single project to KV + per-project /tmp file + monolithic file."""
        if project_id in self._projects:
            _kv_set(f"project:{project_id}", self._projects[project_id])
            # Always save per-project file as fallback
            _save_project_file(project_id, self._projects[project_id])
        self._save_to_disk()

    def _load_project_from_kv(self, project_id):
        """Try to load a project from KV if not in memory."""
        if not KV_AVAILABLE:
            # Fallback: try per-project /tmp file
            data = _load_project_file(project_id)
            if data:
                print(f"[STORE] Loaded project {project_id} from /tmp file (name={data.get('name', '?')})")
                self._projects[project_id] = data
                return True
            print(f"[STORE] KV not available and no /tmp file for project {project_id}")
            return False
        data = _kv_get(f"project:{project_id}")
        if data:
            print(f"[STORE] Loaded project {project_id} from KV (name={data.get('name', '?')})")
            self._projects[project_id] = data
            return True
        # Even with KV available, check /tmp file as last resort
        data = _load_project_file(project_id)
        if data:
            print(f"[STORE] Loaded project {project_id} from /tmp file fallback (name={data.get('name', '?')})")
            self._projects[project_id] = data
            return True
        print(f"[STORE] Project {project_id} not found in KV or /tmp")
        return False

    def create_project(self, project: Project) -> Project:
        project_dict = project.model_dump()
        for key in ["created_at", "updated_at"]:
            if isinstance(project_dict.get(key), datetime):
                project_dict[key] = project_dict[key].isoformat()
        self._projects[project.id] = project_dict
        self._persist_project(project.id)
        return project

    def get_project(self, project_id: str) -> Optional[Project]:
        if project_id not in self._projects:
            self._load_from_disk()
        if project_id not in self._projects:
            self._load_project_from_kv(project_id)
        if project_id not in self._projects:
            print(f"[STORE] Project {project_id} not found in memory, disk, or KV | KV_AVAILABLE={KV_AVAILABLE}")
            return None
        try:
            return Project(**self._projects[project_id])
        except Exception as e:
            # Don't silently lose projects — try to repair the data
            print(f"[STORE] Project {project_id} deserialization error: {e}")
            try:
                # Try to fix common issues: datetime strings, missing fields
                data = dict(self._projects[project_id])
                # Ensure uploaded_files is a list of proper dicts
                if "uploaded_files" in data and isinstance(data["uploaded_files"], list):
                    cleaned = []
                    for uf in data["uploaded_files"]:
                        if isinstance(uf, dict):
                            # Ensure upload_time is a string if present
                            if "upload_time" in uf and isinstance(uf["upload_time"], str):
                                try:
                                    uf["upload_time"] = datetime.fromisoformat(uf["upload_time"].replace("Z", "+00:00"))
                                except Exception:
                                    uf["upload_time"] = datetime.utcnow()
                            cleaned.append(uf)
                    data["uploaded_files"] = cleaned
                # Convert datetime strings
                for dt_field in ["created_at", "updated_at"]:
                    if dt_field in data and isinstance(data[dt_field], str):
                        try:
                            data[dt_field] = datetime.fromisoformat(data[dt_field].replace("Z", "+00:00"))
                        except Exception:
                            data[dt_field] = datetime.utcnow()
                return Project(**data)
            except Exception as e2:
                print(f"[STORE] Project {project_id} repair also failed: {e2}")
                # Last resort: return a minimal project with at least the ID and basic data
                try:
                    raw = self._projects[project_id]
                    return Project(
                        id=project_id,
                        name=raw.get("name", "Unknown"),
                        customer_name=raw.get("customer_name", "Unknown"),
                        status=raw.get("status", "setup"),
                        community_url=raw.get("community_url", ""),
                    )
                except Exception as e3:
                    print(f"[STORE] Project {project_id} minimal reconstruction failed: {e3}")
                    return None

    def list_projects(self) -> List[Project]:
        self._load_from_disk()
        # Also load from KV if available
        if KV_AVAILABLE:
            try:
                keys = _kv_keys("project:*")
                for key in keys:
                    pid = key.replace("project:", "")
                    if pid not in self._projects:
                        self._load_project_from_kv(pid)
            except Exception:
                pass
        projects = []
        for p in self._projects.values():
            try:
                projects.append(Project(**p))
            except Exception:
                continue
        return projects

    def update_project(self, project_id: str, **updates) -> Project:
        if project_id not in self._projects:
            self._load_from_disk()
        if project_id not in self._projects:
            self._load_project_from_kv(project_id)
        if project_id not in self._projects:
            raise ValueError(f"Project {project_id} not found")
        self._projects[project_id].update(updates)
        self._projects[project_id]["updated_at"] = datetime.utcnow().isoformat()
        self._persist_project(project_id)
        return Project(**self._projects[project_id])

    def delete_project(self, project_id: str) -> bool:
        if project_id in self._projects:
            del self._projects[project_id]
            _kv_delete(f"project:{project_id}")
            self._save_to_disk()
            return True
        return False

    def save_configuration(self, project_id: str, config: Configuration) -> None:
        if project_id not in self._projects:
            self._load_from_disk()
        if project_id not in self._projects:
            self._load_project_from_kv(project_id)
        if project_id not in self._projects:
            raise ValueError(f"Project {project_id} not found")
        config_dict = config.model_dump()
        self._projects[project_id]["configuration"] = config_dict
        self._projects[project_id]["updated_at"] = datetime.utcnow().isoformat()
        self._persist_project(project_id)

    def _ensure_project(self, project_id):
        """Ensure project is loaded from all sources."""
        if project_id not in self._projects:
            self._load_from_disk()
        if project_id not in self._projects:
            self._load_project_from_kv(project_id)

    def update_record_type(self, project_id: str, rt_id: str, updates: dict) -> Optional[RecordType]:
        self._ensure_project(project_id)
        if project_id not in self._projects:
            return None
        config = self._projects[project_id].get("configuration")
        if not config:
            return None
        for rt in config.get("record_types", []):
            if rt["id"] == rt_id:
                rt.update(updates)
                self._projects[project_id]["updated_at"] = datetime.utcnow().isoformat()
                self._persist_project(project_id)
                return RecordType(**rt)
        return None

    def update_department(self, project_id: str, dept_id: str, updates: dict) -> Optional[Department]:
        self._ensure_project(project_id)
        if project_id not in self._projects:
            return None
        config = self._projects[project_id].get("configuration")
        if not config:
            return None
        for dept in config.get("departments", []):
            if dept["id"] == dept_id:
                dept.update(updates)
                self._projects[project_id]["updated_at"] = datetime.utcnow().isoformat()
                self._persist_project(project_id)
                return Department(**dept)
        return None

    def update_role(self, project_id: str, role_id: str, updates: dict) -> Optional[UserRole]:
        self._ensure_project(project_id)
        if project_id not in self._projects:
            return None
        config = self._projects[project_id].get("configuration")
        if not config:
            return None
        for role in config.get("user_roles", []):
            if role["id"] == role_id:
                role.update(updates)
                self._projects[project_id]["updated_at"] = datetime.utcnow().isoformat()
                self._persist_project(project_id)
                return UserRole(**role)
        return None

    def add_record_type(self, project_id: str, record_type: RecordType) -> Optional[RecordType]:
        self._ensure_project(project_id)
        if project_id not in self._projects:
            return None
        config = self._projects[project_id].get("configuration")
        if not config:
            return None
        rt_dict = record_type.model_dump()
        config["record_types"].append(rt_dict)
        self._projects[project_id]["updated_at"] = datetime.utcnow().isoformat()
        self._persist_project(project_id)
        return RecordType(**rt_dict)

    def delete_record_type(self, project_id: str, rt_id: str) -> bool:
        self._ensure_project(project_id)
        if project_id not in self._projects:
            return False
        config = self._projects[project_id].get("configuration")
        if not config:
            return False
        for i, rt in enumerate(config.get("record_types", [])):
            if rt["id"] == rt_id:
                config["record_types"].pop(i)
                self._projects[project_id]["updated_at"] = datetime.utcnow().isoformat()
                self._persist_project(project_id)
                return True
        return False

    def add_uploaded_file(self, project_id: str, file_info: UploadedFile) -> None:
        self._ensure_project(project_id)
        if project_id not in self._projects:
            raise ValueError(f"Project {project_id} not found")
        file_dict = file_info.model_dump()
        for key in ["upload_time"]:
            if isinstance(file_dict.get(key), datetime):
                file_dict[key] = file_dict[key].isoformat()
        self._projects[project_id]["uploaded_files"].append(file_dict)
        self._projects[project_id]["updated_at"] = datetime.utcnow().isoformat()
        self._persist_project(project_id)


store = InMemoryStore()


# ============================================================================
# CSV PARSER
# ============================================================================

class CSVParser:
    @staticmethod
    def parse(content: str) -> dict:
        """Parse CSV and extract metadata for AI analysis"""
        reader = csv.DictReader(io.StringIO(content))
        rows = list(reader)
        columns = reader.fieldnames or []

        column_analysis = {}
        for col in columns:
            values = [r.get(col, "") for r in rows if r.get(col)]
            unique_count = len(set(values))
            sample_values = list(set(values))[:10]

            is_numeric = True
            if values:
                for v in values[:20]:
                    if v and not v.replace('.', '', 1).replace('-', '', 1).replace(',', '').isdigit():
                        is_numeric = False
                        break

            is_date = any(kw in col.lower() for kw in ["date", "time", "created", "submitted", "approved"])

            column_analysis[col] = {
                "unique_count": unique_count,
                "total_count": len(values),
                "sample_values": sample_values,
                "appears_numeric": is_numeric,
                "appears_date": is_date,
            }

        sample_rows = rows[:15]

        return {
            "columns": columns,
            "total_rows": len(rows),
            "sample_rows": sample_rows,
            "column_analysis": column_analysis,
        }

    @staticmethod
    def to_summary_string(metadata: dict) -> str:
        """Convert metadata to a string suitable for Claude prompt"""
        lines = []
        lines.append(f"Total rows: {metadata['total_rows']}")
        lines.append(f"Columns ({len(metadata['columns'])}): {', '.join(metadata['columns'])}")
        lines.append("")
        lines.append("Column Analysis:")
        for col, info in metadata['column_analysis'].items():
            sample = info['sample_values'][:5]
            lines.append(f"  {col}: {info['unique_count']} unique values, samples: {sample}")
        lines.append("")
        lines.append("Sample Rows (first 5):")
        for i, row in enumerate(metadata['sample_rows'][:5]):
            lines.append(f"  Row {i+1}: {dict(row)}")
        return "\n".join(lines)


# ============================================================================
# WEB RESEARCHER
# ============================================================================

class WebResearcher:
    """Researches local government websites to gather ordinances, fees, and processes."""

    def __init__(self):
        self.api_key = os.getenv("ANTHROPIC_API_KEY", "")

    def research_community(self, community_url: str, community_name: str = "") -> dict:
        """Research a community's government website for PLC configuration data."""
        return self._generate_mock_research(community_url, community_name)

    def _generate_mock_research(self, url: str, name: str) -> dict:
        """Generate realistic mock research data for demo purposes."""
        community = name or "the community"

        return {
            "community_name": community,
            "website_url": url,
            "research_summary": f"Comprehensive research of {community}'s local government website completed. Found detailed information about permit processes, fee schedules, municipal codes, and departmental structure.",
            "permits_found": [
                {"name": "Building Permit", "description": f"{community} requires building permits for new construction, additions, renovations over $5,000, and structural modifications. Applications reviewed by Planning & Building Dept.", "typical_timeline": "2-6 weeks"},
                {"name": "Business License", "description": f"All businesses operating within {community} city limits must obtain an annual business license. Home-based businesses included.", "typical_timeline": "5-10 business days"},
                {"name": "Encroachment Permit", "description": f"{community} Public Works requires encroachment permits for any work within the public right-of-way including sidewalks, curbs, and utilities.", "typical_timeline": "1-3 weeks"},
                {"name": "Sign Permit", "description": f"Required for all new signs, changes to existing signs, or temporary signs in {community}. Must comply with sign ordinance Chapter 17.40.", "typical_timeline": "1-2 weeks"},
                {"name": "Grading Permit", "description": f"Required for earth-moving activities over 50 cubic yards in {community}.", "typical_timeline": "2-4 weeks"},
                {"name": "Conditional Use Permit", "description": f"Required for uses not permitted by right in specific zoning districts per {community} Zoning Code.", "typical_timeline": "6-12 weeks (requires public hearing)"},
            ],
            "fee_schedule": [
                {"permit_type": "Building Permit", "fee_name": "Plan Check Fee", "amount": "65% of building permit fee", "notes": "Based on project valuation"},
                {"permit_type": "Building Permit", "fee_name": "Building Permit Fee", "amount": "Per CBC Table 1-A", "notes": "Based on project valuation using ICC valuation table"},
                {"permit_type": "Building Permit", "fee_name": "Technology Fee", "amount": "$30.00", "notes": "Flat fee per application"},
                {"permit_type": "Building Permit", "fee_name": "SMIP Fee", "amount": "$0.13 per $1,000 valuation", "notes": "Strong Motion Instrumentation Program"},
                {"permit_type": "Building Permit", "fee_name": "Green Building Fee", "amount": "$4.50 per $1,000 valuation", "notes": "CalGreen compliance review"},
                {"permit_type": "Business License", "fee_name": "Application Fee", "amount": "$125.00", "notes": "Non-refundable"},
                {"permit_type": "Business License", "fee_name": "Annual Renewal", "amount": "$75.00 - $500.00", "notes": "Based on number of employees"},
                {"permit_type": "Business License", "fee_name": "Home Occupation", "amount": "$50.00", "notes": "Annual fee for home-based businesses"},
                {"permit_type": "Encroachment", "fee_name": "Permit Fee", "amount": "$275.00", "notes": "Base fee"},
                {"permit_type": "Encroachment", "fee_name": "Inspection Deposit", "amount": "$1,500.00", "notes": "Refundable upon satisfactory completion"},
                {"permit_type": "Fire Prevention", "fee_name": "Fire Alarm Permit", "amount": "$250.00", "notes": "New installations"},
                {"permit_type": "Fire Prevention", "fee_name": "Sprinkler Plan Review", "amount": "$175.00", "notes": "Per plan set"},
            ],
            "departments": [
                {"name": "Community Development", "description": f"Oversees planning, building, and code enforcement for {community}. Manages building permits, plan reviews, and inspections.", "phone": "(555) 555-0100"},
                {"name": "Business License Division", "description": f"Part of the Finance Department. Processes all business licenses and renewals for {community}.", "phone": "(555) 555-0200"},
                {"name": "Public Works", "description": f"Manages {community}'s infrastructure including streets, sidewalks, storm drains, and rights-of-way.", "phone": "(555) 555-0300"},
                {"name": "Fire Prevention Bureau", "description": f"Part of {community} Fire Department. Handles fire prevention permits, inspections, and plan reviews.", "phone": "(555) 555-0400"},
                {"name": "Code Enforcement", "description": f"Ensures compliance with {community}'s municipal codes, property maintenance standards, and zoning regulations.", "phone": "(555) 555-0500"},
            ],
            "ordinances": [
                {"code": "Title 15 - Buildings and Construction", "summary": f"Adopts California Building Code with {community}-specific amendments. Covers building permits, plan reviews, inspections, and compliance.", "key_provisions": ["Permit required for work over $500", "Plans required for projects over $5,000", "Licensed contractor required for projects over $500"]},
                {"code": "Title 17 - Zoning", "summary": f"{community} Zoning Ordinance establishing land use districts, permitted uses, development standards, and approval processes.", "key_provisions": ["7 residential zones", "5 commercial zones", "3 industrial zones", "Overlay districts for historic and flood areas"]},
                {"code": "Title 5 - Business Licenses and Regulations", "summary": f"Requires all businesses within {community} to obtain a business license. Establishes fee schedule and renewal requirements.", "key_provisions": ["Annual renewal required", "Home occupation permits available", "Penalties for operating without license"]},
                {"code": "Title 8 - Health and Safety", "summary": f"{community}'s health and safety codes including fire prevention, hazardous materials, and property maintenance.", "key_provisions": ["Adopts California Fire Code", "Annual fire inspections for commercial", "Weed abatement program"]},
            ],
            "processes": [
                {"name": "Building Permit Process", "steps": ["Submit application with plans and fees", "Completeness check (3 business days)", "Plan review by multiple departments (2-4 weeks)", "Corrections cycle if needed", "Permit issuance upon approval", "Inspections during construction", "Final inspection and certificate of occupancy"]},
                {"name": "Business License Process", "steps": ["Complete application form", "Pay application fee", "Zoning verification", "Fire inspection (if applicable)", "License issued", "Annual renewal notice sent 30 days before expiration"]},
                {"name": "Code Enforcement Process", "steps": ["Complaint received or violation observed", "Case opened and assigned", "Initial inspection within 5 business days", "Notice of violation sent to property owner", "30-day compliance period", "Re-inspection", "Administrative citation if not corrected", "Hearing process for appeals"]},
            ],
            "documents_commonly_required": [
                "Completed application form",
                "Site plan or plot plan",
                "Architectural/construction plans (3 sets)",
                "Structural calculations (sealed by licensed engineer)",
                "Title 24 Energy compliance forms",
                "Soils/geotechnical report (for new construction)",
                "Proof of property ownership or authorization letter",
                "Licensed contractor information",
                "Proof of insurance",
                "Environmental review documentation (CEQA)",
                "School district fee receipt",
                "Water/sewer availability letter",
            ],
            "sources_reviewed": [
                {
                    "category": "Municipal Code & Ordinances",
                    "sources": [
                        {"title": "Municipal Code — Title 15: Buildings and Construction", "url": f"{url}/municipal-code/title-15", "type": "municipal_code", "description": "Full text of building code adoptions, local amendments, permit requirements, and enforcement provisions. Includes CBC amendments specific to the jurisdiction."},
                        {"title": "Municipal Code — Title 17: Zoning", "url": f"{url}/municipal-code/title-17", "type": "municipal_code", "description": "Zoning districts, permitted/conditional uses, development standards, setbacks, height limits, parking requirements, and overlay zones."},
                        {"title": "Municipal Code — Title 5: Business Licenses & Regulations", "url": f"{url}/municipal-code/title-5", "type": "municipal_code", "description": "Business license requirements, fee schedules, home occupation permits, solicitor permits, and regulatory compliance."},
                        {"title": "Municipal Code — Title 8: Health and Safety", "url": f"{url}/municipal-code/title-8", "type": "municipal_code", "description": "Fire prevention codes, hazardous materials storage, property maintenance standards, and weed abatement programs."},
                        {"title": "Municipal Code — Title 12: Streets, Sidewalks and Public Places", "url": f"{url}/municipal-code/title-12", "type": "municipal_code", "description": "Encroachment permits, right-of-way regulations, sidewalk repair requirements, and public infrastructure standards."},
                        {"title": "Municipal Code — Title 16: Subdivisions", "url": f"{url}/municipal-code/title-16", "type": "municipal_code", "description": "Subdivision map requirements, lot line adjustments, parcel mergers, and development agreement provisions."},
                    ]
                },
                {
                    "category": "Permit Applications & Forms",
                    "sources": [
                        {"title": "Building Permit Application (Form B-100)", "url": f"{url}/departments/community-development/forms/building-permit-application", "type": "application_form", "description": "Standard building permit application form with project details, contractor information, and valuation worksheet."},
                        {"title": "Business License Application (Form BL-01)", "url": f"{url}/departments/finance/forms/business-license-application", "type": "application_form", "description": "New business license application including business classification, owner information, and zoning verification."},
                        {"title": "Encroachment Permit Application (Form PW-200)", "url": f"{url}/departments/public-works/forms/encroachment-permit", "type": "application_form", "description": "Right-of-way work permit with traffic control plan requirements, insurance certificates, and restoration deposit."},
                        {"title": "Sign Permit Application (Form P-300)", "url": f"{url}/departments/planning/forms/sign-permit", "type": "application_form", "description": "Sign installation permit with dimensions, illumination details, and sign code compliance checklist."},
                        {"title": "Conditional Use Permit Application (Form P-400)", "url": f"{url}/departments/planning/forms/cup-application", "type": "application_form", "description": "Conditional use permit with project narrative, site plan requirements, environmental review checklist, and public hearing notice."},
                        {"title": "Grading Permit Application (Form E-100)", "url": f"{url}/departments/engineering/forms/grading-permit", "type": "application_form", "description": "Grading and earthwork permit with volume calculations, erosion control plan, and SWPPP requirements."},
                        {"title": "Demolition Permit Application (Form B-150)", "url": f"{url}/departments/community-development/forms/demolition-permit", "type": "application_form", "description": "Demolition permit with asbestos survey requirements, utility disconnect confirmation, and site restoration plan."},
                    ]
                },
                {
                    "category": "Fee Schedules & Rate Tables",
                    "sources": [
                        {"title": "Master Fee Schedule — Community Development", "url": f"{url}/departments/community-development/fees", "type": "fee_schedule", "description": "Comprehensive fee schedule for building permits, plan review, inspections, re-inspections, and after-hours charges. Based on ICC valuation tables."},
                        {"title": "Master Fee Schedule — Planning Division", "url": f"{url}/departments/planning/fees", "type": "fee_schedule", "description": "Planning application fees including CUPs, variances, zone changes, general plan amendments, environmental review, and public hearing deposits."},
                        {"title": "Master Fee Schedule — Business License", "url": f"{url}/departments/finance/business-license-fees", "type": "fee_schedule", "description": "Business license fee tiers based on employee count, business type, and gross receipts. Includes home occupation and temporary event fees."},
                        {"title": "Master Fee Schedule — Public Works & Engineering", "url": f"{url}/departments/public-works/fees", "type": "fee_schedule", "description": "Encroachment permits, utility connection fees, traffic control plan review, improvement plan check, and inspection deposits."},
                        {"title": "Master Fee Schedule — Fire Prevention", "url": f"{url}/departments/fire/prevention-fees", "type": "fee_schedule", "description": "Fire alarm permits, sprinkler system review, hazardous materials permits, special event fire safety, and annual inspection fees."},
                        {"title": "Impact Fee Schedule (Resolution 2024-45)", "url": f"{url}/departments/community-development/impact-fees", "type": "fee_schedule", "description": "Development impact fees for parks, traffic, drainage, schools, and public facilities. Updated annually per CPI."},
                    ]
                },
                {
                    "category": "Department Pages & Staff Directories",
                    "sources": [
                        {"title": "Community Development Department", "url": f"{url}/departments/community-development", "type": "department_page", "description": "Main department page with building, planning, and code enforcement divisions. Includes staff directory, office hours, and counter service information."},
                        {"title": "Building Division — Inspection Scheduling", "url": f"{url}/departments/community-development/building/inspections", "type": "department_page", "description": "Online inspection request system, inspection types, scheduling requirements, and inspector contact information."},
                        {"title": "Planning Division — Current Projects", "url": f"{url}/departments/planning/current-projects", "type": "department_page", "description": "Active planning applications, public hearing calendar, environmental reviews in progress, and project status tracker."},
                        {"title": "Code Enforcement Division", "url": f"{url}/departments/code-enforcement", "type": "department_page", "description": "Code complaint submission, common violations, compliance timelines, citation and fine schedule, and appeal process."},
                        {"title": "Finance Department — Business Licensing", "url": f"{url}/departments/finance/business-license", "type": "department_page", "description": "Business license portal, renewal process, business classification guide, and home occupation guidelines."},
                        {"title": "Public Works Department", "url": f"{url}/departments/public-works", "type": "department_page", "description": "Infrastructure projects, encroachment permits, utility services, street maintenance, and capital improvement plan."},
                        {"title": "Fire Prevention Bureau", "url": f"{url}/departments/fire/prevention", "type": "department_page", "description": "Fire permits, annual inspections, plan review process, fire sprinkler requirements, and hazardous materials program."},
                    ]
                },
                {
                    "category": "Plans, Policies & Studies",
                    "sources": [
                        {"title": "General Plan 2040", "url": f"{url}/departments/planning/general-plan", "type": "policy_document", "description": "Comprehensive long-range plan covering land use, circulation, housing, conservation, open space, safety, and noise elements."},
                        {"title": "Housing Element (2023-2031 Cycle)", "url": f"{url}/departments/planning/housing-element", "type": "policy_document", "description": "RHNA allocation, site inventory, housing programs, affirmatively furthering fair housing analysis, and ADU provisions."},
                        {"title": "Climate Action Plan", "url": f"{url}/departments/community-development/climate-action", "type": "policy_document", "description": "GHG reduction targets, green building requirements, EV charging mandates, water conservation standards, and sustainability metrics."},
                        {"title": "Capital Improvement Program (CIP) FY 2024-2029", "url": f"{url}/departments/public-works/cip", "type": "policy_document", "description": "Five-year infrastructure investment plan covering streets, water, sewer, parks, and public facilities with funding sources and timelines."},
                        {"title": "Specific Plan — Downtown Revitalization", "url": f"{url}/departments/planning/downtown-specific-plan", "type": "policy_document", "description": "Design standards, permitted uses, parking reductions, mixed-use incentives, and streamlined review process for downtown area."},
                    ]
                },
                {
                    "category": "Public Meeting Records",
                    "sources": [
                        {"title": "City Council Meeting Agendas & Minutes", "url": f"{url}/city-council/agendas-minutes", "type": "meeting_record", "description": "Regular and special meeting records including ordinance adoptions, fee schedule updates, development project approvals, and policy changes."},
                        {"title": "Planning Commission Agendas & Minutes", "url": f"{url}/boards-commissions/planning-commission", "type": "meeting_record", "description": "Public hearing records for CUPs, variances, subdivisions, environmental reviews, and general plan amendments."},
                        {"title": "Building Board of Appeals", "url": f"{url}/boards-commissions/building-board-of-appeals", "type": "meeting_record", "description": "Appeal decisions on building code interpretations, alternative materials requests, and permit denials."},
                    ]
                },
                {
                    "category": "Online Services & Portals",
                    "sources": [
                        {"title": "Online Permit Portal", "url": f"{url}/services/permit-portal", "type": "online_service", "description": "Electronic permit application submission, document upload, fee payment, permit status tracking, and inspection scheduling."},
                        {"title": "GIS / Interactive Zoning Map", "url": f"{url}/services/gis-maps", "type": "online_service", "description": "Parcel lookup, zoning designation, flood zone determination, assessor data, aerial imagery, and utility infrastructure layers."},
                        {"title": "Document Center / Records Request", "url": f"{url}/services/document-center", "type": "online_service", "description": "Public records archive including resolutions, ordinances, contracts, EIRs, inspection reports, and historical permits."},
                        {"title": "Code Enforcement Complaint Portal", "url": f"{url}/services/report-violation", "type": "online_service", "description": "Online code violation reporting with photo upload, anonymous option, case tracking, and status notifications."},
                    ]
                },
            ],
            "research_depth": {
                "pages_analyzed": 47,
                "documents_reviewed": 23,
                "forms_cataloged": 7,
                "fee_tables_extracted": 6,
                "last_updated": datetime.now().isoformat(),
            }
        }

    def format_for_analysis(self, research: dict) -> str:
        """Format research data as context string for Claude analysis prompt."""
        parts = []

        parts.append(f"## Community Research: {research.get('community_name', 'Unknown')}")
        parts.append(f"Website: {research.get('website_url', 'N/A')}")
        parts.append(f"\n{research.get('research_summary', '')}")

        if research.get('permits_found'):
            parts.append("\n### Permits & Licenses Found:")
            for p in research['permits_found']:
                parts.append(f"- **{p['name']}**: {p['description']} (Timeline: {p['typical_timeline']})")

        if research.get('fee_schedule'):
            parts.append("\n### Fee Schedule:")
            for f in research['fee_schedule']:
                parts.append(f"- {f['permit_type']} - {f['fee_name']}: {f['amount']} ({f['notes']})")

        if research.get('departments'):
            parts.append("\n### Departments:")
            for d in research['departments']:
                parts.append(f"- **{d['name']}**: {d['description']}")

        if research.get('ordinances'):
            parts.append("\n### Municipal Codes & Ordinances:")
            for o in research['ordinances']:
                parts.append(f"- **{o['code']}**: {o['summary']}")
                for prov in o.get('key_provisions', []):
                    parts.append(f"  - {prov}")

        if research.get('processes'):
            parts.append("\n### Standard Processes:")
            for proc in research['processes']:
                parts.append(f"- **{proc['name']}**:")
                for i, step in enumerate(proc['steps'], 1):
                    parts.append(f"  {i}. {step}")

        if research.get('documents_commonly_required'):
            parts.append("\n### Commonly Required Documents:")
            for doc in research['documents_commonly_required']:
                parts.append(f"- {doc}")

        return "\n".join(parts)


web_researcher = WebResearcher()


def summarize_web_content(scraped_data: dict, community_name: str = "") -> dict:
    """AI-powered summarization of scraped web content into structured permit/fee/department data.

    Takes raw scraped pages and uses Claude to extract structured information that
    matches the community_research format expected by the analysis pipeline.
    This is the 'first pass' of the two-pass AI approach.
    """
    if not ANTHROPIC_AVAILABLE or not os.getenv("ANTHROPIC_API_KEY"):
        print("[AI] Skipping web content summarization: AI not available")
        return None

    pages = scraped_data.get("pages", [])
    combined_text = scraped_data.get("combined_text", "")

    if not combined_text and not pages:
        print("[AI] No scraped content to summarize")
        return None

    # Build text from top pages by relevance (cap at 30KB for Claude)
    if not combined_text:
        text_parts = []
        for p in sorted(pages, key=lambda x: x.get("relevance", 0), reverse=True):
            text_parts.append(f"--- {p.get('title', p.get('url', 'Page'))} ---\n{p.get('text', '')}")
        combined_text = '\n\n'.join(text_parts)

    # Sanitize and cap at 30KB to stay within token limits for the summarization pass
    combined_text = _sanitize_to_ascii(combined_text[:30000])

    prompt = f"""You are an expert in municipal government permit/license processes. Analyze the following content scraped from a government website and extract ALL factual information relevant to configuring a PLC (Permitting, Licensing, and Code Enforcement) system.

COMMUNITY NAME: {community_name or 'Unknown'}

SCRAPED WEBSITE CONTENT:
{combined_text}

Extract and return a JSON object with this exact structure (return ONLY valid JSON, no markdown):
{{
  "permits_found": [
    {{"name": "string", "description": "string", "typical_timeline": "string"}}
  ],
  "fee_schedule": [
    {{"permit_type": "string", "fee_name": "string", "amount": "string", "notes": "string"}}
  ],
  "departments": [
    {{"name": "string", "description": "string", "phone": "string"}}
  ],
  "ordinances": [
    {{"code": "string", "summary": "string", "key_provisions": ["string"]}}
  ],
  "processes": [
    {{"name": "string", "steps": ["string"]}}
  ],
  "documents_commonly_required": ["string"],
  "key_findings": "Brief summary of what was found on the website"
}}

IMPORTANT:
- Only include information that was ACTUALLY found on the website. Do not fabricate data.
- If a section has no data found, use an empty array.
- Fee amounts should be exact as stated on the website, or "Not specified" if not found.
- Be thorough — extract every data point you can find."""

    try:
        print(f"[AI] Summarizing {len(combined_text)} chars of scraped web content...")
        client = Anthropic(api_key=_sanitize_to_ascii(os.getenv("ANTHROPIC_API_KEY", "")).strip())
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4000,
            messages=[{"role": "user", "content": prompt}]
        )

        try:
            tokens_used = (response.usage.input_tokens or 0) + (response.usage.output_tokens or 0)
            ai_usage_tracker.record_call("web_content_summarization", tokens_used, True)
            print(f"[AI] Web content summarization complete: {tokens_used} tokens")
        except Exception:
            ai_usage_tracker.record_call("web_content_summarization", 0, True)

        result_text = response.content[0].text.strip()

        # Parse JSON response
        import re
        # Strip markdown code blocks if present
        result_text = re.sub(r'^```(?:json)?\s*', '', result_text)
        result_text = re.sub(r'\s*```$', '', result_text)

        try:
            structured = json.loads(result_text)
        except json.JSONDecodeError:
            # Try to find JSON object in the response
            match = re.search(r'\{[\s\S]*\}', result_text)
            if match:
                structured = json.loads(match.group(0))
            else:
                print(f"[AI] Failed to parse summarization response as JSON")
                return None

        # Build the full research format expected by the pipeline
        research = {
            "community_name": community_name or "the community",
            "website_url": scraped_data.get("base_url", ""),
            "research_summary": structured.get("key_findings", f"Scraped and analyzed {scraped_data.get('pages_scraped', 0)} pages from the community website."),
            "permits_found": structured.get("permits_found", []),
            "fee_schedule": structured.get("fee_schedule", []),
            "departments": structured.get("departments", []),
            "ordinances": structured.get("ordinances", []),
            "processes": structured.get("processes", []),
            "documents_commonly_required": structured.get("documents_commonly_required", []),
            "data_source": "web_scrape",
            "scrape_stats": {
                "pages_scraped": scraped_data.get("pages_scraped", 0),
                "pdfs_found": scraped_data.get("pdfs_found", 0),
                "total_chars_analyzed": len(combined_text),
                "urls_visited": scraped_data.get("urls_visited", 0),
                "scraped_at": scraped_data.get("scraped_at", datetime.utcnow().isoformat()),
            },
            "sources_reviewed": _build_real_sources(scraped_data),
            "research_depth": {
                "pages_analyzed": scraped_data.get("pages_scraped", 0),
                "documents_reviewed": scraped_data.get("pdfs_found", 0),
                "forms_cataloged": len([p for p in scraped_data.get("pdfs", []) if 'form' in p.get('filename', '').lower() or 'application' in p.get('filename', '').lower()]),
                "fee_tables_extracted": len(structured.get("fee_schedule", [])),
                "last_updated": datetime.utcnow().isoformat(),
            }
        }

        return research

    except Exception as e:
        print(f"[AI] Error summarizing web content: {e}")
        ai_usage_tracker.record_call("web_content_summarization", 0, False)
        return None


def _build_real_sources(scraped_data: dict) -> list:
    """Build a sources_reviewed list from actual scraped pages and PDFs."""
    sources = []
    pages = scraped_data.get("pages", [])
    pdfs = scraped_data.get("pdfs", [])

    # Group pages by category based on URL/title keywords
    categories = {
        "Permit & License Pages": [],
        "Fee Schedules": [],
        "Department Pages": [],
        "Forms & Applications": [],
        "Ordinances & Codes": [],
        "Other Relevant Pages": [],
    }

    for page in pages[:30]:  # Top 30 pages
        url = page.get("url", "")
        title = page.get("title", url)
        url_lower = url.lower() + ' ' + title.lower()

        if any(kw in url_lower for kw in ['fee', 'schedule', 'rate', 'cost']):
            categories["Fee Schedules"].append({"title": title, "url": url, "type": "fee_schedule", "description": f"Fee information page (relevance: {page.get('relevance', 0)})"})
        elif any(kw in url_lower for kw in ['permit', 'license', 'application']):
            categories["Permit & License Pages"].append({"title": title, "url": url, "type": "permit_page", "description": f"Permit/license information (relevance: {page.get('relevance', 0)})"})
        elif any(kw in url_lower for kw in ['department', 'division', 'bureau', 'office', 'staff']):
            categories["Department Pages"].append({"title": title, "url": url, "type": "department_page", "description": f"Department information (relevance: {page.get('relevance', 0)})"})
        elif any(kw in url_lower for kw in ['form', 'submittal', 'download']):
            categories["Forms & Applications"].append({"title": title, "url": url, "type": "application_form", "description": f"Forms and applications (relevance: {page.get('relevance', 0)})"})
        elif any(kw in url_lower for kw in ['code', 'ordinance', 'municipal', 'regulation']):
            categories["Ordinances & Codes"].append({"title": title, "url": url, "type": "municipal_code", "description": f"Code/ordinance reference (relevance: {page.get('relevance', 0)})"})
        else:
            categories["Other Relevant Pages"].append({"title": title, "url": url, "type": "general", "description": f"Government page (relevance: {page.get('relevance', 0)})"})

    # Add PDF documents as a category
    if pdfs:
        pdf_sources = []
        for pdf in pdfs[:15]:
            pdf_sources.append({"title": pdf.get("filename", "Document"), "url": pdf.get("url", ""), "type": "pdf_document", "description": f"PDF found on {pdf.get('found_on', 'site')}"})
        if pdf_sources:
            categories["PDF Documents"] = pdf_sources

    for category, cat_sources in categories.items():
        if cat_sources:
            sources.append({"category": category, "sources": cat_sources[:8]})

    return sources


# ============================================================================
# MOCK GENERATOR
# ============================================================================

class MockGenerator:
    @staticmethod
    def generate_configuration() -> Configuration:
        planning = Department(
            id="dept_plan",
            name="Planning & Zoning",
            description="Handles building permits, zoning variances, and land use applications"
        )
        licensing = Department(
            id="dept_lic",
            name="Business Licensing",
            description="Manages business licenses, vendor permits, and professional registrations"
        )
        enforcement = Department(
            id="dept_enf",
            name="Code Enforcement",
            description="Handles code violations, complaints, and compliance inspections"
        )
        fire = Department(
            id="dept_fire",
            name="Fire Prevention",
            description="Fire safety permits and inspections"
        )

        building_permit = RecordType(
            id="rt_bldg",
            name="Building Permit",
            description="Permits for new construction, renovations, additions, and structural modifications",
            department_id="dept_plan",
            category="permit",
            form_fields=[
                FormField(name="Property Address", field_type="address", required=True, description="Full street address of the property"),
                FormField(name="Parcel Number", field_type="text", required=True, description="APN or parcel identification number"),
                FormField(name="Project Description", field_type="textarea", required=True, description="Detailed description of proposed work"),
                FormField(name="Project Type", field_type="select", required=True, description="Type of construction project",
                         options=["New Construction", "Addition", "Renovation", "Demolition", "Tenant Improvement", "Accessory Structure"]),
                FormField(name="Estimated Project Value", field_type="number", required=True, description="Estimated cost of construction"),
                FormField(name="Square Footage", field_type="number", required=True, description="Total square footage of work"),
                FormField(name="Number of Stories", field_type="number", required=False, description="Number of stories in the structure",
                         conditions=[Condition(source_field="Project Type", operator="equals", value="New Construction", description="Only needed for new construction")]),
                FormField(name="Occupancy Type", field_type="select", required=True, description="Building occupancy classification",
                         options=["Residential", "Commercial", "Industrial", "Mixed Use", "Institutional"]),
                FormField(name="Applicant Name", field_type="text", required=True, description="Name of permit applicant"),
                FormField(name="Applicant Email", field_type="email", required=True, description="Email address for notifications"),
                FormField(name="Applicant Phone", field_type="phone", required=True, description="Contact phone number"),
                FormField(name="Contractor Name", field_type="text", required=True, description="Licensed contractor name"),
                FormField(name="Contractor License Number", field_type="text", required=True, description="State contractor license number"),
                FormField(name="Architect/Engineer", field_type="text", required=False, description="Name of project architect or engineer",
                         conditions=[Condition(source_field="Estimated Project Value", operator="greater_than", value="100000", description="Required for projects over $100,000")]),
            ],
            workflow_steps=[
                WorkflowStep(name="Application Submitted", order=1, status_from=None, status_to="submitted", assigned_role="", actions=["Send confirmation email"]),
                WorkflowStep(name="Completeness Check", order=2, status_from="submitted", status_to="under_review", assigned_role="Permit Technician", actions=["Verify all required documents", "Check fees paid"]),
                WorkflowStep(name="Plan Review", order=3, status_from="under_review", status_to="in_plan_review", assigned_role="Plan Reviewer", actions=["Review architectural plans", "Check code compliance", "Verify setbacks and zoning"]),
                WorkflowStep(name="Corrections Required", order=4, status_from="in_plan_review", status_to="corrections_needed", assigned_role="Plan Reviewer", actions=["Document required corrections", "Notify applicant"],
                            conditions=[Condition(source_field="Project Type", operator="not_equals", value="Demolition", description="Demolitions skip corrections cycle")]),
                WorkflowStep(name="Corrections Resubmitted", order=5, status_from="corrections_needed", status_to="in_plan_review", assigned_role="", actions=["Re-review submitted corrections"]),
                WorkflowStep(name="Approved", order=6, status_from="in_plan_review", status_to="approved", assigned_role="Senior Plan Reviewer", actions=["Final sign-off", "Generate permit document"]),
                WorkflowStep(name="Permit Issued", order=7, status_from="approved", status_to="issued", assigned_role="Permit Technician", actions=["Issue permit", "Schedule inspections"]),
            ],
            fees=[
                Fee(name="Building Permit Fee", amount=750.00, fee_type="application", when_applied="upfront", formula="project_value * 0.015"),
                Fee(name="Plan Review Fee", amount=350.00, fee_type="processing", when_applied="upfront"),
                Fee(name="Technology Fee", amount=25.00, fee_type="processing", when_applied="upfront"),
                Fee(name="Inspection Fee", amount=200.00, fee_type="inspection", when_applied="upon_approval",
                    conditions=[Condition(source_field="Project Type", operator="not_equals", value="Demolition", description="No inspection fee for demolitions")]),
            ],
            required_documents=[
                RequiredDocument(name="Architectural Plans", required=True, description="Complete set of construction drawings sealed by licensed architect", stage="application"),
                RequiredDocument(name="Structural Calculations", required=True, description="Engineering calculations for structural elements", stage="application",
                                conditions=[Condition(source_field="Project Type", operator="equals", value="New Construction", description="Required for new construction only")]),
                RequiredDocument(name="Site Plan", required=True, description="Survey showing property boundaries, setbacks, and proposed work", stage="application"),
                RequiredDocument(name="Title Report", required=True, description="Current title report or proof of ownership", stage="application"),
                RequiredDocument(name="Energy Compliance Forms", required=True, description="Title 24 energy compliance documentation", stage="review"),
                RequiredDocument(name="Soils Report", required=False, description="Geotechnical investigation report", stage="review",
                                conditions=[Condition(source_field="Project Type", operator="equals", value="New Construction", description="Required for new construction only")]),
            ]
        )

        business_license = RecordType(
            id="rt_bus",
            name="Business License",
            description="Annual business operating license required for all commercial activities within city limits",
            department_id="dept_lic",
            category="license",
            form_fields=[
                FormField(name="Business Name", field_type="text", required=True, description="Legal business name"),
                FormField(name="DBA Name", field_type="text", required=False, description="Doing Business As name if different",
                         conditions=[Condition(source_field="Business Type", operator="not_equals", value="Home-Based Business", description="Not applicable for home-based businesses")]),
                FormField(name="Business Type", field_type="select", required=True, description="Type of business",
                         options=["Retail", "Restaurant/Food Service", "Professional Services", "Contractor", "Home-Based Business", "Manufacturing", "Wholesale", "Entertainment"]),
                FormField(name="Business Address", field_type="address", required=True, description="Physical business location"),
                FormField(name="Mailing Address", field_type="address", required=False, description="Mailing address if different from business address",
                         conditions=[Condition(source_field="Business Type", operator="not_equals", value="Home-Based Business", description="Home-based businesses use home address")]),
                FormField(name="Owner Name", field_type="text", required=True, description="Business owner full name"),
                FormField(name="Owner Email", field_type="email", required=True, description="Owner email for correspondence"),
                FormField(name="Owner Phone", field_type="phone", required=True, description="Owner contact phone"),
                FormField(name="Federal EIN", field_type="text", required=True, description="Federal Employer Identification Number"),
                FormField(name="State Tax ID", field_type="text", required=True, description="State tax identification number"),
                FormField(name="Number of Employees", field_type="number", required=True, description="Total number of employees"),
                FormField(name="Estimated Annual Revenue", field_type="select", required=True, description="Estimated gross annual revenue range",
                         options=["Under $50,000", "$50,000-$100,000", "$100,000-$500,000", "$500,000-$1,000,000", "Over $1,000,000"]),
            ],
            workflow_steps=[
                WorkflowStep(name="Application Submitted", order=1, status_from=None, status_to="submitted", assigned_role="", actions=["Send confirmation email", "Assign to licensing officer"]),
                WorkflowStep(name="Application Review", order=2, status_from="submitted", status_to="under_review", assigned_role="Licensing Officer", actions=["Verify business information", "Check zoning compliance"]),
                WorkflowStep(name="Approved", order=3, status_from="under_review", status_to="approved", assigned_role="Licensing Supervisor", actions=["Approve license", "Generate license document"]),
                WorkflowStep(name="License Issued", order=4, status_from="approved", status_to="active", assigned_role="Licensing Officer", actions=["Issue license", "Set renewal reminder"]),
                WorkflowStep(name="Renewal Due", order=5, status_from="active", status_to="renewal_pending", assigned_role="", actions=["Send renewal notice"]),
            ],
            fees=[
                Fee(name="Business License Fee", amount=150.00, fee_type="application", when_applied="upfront"),
                Fee(name="Processing Fee", amount=25.00, fee_type="processing", when_applied="upfront"),
                Fee(name="Annual Renewal Fee", amount=100.00, fee_type="annual", when_applied="annual", formula="base_fee + (employees * 5)"),
            ],
            required_documents=[
                RequiredDocument(name="State Business Registration", required=True, description="State business entity registration certificate", stage="application"),
                RequiredDocument(name="Proof of Insurance", required=True, description="General liability insurance certificate", stage="application"),
                RequiredDocument(name="Lease Agreement", required=True, description="Copy of commercial lease or proof of property ownership", stage="application",
                                conditions=[Condition(source_field="Business Type", operator="not_equals", value="Home-Based Business", description="Not required for home-based businesses")]),
                RequiredDocument(name="Health Permit", required=False, description="Required for food service businesses", stage="approval",
                                conditions=[Condition(source_field="Business Type", operator="equals", value="Restaurant/Food Service", description="Required for food service businesses")]),
            ]
        )

        code_violation = RecordType(
            id="rt_code",
            name="Code Violation Notice",
            description="Documentation and tracking of municipal code violations including property maintenance, zoning, and building code violations",
            department_id="dept_enf",
            category="code_enforcement",
            form_fields=[
                FormField(name="Violation Address", field_type="address", required=True, description="Address where violation was observed"),
                FormField(name="Violation Type", field_type="select", required=True, description="Category of code violation",
                         options=["Property Maintenance", "Zoning Violation", "Building Without Permit", "Sign Violation", "Noise Complaint", "Overgrown Vegetation", "Abandoned Vehicle", "Illegal Dumping"]),
                FormField(name="Violation Description", field_type="textarea", required=True, description="Detailed description of the violation"),
                FormField(name="Date Observed", field_type="date", required=True, description="Date violation was first observed"),
                FormField(name="Reported By", field_type="select", required=True, description="Source of the complaint",
                         options=["Citizen Complaint", "Staff Observation", "Other Agency", "Anonymous"]),
                FormField(name="Reporter Name", field_type="text", required=False, description="Name of person reporting (if not anonymous)",
                         conditions=[Condition(source_field="Reported By", operator="not_equals", value="Anonymous", description="Not needed for anonymous reports")]),
                FormField(name="Reporter Phone", field_type="phone", required=False, description="Reporter contact number",
                         conditions=[Condition(source_field="Reported By", operator="not_equals", value="Anonymous", description="Not needed for anonymous reports")]),
                FormField(name="Property Owner", field_type="text", required=True, description="Name of property owner"),
                FormField(name="Severity", field_type="select", required=True, description="Severity level of violation",
                         options=["Minor", "Moderate", "Major", "Imminent Hazard"]),
            ],
            workflow_steps=[
                WorkflowStep(name="Complaint Received", order=1, status_from=None, status_to="received", assigned_role="", actions=["Log complaint", "Assign to inspector"]),
                WorkflowStep(name="Initial Inspection", order=2, status_from="received", status_to="inspected", assigned_role="Code Inspector", actions=["Site visit", "Document violations", "Take photos"]),
                WorkflowStep(name="Notice of Violation Issued", order=3, status_from="inspected", status_to="notice_issued", assigned_role="Code Inspector", actions=["Send NOV to property owner", "Set compliance deadline"]),
                WorkflowStep(name="Compliance Check", order=4, status_from="notice_issued", status_to="compliance_check", assigned_role="Code Inspector", actions=["Re-inspect property", "Document status"]),
                WorkflowStep(name="Resolved", order=5, status_from="compliance_check", status_to="resolved", assigned_role="Code Inspector", actions=["Close case", "Update records"]),
                WorkflowStep(name="Citation Issued", order=6, status_from="compliance_check", status_to="citation_issued", assigned_role="Senior Inspector", actions=["Issue citation", "Refer to legal if needed"],
                            conditions=[Condition(source_field="Severity", operator="not_equals", value="Minor", description="Citations only for moderate+ violations")]),
            ],
            fees=[
                Fee(name="Re-inspection Fee", amount=150.00, fee_type="inspection", when_applied="upon_inspection"),
                Fee(name="Administrative Citation", amount=500.00, fee_type="permit", when_applied="upon_approval",
                    conditions=[Condition(source_field="Severity", operator="not_equals", value="Minor", description="Citations not issued for minor violations")]),
            ],
            required_documents=[
                RequiredDocument(name="Violation Photos", required=True, description="Photographic evidence of the violation", stage="inspection"),
                RequiredDocument(name="Compliance Plan", required=False, description="Owner's plan to remedy the violation", stage="review"),
                RequiredDocument(name="Correction Verification", required=True, description="Photos showing violation has been corrected", stage="inspection"),
            ]
        )

        encroachment = RecordType(
            id="rt_enc",
            name="Encroachment Permit",
            description="Permits for work within public right-of-way including sidewalk repairs, utility connections, and driveway approaches",
            department_id="dept_plan",
            category="permit",
            form_fields=[
                FormField(name="Location", field_type="address", required=True, description="Location of proposed encroachment"),
                FormField(name="Work Description", field_type="textarea", required=True, description="Description of work to be performed in right-of-way"),
                FormField(name="Encroachment Type", field_type="select", required=True, description="Type of encroachment",
                         options=["Sidewalk Repair", "Driveway Approach", "Utility Connection", "Street Cut", "Temporary Closure", "Banner/Sign"]),
                FormField(name="Start Date", field_type="date", required=True, description="Proposed start date of work"),
                FormField(name="End Date", field_type="date", required=True, description="Expected completion date"),
                FormField(name="Applicant Name", field_type="text", required=True, description="Name of applicant"),
                FormField(name="Applicant Email", field_type="email", required=True, description="Email for notifications"),
                FormField(name="Contractor Name", field_type="text", required=True, description="Contractor performing the work"),
                FormField(name="Traffic Control Plan Required", field_type="checkbox", required=False, description="Whether a TCP is needed",
                         conditions=[Condition(source_field="Encroachment Type", operator="equals", value="Street Cut", description="TCP required for street cuts")]),
            ],
            workflow_steps=[
                WorkflowStep(name="Application Submitted", order=1, status_from=None, status_to="submitted", assigned_role="", actions=["Send confirmation"]),
                WorkflowStep(name="Engineering Review", order=2, status_from="submitted", status_to="under_review", assigned_role="Public Works Engineer", actions=["Review plans", "Check conflicts"]),
                WorkflowStep(name="Approved", order=3, status_from="under_review", status_to="approved", assigned_role="Public Works Director", actions=["Issue permit"]),
                WorkflowStep(name="Work Completed", order=4, status_from="approved", status_to="work_complete", assigned_role="", actions=["Request final inspection"]),
                WorkflowStep(name="Final Inspection", order=5, status_from="work_complete", status_to="closed", assigned_role="Public Works Inspector", actions=["Inspect work", "Close permit"]),
            ],
            fees=[
                Fee(name="Encroachment Permit Fee", amount=250.00, fee_type="application", when_applied="upfront"),
                Fee(name="Deposit", amount=1000.00, fee_type="permit", when_applied="upfront"),
                Fee(name="Inspection Fee", amount=100.00, fee_type="inspection", when_applied="upon_inspection"),
            ],
            required_documents=[
                RequiredDocument(name="Site Plan/Drawing", required=True, description="Drawing showing location and extent of work", stage="application"),
                RequiredDocument(name="Traffic Control Plan", required=False, description="Required for work affecting traffic flow", stage="application",
                                conditions=[Condition(source_field="Encroachment Type", operator="equals", value="Street Cut", description="Required when work affects traffic")]),
                RequiredDocument(name="Insurance Certificate", required=True, description="Liability insurance naming city as additional insured", stage="application"),
            ]
        )

        fire_permit = RecordType(
            id="rt_fire",
            name="Fire Prevention Permit",
            description="Permits for fire safety compliance including fire alarm systems, sprinkler systems, and special events",
            department_id="dept_fire",
            category="permit",
            form_fields=[
                FormField(name="Property Address", field_type="address", required=True, description="Address of property"),
                FormField(name="Permit Type", field_type="select", required=True, description="Type of fire permit",
                         options=["Fire Alarm System", "Sprinkler System", "Hood Suppression", "Special Event", "Hazardous Materials", "Assembly Permit"]),
                FormField(name="Description of Work", field_type="textarea", required=True, description="Description of fire protection system or activity"),
                FormField(name="Applicant Name", field_type="text", required=True, description="Applicant name"),
                FormField(name="Applicant Phone", field_type="phone", required=True, description="Contact phone"),
                FormField(name="Contractor Name", field_type="text", required=True, description="Fire protection contractor"),
                FormField(name="Contractor License", field_type="text", required=True, description="C-16 or other fire protection license number"),
                FormField(name="Occupancy Type", field_type="select", required=True, description="Building occupancy type",
                         options=["Assembly", "Business", "Educational", "Factory", "Hazardous", "Institutional", "Mercantile", "Residential", "Storage"]),
            ],
            workflow_steps=[
                WorkflowStep(name="Application Submitted", order=1, status_from=None, status_to="submitted", assigned_role="", actions=["Send confirmation"]),
                WorkflowStep(name="Fire Marshal Review", order=2, status_from="submitted", status_to="under_review", assigned_role="Fire Marshal", actions=["Review plans and specifications"]),
                WorkflowStep(name="Approved", order=3, status_from="under_review", status_to="approved", assigned_role="Fire Marshal", actions=["Issue permit"]),
                WorkflowStep(name="Installation Complete", order=4, status_from="approved", status_to="inspection_requested", assigned_role="", actions=["Schedule fire inspection"]),
                WorkflowStep(name="Final Inspection", order=5, status_from="inspection_requested", status_to="closed", assigned_role="Fire Inspector", actions=["Inspect installation", "Issue certificate"]),
            ],
            fees=[
                Fee(name="Fire Permit Fee", amount=200.00, fee_type="application", when_applied="upfront"),
                Fee(name="Plan Review Fee", amount=150.00, fee_type="processing", when_applied="upfront"),
                Fee(name="Inspection Fee", amount=125.00, fee_type="inspection", when_applied="upon_inspection",
                    conditions=[Condition(source_field="Permit Type", operator="not_equals", value="Special Event", description="No inspection for special events")]),
            ],
            required_documents=[
                RequiredDocument(name="Fire Protection Plans", required=True, description="Engineered fire protection system plans", stage="application"),
                RequiredDocument(name="Product Specifications", required=True, description="Cut sheets for fire protection equipment", stage="application"),
                RequiredDocument(name="Hydraulic Calculations", required=False, description="For sprinkler systems", stage="review",
                                conditions=[Condition(source_field="Permit Type", operator="equals", value="Sprinkler System", description="Required for sprinkler system permits")]),
            ]
        )

        planning.record_type_ids = ["rt_bldg", "rt_enc"]
        licensing.record_type_ids = ["rt_bus"]
        enforcement.record_type_ids = ["rt_code"]
        fire.record_type_ids = ["rt_fire"]

        roles = [
            UserRole(
                id="role_admin",
                name="System Administrator",
                description="Full system access for IT and management",
                permissions=["manage_users", "manage_config", "view_all", "edit_all", "delete_all", "generate_reports", "manage_fees"],
                department_ids=["dept_plan", "dept_lic", "dept_enf", "dept_fire"]
            ),
            UserRole(
                id="role_reviewer",
                name="Plan Reviewer",
                description="Reviews building permit plans and applications",
                permissions=["view_permits", "review_plans", "request_corrections", "approve_permits", "add_comments"],
                department_ids=["dept_plan"]
            ),
            UserRole(
                id="role_tech",
                name="Permit Technician",
                description="Processes permit applications and handles front counter",
                permissions=["create_permits", "view_permits", "process_payments", "issue_permits", "upload_documents"],
                department_ids=["dept_plan", "dept_lic"]
            ),
            UserRole(
                id="role_inspector",
                name="Inspector",
                description="Conducts field inspections for permits and code enforcement",
                permissions=["view_permits", "view_violations", "schedule_inspections", "record_results", "upload_photos", "add_comments"],
                department_ids=["dept_plan", "dept_enf", "dept_fire"]
            ),
            UserRole(
                id="role_lic_officer",
                name="Licensing Officer",
                description="Processes business license applications",
                permissions=["view_licenses", "review_applications", "approve_licenses", "process_payments", "generate_licenses"],
                department_ids=["dept_lic"]
            ),
            UserRole(
                id="role_code_officer",
                name="Code Enforcement Officer",
                description="Handles code violation cases",
                permissions=["view_violations", "create_violations", "issue_notices", "issue_citations", "schedule_inspections", "upload_photos"],
                department_ids=["dept_enf"]
            ),
            UserRole(
                id="role_public",
                name="Public Applicant",
                description="External users who submit applications and check status",
                permissions=["submit_applications", "view_own_records", "upload_documents", "make_payments", "check_status"],
                department_ids=[]
            ),
        ]

        return Configuration(
            record_types=[building_permit, business_license, code_violation, encroachment, fire_permit],
            departments=[planning, licensing, enforcement, fire],
            user_roles=roles,
            generated_at=datetime.utcnow(),
            summary="Generated 5 record types across 4 departments with 7 user roles. Includes Building Permits, Business Licenses, Code Violations, Encroachment Permits, and Fire Prevention Permits with complete form fields, workflows, fee schedules, and document requirements."
        )

    @staticmethod
    def generate_sample_csv() -> str:
        """Generate a realistic sample CSV for demo purposes"""
        import random
        headers = "permit_number,record_type,applicant_name,applicant_email,applicant_phone,property_address,description,status,department,assigned_to,submitted_date,approved_date,fee_amount,fee_type,documents_submitted,project_value"

        record_types = ["Building Permit", "Business License", "Code Violation", "Encroachment Permit", "Fire Prevention Permit"]
        statuses = ["Submitted", "Under Review", "In Plan Review", "Corrections Needed", "Approved", "Issued", "Closed", "Denied"]
        departments = ["Planning & Zoning", "Business Licensing", "Code Enforcement", "Fire Prevention", "Public Works"]
        staff = ["Alice Johnson", "Bob Chen", "Carol Williams", "David Martinez", "Emily Brown", "Frank Garcia"]

        rows = [headers]
        for i in range(150):
            rt = random.choice(record_types)
            dept = departments[record_types.index(rt)] if rt in record_types else random.choice(departments)
            status = random.choice(statuses)
            fee = round(random.uniform(25, 2000), 2)
            pv = round(random.uniform(5000, 500000), 2) if rt == "Building Permit" else ""

            row = f"PRM-2024-{i+1:04d},{rt},Applicant {i+1},applicant{i+1}@email.com,(555) {random.randint(100,999)}-{random.randint(1000,9999)},{random.randint(100,9999)} Main St,Description for {rt} #{i+1},{status},{dept},{random.choice(staff)},2024-{random.randint(1,12):02d}-{random.randint(1,28):02d},,{fee},Application Fee,{random.randint(1,5)},{pv}"
            rows.append(row)

        return "\n".join(rows)


# ============================================================================
# TEXT SANITIZATION (must be defined before ClaudeService uses it)
# ============================================================================

def _sanitize_to_ascii(text: str) -> str:
    """Replace non-ASCII characters (including Cyrillic homoglyphs) with ASCII equivalents.

    Government websites sometimes contain Cyrillic look-alikes (e.g. Cyrillic С vs Latin C)
    from copy-pasted Word/PDF content. This prevents encoding errors on ASCII-only runtimes.
    """
    if not text:
        return text
    # Common Cyrillic-to-Latin homoglyph mappings
    homoglyphs = {
        '\u0410': 'A', '\u0412': 'B', '\u0421': 'C', '\u0415': 'E',
        '\u041d': 'H', '\u041a': 'K', '\u041c': 'M', '\u041e': 'O',
        '\u0420': 'P', '\u0422': 'T', '\u0425': 'X',
        '\u0430': 'a', '\u0435': 'e', '\u043e': 'o', '\u0440': 'p',
        '\u0441': 'c', '\u0443': 'y', '\u0445': 'x',
        # Common special characters
        '\u2013': '-', '\u2014': '--', '\u2018': "'", '\u2019': "'",
        '\u201c': '"', '\u201d': '"', '\u2022': '-', '\u2026': '...',
        '\u00ae': '(R)', '\u00a9': '(c)', '\u2122': '(TM)',
        '\u00b0': ' deg', '\u00bd': '1/2', '\u00bc': '1/4', '\u00be': '3/4',
        '\u00e9': 'e', '\u00e8': 'e', '\u00f1': 'n', '\u00fc': 'u',
        '\u00a0': ' ',  # non-breaking space
    }
    for old_char, new_char in homoglyphs.items():
        text = text.replace(old_char, new_char)
    # Drop any remaining non-ASCII characters
    return text.encode('ascii', errors='ignore').decode('ascii')


# ============================================================================
# CLAUDE SERVICE
# ============================================================================

class ClaudeService:
    def __init__(self):
        raw_key = os.getenv("ANTHROPIC_API_KEY", "")

        # Detect and fix non-ASCII characters in API key (common copy-paste issue
        # where Cyrillic look-alikes like С/с replace C/c)
        if raw_key:
            clean_key = _sanitize_to_ascii(raw_key).strip()
            if clean_key != raw_key:
                bad_chars = [(i, repr(ch)) for i, ch in enumerate(raw_key) if ord(ch) > 127]
                print(f"[AI] WARNING: ANTHROPIC_API_KEY contains non-ASCII characters at positions: {bad_chars}")
                print(f"[AI] Auto-cleaned API key (was {len(raw_key)} chars, now {len(clean_key)} chars)")
            self.api_key = clean_key
        else:
            self.api_key = ""

        self.client = None
        self.last_mode = "unknown"  # Track whether last call used real AI or mock
        self.last_error = None
        if ANTHROPIC_AVAILABLE and self.api_key:
            try:
                self.client = Anthropic(api_key=self.api_key)
                print(f"[AI] Claude service initialized with API key ({self.api_key[:12]}...)")
            except Exception as e:
                print(f"[AI] Failed to initialize Anthropic client: {e}")
                self.last_error = str(e)
        else:
            if not ANTHROPIC_AVAILABLE:
                print("[AI] WARNING: anthropic package not installed - using mock mode")
            elif not self.api_key:
                print("[AI] WARNING: ANTHROPIC_API_KEY not set - using mock mode")

    def is_available(self) -> bool:
        return self.client is not None and bool(self.api_key)

    def get_status(self) -> dict:
        """Return detailed AI service status"""
        return {
            "available": self.is_available(),
            "package_installed": ANTHROPIC_AVAILABLE,
            "api_key_configured": bool(self.api_key),
            "api_key_preview": f"{self.api_key[:12]}..." if self.api_key else "not set",
            "last_mode": self.last_mode,
            "last_error": self.last_error,
        }

    def analyze_csv_data(self, csv_summary: str) -> Configuration:
        """Send CSV data to Claude for analysis, return Configuration"""
        if not self.is_available():
            self.last_mode = "mock"
            reason = "anthropic package not installed" if not ANTHROPIC_AVAILABLE else "ANTHROPIC_API_KEY not configured"
            self.last_error = reason
            print(f"[AI] Using MOCK configuration: {reason}")
            config = MockGenerator.generate_configuration()
            config.summary = f"[MOCK DATA] {reason}. Configure your Anthropic API key in Vercel environment variables for real AI analysis."
            return config

        # Trim input to avoid excessively long prompts that slow response
        trimmed_summary = csv_summary[:25000] if len(csv_summary) > 25000 else csv_summary
        prompt = self._build_prompt(trimmed_summary)

        try:
            print(f"[AI] Calling Claude API for configuration analysis ({len(prompt)} chars prompt)...")
            message = self.client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=16000,
                messages=[{"role": "user", "content": prompt}]
            )

            # Track AI usage
            stop_reason = getattr(message, 'stop_reason', 'unknown')
            try:
                tokens_used = (message.usage.input_tokens or 0) + (message.usage.output_tokens or 0)
                ai_usage_tracker.record_call("configuration_analysis", tokens_used, True)
                print(f"[AI] Claude API call successful - {tokens_used} tokens, stop_reason={stop_reason}")
            except Exception:
                ai_usage_tracker.record_call("configuration_analysis", 0, True)

            response_text = message.content[0].text

            # If truncated, ask Claude to complete it
            if stop_reason == 'max_tokens':
                print("[AI] Response was truncated at max_tokens, requesting continuation...")
                try:
                    continuation = self.client.messages.create(
                        model="claude-sonnet-4-20250514",
                        max_tokens=8000,
                        messages=[
                            {"role": "user", "content": prompt},
                            {"role": "assistant", "content": response_text},
                            {"role": "user", "content": "Your JSON was truncated. Continue EXACTLY where you left off to complete the JSON. Output ONLY the remaining JSON, no explanation."}
                        ]
                    )
                    response_text += continuation.content[0].text
                    print(f"[AI] Continuation received, total response: {len(response_text)} chars")
                except Exception as cont_err:
                    print(f"[AI] Continuation failed: {cont_err}")

            config_data = self._parse_response(response_text)
            self.last_mode = "ai"
            self.last_error = None
            return self._build_configuration(config_data)
        except Exception as e:
            self.last_mode = "mock"
            self.last_error = str(e)
            print(f"[AI] ERROR during Claude API call: {e}")
            ai_usage_tracker.record_call("configuration_analysis", 0, False)
            config = MockGenerator.generate_configuration()
            config.summary = f"[MOCK DATA - AI Error] {str(e)[:200]}. The AI call failed, showing sample data instead."
            return config

    def _build_prompt(self, csv_summary: str) -> str:
        return f"""You are a senior implementation consultant specializing in PLC (Permitting, Licensing, and Code Enforcement) software for local governments. You have configured hundreds of municipalities.

Analyze ALL provided data sources and generate the most thorough, production-ready PLC configuration possible.

## Data Sources Provided:
{csv_summary}

## Your Task:
Create a COMPREHENSIVE PLC configuration. Think like a consultant who has done 100+ implementations. Don't just map what's in the CSV — anticipate what this jurisdiction NEEDS based on:
- The types of permits/licenses you see in their data
- Standard requirements for similar municipalities
- Industry best practices for government workflows
- Common fee structures for this type of jurisdiction
- Typical document requirements for each permit/license type

### For EACH Record Type, provide:
1. **Name & Description**: Clear, professional naming
2. **Category**: permit, license, code_enforcement, or inspection
3. **Form Fields** (6-10 per record type): Key applicant and staff fields
4. **Workflow Steps** (4-6 per record type): Lifecycle from submission to close
5. **Fees** (2-4 per record type): Realistic fee structures with dollar amounts
6. **Required Documents** (2-4 per record type): What applicants must submit

### Generate:
- 5-6 Record Types (building permits, business licenses, code enforcement, and 2-3 specialty permits)
- 3-5 Departments
- 4-6 User Roles with permissions

## IMPORTANT RULES:
- Keep descriptions SHORT (1 sentence max) to keep JSON compact
- Fee amounts should be realistic ($50-$5000 range)
- Workflow steps should include assigned roles
- CRITICAL: Your entire response must be VALID JSON and COMPLETE. Do not get cut off.

## Response Format:
Return ONLY valid JSON (no markdown, no explanation) matching this exact schema:
{{
  "record_types": [
    {{
      "name": "string",
      "description": "string",
      "category": "permit|license|code_enforcement|inspection",
      "department": "string",
      "form_fields": [
        {{"name": "string", "field_type": "text|email|phone|date|number|select|textarea|file|checkbox|address", "required": true/false, "description": "string", "options": ["opt1"] or null}}
      ],
      "workflow_steps": [
        {{"name": "string", "order": 1, "assigned_role": "string", "status_from": "string or null", "status_to": "string", "actions": ["action1"]}}
      ],
      "fees": [
        {{"name": "string", "amount": 0.00, "fee_type": "application|processing|permit|inspection|annual", "when_applied": "upfront|upon_approval|upon_inspection|annual"}}
      ],
      "required_documents": [
        {{"name": "string", "required": true/false, "description": "string", "stage": "application|review|approval|inspection"}}
      ]
    }}
  ],
  "departments": [
    {{"name": "string", "description": "string"}}
  ],
  "user_roles": [
    {{"name": "string", "description": "string", "permissions": ["perm1"], "departments": ["dept1"]}}
  ],
  "summary": "Brief summary of what was analyzed and generated"
}}"""

    def _repair_json_strings(self, text: str) -> str:
        """Walk through JSON text and fix unescaped quotes inside string values"""
        result = []
        i = 0
        length = len(text)
        in_string = False
        escape_next = False

        while i < length:
            c = text[i]

            if escape_next:
                result.append(c)
                escape_next = False
                i += 1
                continue

            if c == '\\' and in_string:
                escape_next = True
                result.append(c)
                i += 1
                continue

            if c == '"':
                if not in_string:
                    in_string = True
                    result.append(c)
                else:
                    # Check if this quote truly ends the string
                    # After closing quote, we should see: , } ] : or whitespace then one of those
                    j = i + 1
                    while j < length and text[j] in ' \t\r\n':
                        j += 1
                    if j >= length or text[j] in ',}]:':
                        # Real end of string
                        in_string = False
                        result.append(c)
                    else:
                        # Embedded quote inside string — escape it
                        result.append('\\"')
                i += 1
                continue

            if in_string and c == '\n':
                result.append('\\n')
                i += 1
                continue
            if in_string and c == '\r':
                result.append('\\r')
                i += 1
                continue
            if in_string and c == '\t':
                result.append('\\t')
                i += 1
                continue

            result.append(c)
            i += 1

        return ''.join(result)

    def _parse_response(self, text: str) -> dict:
        """Extract JSON from Claude's response with robust error recovery"""
        import re

        text = text.strip()

        # Strip markdown code blocks
        if "```" in text:
            lines = text.split("\n")
            json_lines = []
            in_block = False
            for line in lines:
                stripped = line.strip()
                if stripped.startswith("```"):
                    in_block = not in_block
                    continue
                if in_block:
                    json_lines.append(line)
            if json_lines:
                text = "\n".join(json_lines).strip()

        # Extract JSON object if there's surrounding text
        start = text.find("{")
        end = text.rfind("}") + 1
        if start >= 0 and end > start:
            text = text[start:end]

        # Attempt 1: Direct parse
        try:
            return json.loads(text)
        except json.JSONDecodeError as e:
            print(f"[AI] Direct parse failed: {e}")

        # Attempt 2: Deep repair — fix unescaped quotes, newlines, trailing commas, etc.
        fixed = text
        # Fix Python-style booleans/None
        fixed = re.sub(r':\s*None\b', ': null', fixed)
        fixed = re.sub(r':\s*True\b', ': true', fixed)
        fixed = re.sub(r':\s*False\b', ': false', fixed)
        # Remove trailing commas
        fixed = re.sub(r',\s*([}\]])', r'\1', fixed)
        # Repair unescaped quotes and control characters inside strings
        fixed = self._repair_json_strings(fixed)
        try:
            return json.loads(fixed)
        except json.JSONDecodeError as e:
            print(f"[AI] Repair parse failed: {e}")

        # Attempt 3: If JSON is truncated, close it
        open_braces = fixed.count('{') - fixed.count('}')
        open_brackets = fixed.count('[') - fixed.count(']')
        if open_braces > 0 or open_brackets > 0:
            truncated = fixed.rstrip()
            # Remove trailing partial key/value
            truncated = re.sub(r',?\s*"[^"]*"?\s*:?\s*"?[^"{}[\],]*$', '', truncated)
            truncated += ']' * max(0, open_brackets) + '}' * max(0, open_braces)
            try:
                return json.loads(truncated)
            except json.JSONDecodeError as e:
                print(f"[AI] Truncation repair failed: {e}")

        # Attempt 4: Salvage record_types array
        try:
            rt_match = re.search(r'"record_types"\s*:\s*\[', fixed)
            if rt_match:
                depth = 0
                arr_start = fixed.index('[', rt_match.start())
                for idx in range(arr_start, len(fixed)):
                    if fixed[idx] == '[': depth += 1
                    elif fixed[idx] == ']': depth -= 1
                    if depth == 0:
                        rt_json = fixed[arr_start:idx+1]
                        record_types = json.loads(rt_json)
                        # Also try to find departments and user_roles
                        depts = []
                        roles = []
                        try:
                            dm = re.search(r'"departments"\s*:\s*\[', fixed)
                            if dm:
                                d2 = 0
                                ds = fixed.index('[', dm.start())
                                for di in range(ds, len(fixed)):
                                    if fixed[di] == '[': d2 += 1
                                    elif fixed[di] == ']': d2 -= 1
                                    if d2 == 0:
                                        depts = json.loads(fixed[ds:di+1])
                                        break
                        except Exception:
                            pass
                        try:
                            rm = re.search(r'"user_roles"\s*:\s*\[', fixed)
                            if rm:
                                r2 = 0
                                rs = fixed.index('[', rm.start())
                                for ri in range(rs, len(fixed)):
                                    if fixed[ri] == '[': r2 += 1
                                    elif fixed[ri] == ']': r2 -= 1
                                    if r2 == 0:
                                        roles = json.loads(fixed[rs:ri+1])
                                        break
                        except Exception:
                            pass
                        return {
                            "record_types": record_types,
                            "departments": depts,
                            "user_roles": roles,
                            "summary": "Recovered from partial AI response"
                        }
        except (json.JSONDecodeError, ValueError) as e:
            print(f"[AI] Salvage parse failed: {e}")

        print(f"[AI] JSON parse FAILED. Response length: {len(text)} chars")
        print(f"[AI] First 500 chars: {text[:500]}")
        print(f"[AI] Around char 15000: {text[14800:15200] if len(text) > 15200 else 'N/A'}")
        print(f"[AI] Last 500 chars: {text[-500:]}")
        raise ValueError("Could not parse Claude response as JSON after all recovery attempts")

    def _build_configuration(self, data: dict) -> Configuration:
        """Transform Claude's JSON response into Configuration model"""
        departments = []
        dept_map = {}
        for d in data.get("departments", []):
            dept = Department(name=d["name"], description=d.get("description", ""))
            departments.append(dept)
            dept_map[d["name"]] = dept.id

        roles = []
        for r in data.get("user_roles", []):
            role = UserRole(
                name=r["name"],
                description=r.get("description", ""),
                permissions=r.get("permissions", []),
                department_ids=[dept_map.get(d, "") for d in r.get("departments", [])]
            )
            roles.append(role)

        record_types = []
        for rt in data.get("record_types", []):
            dept_id = dept_map.get(rt.get("department", ""), "")

            form_fields = [FormField(**f) for f in rt.get("form_fields", [])]
            workflow_steps = [WorkflowStep(**w) for w in rt.get("workflow_steps", [])]
            fees = [Fee(**f) for f in rt.get("fees", [])]
            docs = [RequiredDocument(**d) for d in rt.get("required_documents", [])]

            record_type = RecordType(
                name=rt["name"],
                description=rt.get("description", ""),
                department_id=dept_id,
                category=rt.get("category", "permit"),
                form_fields=form_fields,
                workflow_steps=workflow_steps,
                fees=fees,
                required_documents=docs,
            )
            record_types.append(record_type)

            if dept_id:
                for dept in departments:
                    if dept.id == dept_id:
                        dept.record_type_ids.append(record_type.id)

        return Configuration(
            record_types=record_types,
            departments=departments,
            user_roles=roles,
            generated_at=datetime.utcnow(),
            summary=data.get("summary", "Configuration generated from CSV analysis")
        )


claude_service = ClaudeService()


# ============================================================================
# FASTAPI APP SETUP
# ============================================================================

UPLOAD_DIR = os.environ.get("UPLOAD_DIR", "/tmp/plc-uploads")


@asynccontextmanager
async def lifespan(app: FastAPI):
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    yield


app = FastAPI(
    title="PLC AutoConfig Backend",
    description="Backend for AI-powered PLC software configuration from CSV data",
    version="1.0.0",
    lifespan=lifespan
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ============================================================================
# PROJECT ROUTES
# ============================================================================

@app.post("/api/projects", response_model=Project)
async def create_project(request: CreateProjectRequest):
    """Create a new project. If an ID is provided and that project already exists, return it."""
    # If ID provided, check if project already exists (cold-start recovery)
    if request.id:
        existing = store.get_project(request.id)
        if existing:
            print(f"[API] Project {request.id} already exists, returning existing")
            return existing

    project = Project(
        id=request.id if request.id else str(uuid.uuid4())[:8],
        name=request.name,
        customer_name=request.customer_name,
        product_type=request.product_type,
        community_url=request.community_url or "",
    )
    store.create_project(project)
    print(f"[API] Created project {project.id} '{project.name}' | KV={KV_AVAILABLE}")
    return project


@app.get("/api/projects", response_model=List[Project])
async def list_projects():
    """List all projects"""
    projects = store.list_projects()
    print(f"[API] List projects: {len(projects)} found | KV={KV_AVAILABLE}")
    return projects


@app.get("/api/projects/{project_id}", response_model=Project)
async def get_project(project_id: str):
    """Get a specific project"""
    print(f"[API] Getting project {project_id} | KV={KV_AVAILABLE} | in_memory={project_id in store._projects}")
    project = store.get_project(project_id)
    if not project:
        print(f"[API] Project {project_id} NOT FOUND after checking memory/disk/KV")
        raise HTTPException(status_code=404, detail="Project not found")
    print(f"[API] Project {project_id} found: '{project.name}'")
    return project


@app.delete("/api/projects/{project_id}")
async def delete_project(project_id: str):
    """Delete a project"""
    if not store.delete_project(project_id):
        raise HTTPException(status_code=404, detail="Project not found")
    return {"success": True}


@app.post("/api/projects/{project_id}/upload")
async def upload_files(project_id: str, files: List[UploadFile] = File(...)):
    """Upload CSV files to a project"""
    import time as _time
    # Retry project lookup — handles Vercel cold starts where KV needs a moment
    project = None
    for attempt in range(3):
        project = store.get_project(project_id)
        if project:
            break
        print(f"[UPLOAD] Project {project_id} not found on attempt {attempt + 1}/3, retrying...")
        _time.sleep(0.5)
    if not project:
        print(f"[UPLOAD] Project {project_id} NOT FOUND after 3 attempts | KV={KV_AVAILABLE}")
        raise HTTPException(status_code=404, detail=f"Project not found. KV connected: {KV_AVAILABLE}. Please try refreshing the page.")

    upload_dir = os.path.join(UPLOAD_DIR, project_id)
    os.makedirs(upload_dir, exist_ok=True)

    uploaded = []
    for file in files:
        content = await file.read()
        file_path = os.path.join(upload_dir, file.filename)
        with open(file_path, "wb") as f:
            f.write(content)

        try:
            csv_text = content.decode("utf-8")
            metadata = CSVParser.parse(csv_text)

            file_info = UploadedFile(
                filename=file.filename,
                size=len(content),
                rows_count=metadata["total_rows"],
                columns=metadata["columns"],
            )
            store.add_uploaded_file(project_id, file_info)
            uploaded.append(file_info)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to parse {file.filename}: {str(e)}")

    store.update_project(project_id, status="uploading")
    return {"files": uploaded, "project_status": "uploading"}


@app.post("/api/projects/{project_id}/research")
async def research_community_endpoint(project_id: str):
    """Scrape community website and extract structured permit/fee/department data using AI."""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    if not project.community_url:
        raise HTTPException(status_code=400, detail="No community URL provided")

    url = project.community_url
    if not url.startswith('http'):
        url = f"https://{url}"

    try:
        # Phase 1: Real web scrape
        print(f"[RESEARCH] Scraping {url} for project {project_id}...")
        scraped_data = scrape_community_website(url, max_pages=35)
        scraped_data["base_url"] = url

        # Cache scrape data in KV
        _kv_set(f"scrape:{project_id}", scraped_data)

        # Phase 2: AI summarization of scraped content
        research = None
        if scraped_data.get("pages_scraped", 0) >= 1:
            print(f"[RESEARCH] Summarizing {scraped_data['pages_scraped']} scraped pages with AI...")
            research = summarize_web_content(scraped_data, community_name=project.customer_name)

        # Fallback to mock only if scrape returned nothing useful AND AI is unavailable
        if not research:
            print(f"[RESEARCH] AI summarization unavailable or no pages found, using mock fallback")
            research = web_researcher.research_community(
                community_url=project.community_url,
                community_name=project.customer_name
            )

        # Save to project
        store.update_project(
            project_id,
            community_name=research.get("community_name", project.customer_name),
            community_research=json.dumps(research)
        )
        return {"status": "complete", "message": f"Scraped {scraped_data.get('pages_scraped', 0)} pages, found {scraped_data.get('pdfs_found', 0)} PDFs", "data": research}
    except Exception as e:
        print(f"[RESEARCH] Error: {e}")
        raise HTTPException(status_code=500, detail=f"Research failed: {str(e)}")


@app.get("/api/projects/{project_id}/research")
async def get_research(project_id: str):
    """Get community research data"""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    if project.community_research:
        try:
            return json.loads(project.community_research)
        except:
            return {"raw": project.community_research}

    return {"status": "no_research", "message": "No community research available yet"}


# ============================================================================
# INTELLIGENCE PIPELINE HELPERS
# ============================================================================

def _match_peer_template(csv_summary: str, community_name: str) -> dict:
    """Auto-match the best peer city template based on CSV data analysis."""
    csv_lower = csv_summary.lower()

    permit_keywords = ["building", "construction", "residential", "commercial", "demolition", "grading", "electrical", "plumbing", "mechanical"]
    license_keywords = ["business", "license", "vendor", "alcohol", "food", "contractor"]
    enforcement_keywords = ["code", "enforcement", "violation", "complaint", "nuisance"]
    land_keywords = ["subdivision", "zoning", "variance", "conditional use", "land use", "planning"]

    permit_count = sum(1 for k in permit_keywords if k in csv_lower)
    license_count = sum(1 for k in license_keywords if k in csv_lower)
    enforcement_count = sum(1 for k in enforcement_keywords if k in csv_lower)
    land_count = sum(1 for k in land_keywords if k in csv_lower)

    total_complexity = permit_count + license_count + enforcement_count + land_count

    if land_count >= 2:
        return next((t for t in PEER_CITY_TEMPLATES if t["id"] == "county-planning"), PEER_CITY_TEMPLATES[1])
    elif total_complexity >= 6:
        return next((t for t in PEER_CITY_TEMPLATES if t["id"] == "mid-city-full"), PEER_CITY_TEMPLATES[1])
    else:
        return next((t for t in PEER_CITY_TEMPLATES if t["id"] == "small-town-basic"), PEER_CITY_TEMPLATES[0])


def _build_intelligence_context(csv_summary: str, community_context: str, matched_template: dict) -> str:
    """Build comprehensive intelligence context for AI analysis."""
    template_context = f"""
## Matched Peer City Template: {matched_template['name']}
Description: {matched_template['description']}
Population Range: {matched_template['population']}

### Template Record Types:
"""
    for rt in matched_template.get("record_types", []):
        template_context += f"\n- **{rt['name']}** ({rt.get('category', 'General')})"
        if rt.get("fees"):
            template_context += f"\n  Fees: {', '.join(f['name'] + ' $' + str(f['amount']) for f in rt['fees'])}"
        if rt.get("workflow_steps"):
            template_context += f"\n  Workflow: {' → '.join(s['name'] for s in rt['workflow_steps'])}"
        if rt.get("form_fields"):
            template_context += f"\n  Fields: {len(rt['form_fields'])} fields including {', '.join(f['name'] for f in rt['form_fields'][:5])}"
        if rt.get("required_documents"):
            template_context += f"\n  Documents: {', '.join(d['name'] for d in rt['required_documents'])}"

    template_context += f"\n\n### Template Departments:\n"
    for dept in matched_template.get("departments", []):
        template_context += f"- {dept['name']}: {dept.get('description', '')}\n"

    template_context += f"\n### Template User Roles:\n"
    for role in matched_template.get("roles", []):
        template_context += f"- {role['name']}: {role.get('description', '')}\n"

    return template_context


def _compute_agent_stats(configuration) -> dict:
    """Compute per-agent domain statistics from configuration"""
    stats = {
        "forms_agent": {"total_form_fields": 0, "record_types_with_forms": 0, "avg_fields_per_type": 0},
        "fees_agent": {"total_fees": 0, "total_fee_amount": 0.0, "record_types_with_fees": 0, "avg_fees_per_type": 0},
        "workflows_agent": {"total_workflow_steps": 0, "record_types_with_workflows": 0, "avg_steps_per_type": 0},
        "documents_agent": {"total_documents": 0, "record_types_with_docs": 0, "avg_docs_per_type": 0},
        "internal_agent": {"total_departments": 0, "total_user_roles": 0},
    }

    if not configuration or not hasattr(configuration, 'record_types'):
        return stats

    for rt in configuration.record_types:
        if hasattr(rt, 'form_fields') and rt.form_fields:
            stats["forms_agent"]["total_form_fields"] += len(rt.form_fields)
            stats["forms_agent"]["record_types_with_forms"] += 1
        if hasattr(rt, 'fees') and rt.fees:
            stats["fees_agent"]["total_fees"] += len(rt.fees)
            stats["fees_agent"]["record_types_with_fees"] += 1
            stats["fees_agent"]["total_fee_amount"] += sum(f.amount for f in rt.fees if hasattr(f, 'amount'))
        if hasattr(rt, 'workflow_steps') and rt.workflow_steps:
            stats["workflows_agent"]["total_workflow_steps"] += len(rt.workflow_steps)
            stats["workflows_agent"]["record_types_with_workflows"] += 1
        if hasattr(rt, 'required_documents') and rt.required_documents:
            stats["documents_agent"]["total_documents"] += len(rt.required_documents)
            stats["documents_agent"]["record_types_with_docs"] += 1

    # Compute averages
    if stats["forms_agent"]["record_types_with_forms"] > 0:
        stats["forms_agent"]["avg_fields_per_type"] = round(stats["forms_agent"]["total_form_fields"] / stats["forms_agent"]["record_types_with_forms"], 1)
    if stats["fees_agent"]["record_types_with_fees"] > 0:
        stats["fees_agent"]["avg_fees_per_type"] = round(stats["fees_agent"]["total_fees"] / stats["fees_agent"]["record_types_with_fees"], 1)
    if stats["workflows_agent"]["record_types_with_workflows"] > 0:
        stats["workflows_agent"]["avg_steps_per_type"] = round(stats["workflows_agent"]["total_workflow_steps"] / stats["workflows_agent"]["record_types_with_workflows"], 1)
    if stats["documents_agent"]["record_types_with_docs"] > 0:
        stats["documents_agent"]["avg_docs_per_type"] = round(stats["documents_agent"]["total_documents"] / stats["documents_agent"]["record_types_with_docs"], 1)

    # Internal agent
    if hasattr(configuration, 'departments'):
        stats["internal_agent"]["total_departments"] = len(configuration.departments)
    if hasattr(configuration, 'user_roles'):
        stats["internal_agent"]["total_user_roles"] = len(configuration.user_roles)

    return stats


def _build_intelligence_report(csv_summary: str, community_context: str, matched_template: dict, configuration) -> dict:
    """Generate an intelligence report summarizing what was analyzed and used."""
    report = {
        "sources_used": [],
        "matched_template": {
            "id": matched_template["id"],
            "name": matched_template["name"],
            "description": matched_template["description"],
            "population": matched_template["population"],
        },
        "analysis_summary": "",
        "completeness_score": 0,
        "auto_enhancements": [],
        "recommendations": [],
        "ai_usage_stats": ai_usage_tracker.get_stats(),
        "agent_stats": {},
    }

    report["sources_used"].append({
        "type": "csv_data",
        "name": "Uploaded CSV Files",
        "status": "analyzed",
        "description": "Historical permit/license/enforcement records parsed and analyzed"
    })

    if community_context:
        report["sources_used"].append({
            "type": "community_research",
            "name": "Community Website Research",
            "status": "analyzed",
            "description": "Government website scraped for department info, services, and contact details"
        })

    report["sources_used"].append({
        "type": "peer_template",
        "name": f"Peer Template: {matched_template['name']}",
        "status": "applied",
        "description": f"Best-match template for {matched_template['population']} population jurisdictions"
    })

    report["sources_used"].append({
        "type": "ai_best_practices",
        "name": "AI Best Practices Engine",
        "status": "applied",
        "description": "Industry best practices for PLC configuration including standard fees, workflows, and document requirements"
    })

    if configuration:
        rt_count = len(configuration.record_types) if hasattr(configuration, 'record_types') else 0
        dept_count = len(configuration.departments) if hasattr(configuration, 'departments') else 0
        role_count = len(configuration.user_roles) if hasattr(configuration, 'user_roles') else 0

        score = 0
        if rt_count >= 3: score += 25
        elif rt_count >= 1: score += 15
        if dept_count >= 2: score += 20
        elif dept_count >= 1: score += 10
        if role_count >= 3: score += 15
        elif role_count >= 1: score += 8

        has_fees = any(len(rt.fees) > 0 for rt in configuration.record_types) if hasattr(configuration, 'record_types') else False
        has_workflows = any(len(rt.workflow_steps) > 0 for rt in configuration.record_types) if hasattr(configuration, 'record_types') else False
        has_docs = any(len(rt.required_documents) > 0 for rt in configuration.record_types) if hasattr(configuration, 'record_types') else False
        has_fields = any(len(rt.form_fields) > 0 for rt in configuration.record_types) if hasattr(configuration, 'record_types') else False

        if has_fees: score += 10
        if has_workflows: score += 10
        if has_docs: score += 10
        if has_fields: score += 10

        report["completeness_score"] = min(score, 100)

        # Count total fees across all record types
        total_fees = sum(len(rt.fees) for rt in configuration.record_types) if hasattr(configuration, 'record_types') else 0

        report["config_overview"] = {
            "record_types_count": rt_count,
            "departments_count": dept_count,
            "user_roles_count": role_count,
            "total_fees": total_fees,
            "has_workflows": has_workflows,
            "has_documents": has_docs,
            "has_form_fields": has_fields,
        }

        report["auto_enhancements"] = []
        if has_fees:
            report["auto_enhancements"].append({"title": "Fee Structures Generated", "description": f"{total_fees} fee items created across all record types based on industry standards"})
        if has_workflows:
            report["auto_enhancements"].append({"title": "Workflow Steps Configured", "description": "Complete approval workflows generated for each record type"})
        if has_docs:
            report["auto_enhancements"].append({"title": "Required Documents Defined", "description": "Document checklists created for each permit/license type"})
        if has_fields:
            report["auto_enhancements"].append({"title": "Form Fields Built", "description": "Application forms designed with all necessary fields for each record type"})
        if dept_count >= 2:
            report["auto_enhancements"].append({"title": "Department Structure", "description": f"{dept_count} departments configured with appropriate assignments"})

        report["analysis_summary"] = f"Generated {rt_count} record types, {dept_count} departments, and {role_count} user roles with {total_fees} fee items. Configuration includes {'fees, ' if has_fees else ''}{'workflows, ' if has_workflows else ''}{'documents, ' if has_docs else ''}{'form fields' if has_fields else ''}."

    report["agent_stats"] = _compute_agent_stats(configuration)
    report["ai_mode"] = claude_service.last_mode  # "ai" or "mock"
    report["ai_error"] = claude_service.last_error
    return report


@app.post("/api/projects/{project_id}/analyze")
async def analyze_project(project_id: str):
    """Unified analysis pipeline: scrape → summarize → parse CSV → AI config → intelligence report"""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    if not project.uploaded_files:
        raise HTTPException(status_code=400, detail="No files uploaded")

    scrape_stats = None  # Track scrape results for intelligence report

    try:
        # ================================================================
        # PHASE 1: Parse CSV files
        # ================================================================
        store.update_project(project_id, status="analyzing", analysis_progress=5, analysis_stage="Parsing uploaded CSV files...")

        upload_dir = os.path.join(UPLOAD_DIR, project_id)
        all_csv_data = ""

        if project.uploaded_files:
            for f in project.uploaded_files:
                file_path = os.path.join(upload_dir, f.filename)
                if os.path.exists(file_path):
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as fp:
                        content = fp.read()
                        metadata = CSVParser.parse(content)
                        all_csv_data += f"\n--- File: {f.filename} ---\n"
                        all_csv_data += CSVParser.to_summary_string(metadata)

        # ================================================================
        # PHASE 2: Web Scrape (if community URL exists)
        # ================================================================
        community_context = ""
        research = None

        if project.community_url:
            url = project.community_url
            if not url.startswith('http'):
                url = f"https://{url}"

            # Check for cached scrape data in KV
            store.update_project(project_id, analysis_progress=10, analysis_stage="Scraping community website...")
            cached_scrape = _kv_get(f"scrape:{project_id}")

            if cached_scrape and cached_scrape.get("pages_scraped", 0) >= 3:
                print(f"[ANALYZE] Using cached scrape data for project {project_id} ({cached_scrape.get('pages_scraped', 0)} pages)")
                scraped_data = cached_scrape
            else:
                print(f"[ANALYZE] Scraping {url} for project {project_id}...")
                scraped_data = scrape_community_website(url, max_pages=35)
                scraped_data["base_url"] = url
                # Cache in KV for reuse
                _kv_set(f"scrape:{project_id}", scraped_data)

            scrape_stats = {
                "pages_scraped": scraped_data.get("pages_scraped", 0),
                "pdfs_found": scraped_data.get("pdfs_found", 0),
                "urls_visited": scraped_data.get("urls_visited", 0),
                "total_chars": scraped_data.get("total_chars", 0),
                "scraped_at": scraped_data.get("scraped_at", datetime.utcnow().isoformat()),
            }

            store.update_project(
                project_id,
                analysis_progress=25,
                analysis_stage=f"Scraped {scrape_stats['pages_scraped']} pages, {scrape_stats['pdfs_found']} PDFs found..."
            )

            # ================================================================
            # PHASE 3: AI Summarize Scraped Content (Two-Pass)
            # ================================================================
            store.update_project(project_id, analysis_progress=30, analysis_stage="Extracting permits, fees & processes from website...")

            if scraped_data.get("pages_scraped", 0) >= 1:
                research = summarize_web_content(scraped_data, community_name=project.customer_name)

            # If AI summarization failed but we have raw scraped text, use it directly
            if not research and scraped_data.get("combined_text"):
                print(f"[ANALYZE] AI summarization unavailable, using raw scraped text as context")
                research = {
                    "community_name": project.customer_name or "the community",
                    "website_url": url,
                    "research_summary": f"Scraped {scraped_data.get('pages_scraped', 0)} pages from the community website. AI summarization unavailable — raw content passed to analysis.",
                    "permits_found": [],
                    "fee_schedule": [],
                    "departments": [],
                    "ordinances": [],
                    "processes": [],
                    "documents_commonly_required": [],
                    "data_source": "raw_scrape",
                    "scrape_stats": {
                        "pages_scraped": scraped_data.get("pages_scraped", 0),
                        "pdfs_found": scraped_data.get("pdfs_found", 0),
                    },
                }

            if research:
                community_context = web_researcher.format_for_analysis(research)
                # Also append raw scraped highlights so the final AI gets real website text
                raw_text = scraped_data.get("combined_text", "")
                if raw_text:
                    community_context += f"\n\n### Raw Website Content Highlights:\n{raw_text[:15000]}"
            else:
                community_context = f"Scraped website but found no relevant content. Raw text:\n{scraped_data.get('combined_text', '')[:10000]}"

            store.update_project(
                project_id,
                community_name=research.get("community_name", project.customer_name),
                community_research=json.dumps(research),
                analysis_progress=40,
                analysis_stage="Website analysis complete..."
            )

        # ================================================================
        # PHASE 4: Match peer city template
        # ================================================================
        store.update_project(project_id, analysis_progress=45, analysis_stage="Matching peer city templates...")
        matched_template = _match_peer_template(all_csv_data, project.customer_name)

        # ================================================================
        # PHASE 5: Build unified context & AI Analysis
        # ================================================================
        store.update_project(project_id, analysis_progress=50, analysis_stage="AI generating record types & workflows...")
        template_context = _build_intelligence_context(all_csv_data, community_context, matched_template)

        # Build the combined context from ALL sources
        web_source_note = ""
        if scrape_stats:
            web_source_note = f"(Extracted from real website scrape of {scrape_stats['pages_scraped']} pages)"

        combined_data = _sanitize_to_ascii(f"""## SOURCE 1: Uploaded CSV Data (Primary Source)
{all_csv_data}

## SOURCE 2: Community Website Research {web_source_note}
{community_context if community_context else "No community URL provided - using best practices for configuration."}

## SOURCE 3: Peer City Reference Data
{template_context}

## SOURCE 4: Industry Best Practices
Use your knowledge of PLC implementations across hundreds of municipalities to fill in gaps and ensure completeness.
Every record type should have comprehensive form fields, realistic fees, complete workflows, and appropriate document requirements.
Do not just map what is in the CSV — build a COMPLETE configuration that this municipality can use immediately.
""")

        store.update_project(project_id, analysis_progress=60, analysis_stage="Building departments & user roles...")
        configuration = claude_service.analyze_csv_data(combined_data)

        # ================================================================
        # PHASE 6: Generate intelligence report with source attribution
        # ================================================================
        store.update_project(project_id, analysis_progress=85, analysis_stage="Finalizing configuration...")
        intel_report = _build_intelligence_report(all_csv_data, community_context, matched_template, configuration)

        # Add scrape stats to intelligence report
        if scrape_stats:
            intel_report["sources_used"].append({
                "type": "web_scrape",
                "name": f"Website Scrape: {project.community_url}",
                "status": "analyzed",
                "description": f"Crawled {scrape_stats['urls_visited']} URLs, found {scrape_stats['pages_scraped']} relevant pages and {scrape_stats['pdfs_found']} PDFs"
            })
            intel_report["scrape_stats"] = scrape_stats
            # Source breakdown
            intel_report["source_breakdown"] = {
                "web_scrape": {
                    "pages_scraped": scrape_stats["pages_scraped"],
                    "pdfs_found": scrape_stats["pdfs_found"],
                    "permits_found": len(research.get("permits_found", [])) if research else 0,
                    "fees_found": len(research.get("fee_schedule", [])) if research else 0,
                    "departments_found": len(research.get("departments", [])) if research else 0,
                },
                "csv_data": {
                    "files_analyzed": len(project.uploaded_files),
                    "total_data_chars": len(all_csv_data),
                },
                "peer_template": {
                    "template_name": matched_template.get("name", ""),
                    "template_id": matched_template.get("id", ""),
                },
            }

        # ================================================================
        # PHASE 7: Save everything
        # ================================================================
        store.save_configuration(project_id, configuration)
        store.update_project(
            project_id,
            status="configured",
            analysis_progress=100,
            analysis_stage="Complete",
            intelligence_report=json.dumps(intel_report)
        )

        return {
            "status": "configured",
            "message": "Analysis complete",
            "intelligence": intel_report
        }

    except Exception as e:
        error_msg = f"Error: {str(e)}"
        store.update_project(project_id, status="error", analysis_stage=error_msg)
        raise HTTPException(status_code=500, detail=error_msg)


@app.get("/api/projects/{project_id}/analysis-status")
async def get_analysis_status(project_id: str):
    """Get current analysis status and progress"""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    return {
        "status": project.status,
        "progress": project.analysis_progress,
        "stage": project.analysis_stage,
    }


@app.get("/api/projects/{project_id}/scrape-status")
async def get_scrape_status(project_id: str):
    """Get the scrape status for a project — whether the website has been scraped and cached."""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    cached_scrape = _kv_get(f"scrape:{project_id}")
    if cached_scrape:
        return {
            "scraped": True,
            "pages_scraped": cached_scrape.get("pages_scraped", 0),
            "pdfs_found": cached_scrape.get("pdfs_found", 0),
            "urls_visited": cached_scrape.get("urls_visited", 0),
            "total_chars": cached_scrape.get("total_chars", 0),
            "scraped_at": cached_scrape.get("scraped_at", ""),
            "pdfs": [{"url": p.get("url", ""), "filename": p.get("filename", "")} for p in cached_scrape.get("pdfs", [])[:20]],
        }
    return {
        "scraped": False,
        "pages_scraped": 0,
        "pdfs_found": 0,
        "urls_visited": 0,
        "total_chars": 0,
        "scraped_at": "",
        "pdfs": [],
    }


@app.get("/api/projects/{project_id}/sample-csv")
async def get_sample_csv(project_id: str):
    """Return a sample CSV for demo purposes"""
    csv_content = MockGenerator.generate_sample_csv()
    return {"content": csv_content, "filename": "sample_permits_data.csv"}


# ============================================================================
# CONFIGURATION ROUTES
# ============================================================================

@app.get("/api/projects/{project_id}/configurations")
async def get_configurations(project_id: str):
    """Get the configuration for a project"""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    if not project.configuration:
        raise HTTPException(status_code=404, detail="No configuration generated yet")
    return project.configuration


@app.put("/api/projects/{project_id}/configurations/record-types/{rt_id}")
async def update_record_type(project_id: str, rt_id: str, request: UpdateRecordTypeRequest):
    """Update a record type in the configuration"""
    project = store.get_project(project_id)
    if not project or not project.configuration:
        raise HTTPException(status_code=404, detail="Project or configuration not found")

    updates = request.model_dump(exclude_none=True)
    updated = store.update_record_type(project_id, rt_id, updates)
    if not updated:
        raise HTTPException(status_code=404, detail="Record type not found")
    return updated


@app.post("/api/projects/{project_id}/configurations/record-types")
async def add_record_type(project_id: str, record_type: RecordType):
    """Add a new record type to the configuration"""
    project = store.get_project(project_id)
    if not project or not project.configuration:
        raise HTTPException(status_code=404, detail="Project or configuration not found")

    result = store.add_record_type(project_id, record_type)
    return result


@app.delete("/api/projects/{project_id}/configurations/record-types/{rt_id}")
async def delete_record_type(project_id: str, rt_id: str):
    """Delete a record type from the configuration"""
    if not store.delete_record_type(project_id, rt_id):
        raise HTTPException(status_code=404, detail="Record type not found")
    return {"success": True}


@app.put("/api/projects/{project_id}/configurations/departments/{dept_id}")
async def update_department(project_id: str, dept_id: str, request: UpdateDepartmentRequest):
    """Update a department in the configuration"""
    project = store.get_project(project_id)
    if not project or not project.configuration:
        raise HTTPException(status_code=404, detail="Project or configuration not found")

    updates = request.model_dump(exclude_none=True)
    updated = store.update_department(project_id, dept_id, updates)
    if not updated:
        raise HTTPException(status_code=404, detail="Department not found")
    return updated


@app.put("/api/projects/{project_id}/configurations/roles/{role_id}")
async def update_role(project_id: str, role_id: str, request: UpdateRoleRequest):
    """Update a user role in the configuration"""
    project = store.get_project(project_id)
    if not project or not project.configuration:
        raise HTTPException(status_code=404, detail="Project or configuration not found")

    updates = request.model_dump(exclude_none=True)
    updated = store.update_role(project_id, role_id, updates)
    if not updated:
        raise HTTPException(status_code=404, detail="Role not found")
    return updated


@app.post("/api/projects/{project_id}/configurations/deploy")
async def deploy_configuration(project_id: str):
    """Placeholder for PLC API deployment"""
    project = store.get_project(project_id)
    if not project or not project.configuration:
        raise HTTPException(status_code=404, detail="Project or configuration not found")

    return {
        "success": False,
        "message": "PLC API integration is not yet configured. This is a placeholder for future deployment functionality.",
        "record_types_count": len(project.configuration.record_types),
        "departments_count": len(project.configuration.departments),
        "roles_count": len(project.configuration.user_roles),
    }


# ============================================================================
# HEALTH CHECK
# ============================================================================

@app.get("/api/health")
async def health_check():
    kv_status = "connected" if KV_AVAILABLE else "not configured"
    kv_test = False
    kv_error = None
    if KV_AVAILABLE:
        try:
            _kv_set("_health_check", {"ts": datetime.utcnow().isoformat()})
            result = _kv_get("_health_check")
            kv_test = result is not None
        except Exception as e:
            kv_error = str(e)

    # Check which env vars are present (masked for security)
    env_check = {
        "KV_REST_API_URL": bool(os.environ.get("KV_REST_API_URL")),
        "KV_REST_API_TOKEN": bool(os.environ.get("KV_REST_API_TOKEN")),
        "UPSTASH_REDIS_REST_URL": bool(os.environ.get("UPSTASH_REDIS_REST_URL")),
        "UPSTASH_REDIS_REST_TOKEN": bool(os.environ.get("UPSTASH_REDIS_REST_TOKEN")),
        "ANTHROPIC_API_KEY": bool(os.environ.get("ANTHROPIC_API_KEY")),
    }

    project_count = len(store._projects)
    # Also count KV projects
    kv_project_count = 0
    if KV_AVAILABLE:
        try:
            kv_keys = _kv_keys("project:*")
            kv_project_count = len(kv_keys)
        except Exception:
            pass

    return {
        "status": "healthy",
        "storage": {
            "kv": kv_status,
            "kv_working": kv_test,
            "kv_error": kv_error,
            "projects_in_memory": project_count,
            "projects_in_kv": kv_project_count,
        },
        "env_vars": env_check,
        "ai": claude_service.get_status()
    }


# ============================================================================
# INTELLIGENCE REPORT ENDPOINTS
# ============================================================================

@app.get("/api/projects/{project_id}/intelligence")
async def get_intelligence_report(project_id: str):
    """Get the auto-generated intelligence report for a project"""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    if project.intelligence_report:
        try:
            report = json.loads(project.intelligence_report)
        except (json.JSONDecodeError, TypeError):
            report = None
    else:
        report = None

    if not report:
        return {
            "status": "not_available",
            "message": "Intelligence analysis has not been run yet. Upload files and run analysis first."
        }

    return {
        "status": "available",
        "report": report
    }


@app.post("/api/projects/{project_id}/re-analyze")
async def re_analyze_with_context(project_id: str, data: dict = {}):
    """Re-run analysis with additional user-provided context"""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    additional_context = data.get("additional_context", "")

    # Get existing CSV data
    upload_dir = os.path.join(UPLOAD_DIR, project_id)
    all_csv_data = ""
    if project.uploaded_files:
        for f in project.uploaded_files:
            file_path = os.path.join(upload_dir, f.filename)
            if os.path.exists(file_path):
                with open(file_path, "r") as fp:
                    content = fp.read()
                    metadata = CSVParser.parse(content)
                    all_csv_data += f"\n--- File: {f.filename} ---\n"
                    all_csv_data += CSVParser.to_summary_string(metadata)

    # Get community research
    community_context = ""
    if project.community_research:
        try:
            research = json.loads(project.community_research)
            community_context = web_researcher.format_for_analysis(research)
        except Exception:
            pass

    # Match template
    matched_template = _match_peer_template(all_csv_data, project.customer_name)
    template_context = _build_intelligence_context(all_csv_data, community_context, matched_template)

    combined_data = f"""## SOURCE 1: Uploaded CSV Data
{all_csv_data}

## SOURCE 2: Community Research
{community_context if community_context else "Not available."}

## SOURCE 3: Peer City Reference
{template_context}

## SOURCE 4: Additional Context from User
{additional_context if additional_context else "None provided."}

## SOURCE 5: Industry Best Practices
Use your knowledge to build the most complete configuration possible.
"""

    store.update_project(project_id, status="analyzing", analysis_progress=50, analysis_stage="Re-analyzing with additional context...")

    try:
        configuration = claude_service.analyze_csv_data(combined_data)
        intel_report = _build_intelligence_report(all_csv_data, community_context, matched_template, configuration)

        store.save_configuration(project_id, configuration)
        store.update_project(
            project_id,
            status="configured",
            analysis_progress=100,
            analysis_stage="Complete",
            intelligence_report=json.dumps(intel_report)
        )

        return {"status": "configured", "message": "Re-analysis complete", "intelligence": intel_report}
    except Exception as e:
        store.update_project(project_id, status="error", analysis_stage=f"Error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


def scrape_community_website(base_url: str, max_pages: int = 35) -> dict:
    """Scrape a government website comprehensively for permit/license/fee content.

    Returns structured scrape data including pages, PDFs, and combined text.
    Used by the unified analysis pipeline.
    """
    import urllib.request
    import urllib.parse
    import time as _time
    from html.parser import HTMLParser

    scraped_pages = []
    pdf_links = []  # [{url, filename, found_on}]
    visited = set()
    # Priority queue: (priority_score, url) — lower score = higher priority
    priority_queue = [(0, base_url)]

    # Simple HTML parser to extract links and text
    class LinkTextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.links = []
            self.text_parts = []
            self.title = ""
            self._in_body = False
            self._in_title = False
            self._skip_tags = {'script', 'style', 'nav', 'footer', 'header', 'noscript'}
            self._skip_depth = 0

        def handle_starttag(self, tag, attrs):
            if tag in self._skip_tags:
                self._skip_depth += 1
            if tag == 'body':
                self._in_body = True
            if tag == 'title':
                self._in_title = True
            attrs_dict = dict(attrs)
            href = attrs_dict.get('href', '')
            if href and not href.startswith('#') and not href.startswith('javascript:') and not href.startswith('mailto:'):
                self.links.append(href)

        def handle_endtag(self, tag):
            if tag in self._skip_tags and self._skip_depth > 0:
                self._skip_depth -= 1
            if tag == 'title':
                self._in_title = False

        def handle_data(self, data):
            if self._in_title:
                self.title += data.strip()
            if self._in_body and self._skip_depth == 0:
                text = data.strip()
                if text:
                    self.text_parts.append(text)

    # Keywords that indicate permit/license content — weighted for URL prioritization
    HIGH_PRIORITY_URL_KEYWORDS = [
        'permit', 'license', 'fee', 'schedule', 'application', 'form',
        'building', 'planning', 'zoning', 'code-enforcement', 'inspection',
        'ordinance', 'municipal-code', 'development-services'
    ]
    MEDIUM_PRIORITY_URL_KEYWORDS = [
        'department', 'public-works', 'services', 'community-development',
        'engineering', 'fire-prevention', 'business', 'contractor',
        'requirements', 'submittal', 'review', 'approval'
    ]
    CONTENT_RELEVANCE_KEYWORDS = [
        'permit', 'license', 'application', 'fee', 'schedule', 'ordinance',
        'building', 'planning', 'zoning', 'code', 'enforcement', 'inspection',
        'department', 'public works', 'development', 'services', 'forms',
        'requirements', 'submittal', 'review', 'approval', 'contractor',
        'business', 'occupation', 'sign', 'demolition', 'grading',
        'electrical', 'plumbing', 'mechanical', 'fire', 'safety',
        'conditional use', 'variance', 'subdivision', 'encroachment',
        'right-of-way', 'impact fee', 'plan check', 'certificate'
    ]

    def _url_priority(link_url: str) -> int:
        """Score a URL for crawl priority (lower = higher priority)."""
        lower = link_url.lower()
        if any(kw in lower for kw in HIGH_PRIORITY_URL_KEYWORDS):
            return 1
        if any(kw in lower for kw in MEDIUM_PRIORITY_URL_KEYWORDS):
            return 2
        return 3

    all_text_content = []
    pages_analyzed = 0
    parsed_base = urllib.parse.urlparse(base_url)

    print(f"[SCRAPE] Starting crawl of {base_url} (max {max_pages} pages)")

    while priority_queue and pages_analyzed < max_pages:
        # Sort by priority and take the best
        priority_queue.sort(key=lambda x: x[0])
        _priority, current_url = priority_queue.pop(0)

        # Normalize URL
        if not current_url.startswith('http'):
            current_url = urllib.parse.urljoin(base_url, current_url)

        # Remove fragments and normalize
        parsed_current = urllib.parse.urlparse(current_url)
        current_url = urllib.parse.urlunparse(parsed_current._replace(fragment=''))

        # Skip if already visited or external
        if current_url in visited:
            continue
        if parsed_current.netloc and parsed_current.netloc != parsed_base.netloc:
            continue
        # Skip non-web resources
        if any(current_url.lower().endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif', '.svg', '.css', '.js', '.ico', '.woff', '.woff2', '.ttf', '.mp4', '.mp3', '.zip']):
            continue

        visited.add(current_url)

        try:
            req = urllib.request.Request(current_url, headers={
                'User-Agent': 'Mozilla/5.0 (compatible; OpenGov-AutoConfig/1.0)'
            })
            with urllib.request.urlopen(req, timeout=10) as resp:
                content_type = resp.headers.get('Content-Type', '')

                if 'pdf' in content_type.lower():
                    filename = current_url.split('/')[-1] or 'document.pdf'
                    pdf_links.append({"url": current_url, "filename": filename, "found_on": "direct"})
                    continue

                if 'html' not in content_type.lower() and 'text' not in content_type.lower():
                    continue

                html = resp.read().decode('utf-8', errors='ignore')

                parser = LinkTextExtractor()
                try:
                    parser.feed(html)
                except Exception:
                    continue

                page_text = _sanitize_to_ascii(' '.join(parser.text_parts))
                page_title = _sanitize_to_ascii(parser.title.strip()) or current_url

                # Check if page has relevant content
                page_lower = page_text.lower()
                relevance_score = sum(1 for kw in CONTENT_RELEVANCE_KEYWORDS if kw in page_lower)

                # Always store the page if it has some relevance (lower threshold for broader coverage)
                if relevance_score >= 1:
                    scraped_pages.append({
                        "url": current_url,
                        "title": page_title,
                        "relevance": relevance_score,
                        "text_length": len(page_text),
                        "text": page_text[:4000]  # Store more text per page
                    })
                    all_text_content.append(
                        f"--- PAGE: {page_title} ({current_url}) [relevance: {relevance_score}] ---\n{page_text[:4000]}\n"
                    )

                # Discover new links
                for link in parser.links:
                    full_link = urllib.parse.urljoin(current_url, link)
                    # Remove fragment
                    full_link = urllib.parse.urlunparse(urllib.parse.urlparse(full_link)._replace(fragment=''))

                    if full_link in visited:
                        continue

                    if full_link.lower().endswith('.pdf'):
                        filename = full_link.split('/')[-1] or 'document.pdf'
                        pdf_links.append({"url": full_link, "filename": filename, "found_on": current_url})
                    else:
                        # Check if same domain
                        parsed_link = urllib.parse.urlparse(full_link)
                        if not parsed_link.netloc or parsed_link.netloc == parsed_base.netloc:
                            priority = _url_priority(full_link)
                            priority_queue.append((priority, full_link))

                pages_analyzed += 1

            # Polite crawling: 0.3s delay between requests
            _time.sleep(0.3)

        except Exception as e:
            continue

    # Deduplicate PDF links by URL
    seen_pdfs = set()
    unique_pdfs = []
    for pdf in pdf_links:
        if pdf["url"] not in seen_pdfs:
            seen_pdfs.add(pdf["url"])
            unique_pdfs.append(pdf)

    # Sort scraped pages by relevance (highest first)
    scraped_pages.sort(key=lambda p: p.get("relevance", 0), reverse=True)

    print(f"[SCRAPE] Complete: {len(scraped_pages)} relevant pages, {len(unique_pdfs)} PDFs, {len(visited)} URLs visited")

    return {
        "pages": scraped_pages,
        "pdfs": unique_pdfs,
        "pages_scraped": len(scraped_pages),
        "pdfs_found": len(unique_pdfs),
        "total_chars": sum(len(p.get("text", "")) for p in scraped_pages),
        "total_text_length": sum(len(t) for t in all_text_content),
        "combined_text": '\n\n'.join(all_text_content)[:60000],  # Cap at 60k chars
        "urls_visited": len(visited),
        "scraped_at": datetime.utcnow().isoformat(),
    }


# Keep backward-compatible alias
def _deep_scrape_site(base_url: str, max_pages: int = 20) -> dict:
    """Backward-compatible wrapper around scrape_community_website."""
    result = scrape_community_website(base_url, max_pages)
    # Map new format to old format for existing consumers
    return {
        "pages_scraped": result["pages_scraped"],
        "pdfs_found": result["pdfs_found"],
        "pdf_links": [p["url"] for p in result.get("pdfs", [])][:10],
        "total_text_length": result["total_text_length"],
        "relevant_pages": [{"url": p["url"], "relevance": p["relevance"], "text_length": p["text_length"], "text_preview": p.get("text", "")[:2000]} for p in result.get("pages", [])[:15]],
        "combined_text": result["combined_text"],
        "urls_visited": result["urls_visited"],
    }


@app.post("/api/projects/{project_id}/deep-scrape")
async def deep_scrape_project(project_id: str, data: dict = {}):
    """Deep scrape the project's community website for permit data, then re-analyze."""
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    agent_id = data.get("agent_id", "all")
    url = data.get("url", "") or project.community_url

    if not url:
        raise HTTPException(status_code=400, detail="No community URL configured. Please add a community website URL to the project.")

    # Ensure URL has protocol
    if not url.startswith('http'):
        url = f"https://{url}"

    try:
        # Phase 1: Deep scrape
        store.update_project(project_id, status="analyzing", analysis_progress=10, analysis_stage=f"Deep scraping {url}...")

        scrape_results = _deep_scrape_site(url, max_pages=20)

        store.update_project(project_id, analysis_progress=40, analysis_stage=f"Found {scrape_results['pages_scraped']} relevant pages, {scrape_results['pdfs_found']} PDFs...")

        # Phase 2: AI analysis of scraped content
        if scrape_results['combined_text'] and ANTHROPIC_AVAILABLE:
            store.update_project(project_id, analysis_progress=50, analysis_stage="AI analyzing scraped content...")

            scrape_summary = _extract_with_ai(
                f"""You are an expert in municipal government permit/license processes. Analyze the following content scraped from a government website and extract ALL information relevant to configuring a PLC (Permitting, Licensing, and Code Enforcement) system.

Focus on identifying:
1. All permit and license types offered
2. Fee amounts and fee structures
3. Application requirements and form fields
4. Workflow/approval processes
5. Required documents for each permit type
6. Department names and contact info
7. Inspection requirements
8. Any ordinance references

SCRAPED WEBSITE CONTENT:
{scrape_results['combined_text'][:30000]}

Provide a comprehensive, structured summary of all findings. Be thorough - extract every data point you can find.""",
                "",
                "deep_scrape_analysis"
            )
        else:
            scrape_summary = f"Scraped {scrape_results['pages_scraped']} pages with {scrape_results['total_text_length']} characters of relevant content."

        # Phase 3: Rebuild configuration with new data
        store.update_project(project_id, analysis_progress=65, analysis_stage="Rebuilding configuration with deep scrape data...")

        # Get existing CSV data
        upload_dir = os.path.join(UPLOAD_DIR, project_id)
        all_csv_data = ""
        if project.uploaded_files:
            for f in project.uploaded_files:
                file_path = os.path.join(upload_dir, f.filename)
                if os.path.exists(file_path):
                    with open(file_path, "r") as fp:
                        content = fp.read()
                        metadata = CSVParser.parse(content)
                        all_csv_data += f"\n--- File: {f.filename} ---\n"
                        all_csv_data += CSVParser.to_summary_string(metadata)

        # Get community research
        community_context = ""
        if project.community_research:
            try:
                research = json.loads(project.community_research)
                community_context = web_researcher.format_for_analysis(research)
            except Exception:
                pass

        # Match template
        matched_template = _match_peer_template(all_csv_data, project.customer_name)
        template_context = _build_intelligence_context(all_csv_data, community_context, matched_template)

        # Build enhanced combined context with deep scrape data
        combined_data = f"""## SOURCE 1: Uploaded CSV Data
{all_csv_data if all_csv_data else "No CSV data uploaded."}

## SOURCE 2: Community Research
{community_context if community_context else "Not available."}

## SOURCE 3: Peer City Reference
{template_context}

## SOURCE 4: DEEP WEBSITE SCRAPE (PRIMARY NEW DATA)
The following data was extracted from a deep scrape of {scrape_results['urls_visited']} URLs on the government website.
{scrape_results['pages_scraped']} relevant pages found, {scrape_results['pdfs_found']} PDF documents discovered.

AI Analysis of Scraped Content:
{scrape_summary or 'No AI analysis available.'}

Raw scraped highlights:
{scrape_results['combined_text'][:15000]}

## SOURCE 5: Industry Best Practices
Use your knowledge to fill any remaining gaps and build the most complete configuration possible.
"""

        store.update_project(project_id, analysis_progress=75, analysis_stage="AI generating enhanced configuration...")

        configuration = claude_service.analyze_csv_data(combined_data)

        # Phase 4: Build updated intelligence report
        store.update_project(project_id, analysis_progress=90, analysis_stage="Finalizing intelligence report...")

        intel_report = _build_intelligence_report(all_csv_data, community_context, matched_template, configuration)

        # Add deep scrape source to the report
        intel_report["sources_used"].append({
            "type": "deep_scrape",
            "name": f"Deep Website Scrape: {url}",
            "status": "analyzed",
            "description": f"Crawled {scrape_results['urls_visited']} URLs, found {scrape_results['pages_scraped']} relevant pages and {scrape_results['pdfs_found']} PDFs"
        })

        # Add scrape stats
        intel_report["deep_scrape_stats"] = {
            "urls_visited": scrape_results["urls_visited"],
            "relevant_pages": scrape_results["pages_scraped"],
            "pdfs_found": scrape_results["pdfs_found"],
            "total_text_analyzed": scrape_results["total_text_length"],
            "data_points_extracted": scrape_results["pages_scraped"] * 15,  # Estimated
            "pdf_links": scrape_results.get("pdf_links", []),
            "scrape_summary": scrape_summary[:2000] if scrape_summary else "",
        }

        # Phase 5: Save
        store.save_configuration(project_id, configuration)
        store.update_project(
            project_id,
            status="configured",
            analysis_progress=100,
            analysis_stage="Deep analysis complete",
            intelligence_report=json.dumps(intel_report)
        )

        # Build response matching frontend expectations (AgentDetailDialog.jsx)
        deep_stats = intel_report.get("deep_scrape_stats", {})
        return {
            "status": "complete",
            "message": f"Deep scrape analyzed {scrape_results['pages_scraped']} pages and {scrape_results['pdfs_found']} PDFs. Configuration has been rebuilt.",
            "deep_scrape_stats": {
                "pages_scraped": deep_stats.get("relevant_pages", scrape_results.get("pages_scraped", 0)),
                "pdfs_found": deep_stats.get("pdfs_found", scrape_results.get("pdfs_found", 0)),
                "data_points": deep_stats.get("data_points_extracted", 0),
                "urls_visited": deep_stats.get("urls_visited", 0),
                "total_text_analyzed": deep_stats.get("total_text_analyzed", 0),
            },
            "summary": deep_stats.get("scrape_summary", "") or scrape_summary or f"Analyzed {scrape_results['pages_scraped']} pages and {scrape_results['pdfs_found']} PDFs from the community website.",
            "intelligence": intel_report,
        }

    except Exception as e:
        store.update_project(project_id, status="configured", analysis_stage=f"Deep scrape error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Deep scrape failed: {str(e)}")


# ============================================================================
# LMS CONTENT GENERATION
# ============================================================================

@app.post("/api/projects/{project_id}/lms/generate/{content_type}")
async def generate_lms_content(project_id: str, content_type: str):
    """Generate LMS content for a project configuration.

    content_type: "training-deck", "faq", "quiz", or "handbook"
    Returns: { filename, content_base64, size_bytes, generated_at }
    """
    valid_types = ["training-deck", "faq", "quiz", "handbook"]
    if content_type not in valid_types:
        raise HTTPException(status_code=400, detail=f"Invalid content_type. Must be one of {valid_types}")

    project = store.get_project(project_id)
    if not project or not project.configuration:
        raise HTTPException(status_code=404, detail="Project or configuration not found")

    if not project.configuration.record_types:
        raise HTTPException(status_code=400, detail="Project has no record types configured")

    generated_at = datetime.utcnow().isoformat()

    try:
        if content_type == "training-deck":
            filename, content_bytes = _generate_training_deck(project)
        elif content_type == "faq":
            filename, content_bytes = _generate_faq(project)
        elif content_type == "quiz":
            filename, content_bytes = _generate_quiz(project)
        else:  # handbook
            filename, content_bytes = _generate_handbook(project)
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        raise HTTPException(status_code=500, detail=f"Generation failed ({content_type}): {str(e)[:300]}")

    content_base64 = base64.b64encode(content_bytes).decode("utf-8")

    return {
        "filename": filename,
        "content_base64": content_base64,
        "size_bytes": len(content_bytes),
        "generated_at": generated_at
    }


def _get_community_branding(project):
    """Extract branding info from project for use in generated materials."""
    return {
        "name": project.customer_name or "Community",
        "url": project.community_url or "",
        "primary_color": (26, 86, 219),  # Blue default
        "dark_color": (15, 23, 42),
        "light_color": (241, 245, 249),
        "accent_color": (16, 185, 129),
    }


def _safe_text(text, max_len=500):
    """Sanitize text for FPDF - remove unicode chars that cause issues."""
    if not text:
        return ""
    cleaned = text[:max_len]
    # Replace common unicode chars with ASCII equivalents
    replacements = {
        '\u2013': '-', '\u2014': '--', '\u2018': "'", '\u2019': "'",
        '\u201c': '"', '\u201d': '"', '\u2022': '-', '\u2026': '...',
        '\u00ae': '(R)', '\u00a9': '(c)', '\u2122': '(TM)',
        '\u00b0': ' deg', '\u00bd': '1/2', '\u00bc': '1/4',
        '\u00be': '3/4', '\u00e9': 'e', '\u00e8': 'e', '\u00f1': 'n',
    }
    for old, new in replacements.items():
        cleaned = cleaned.replace(old, new)
    return cleaned.encode('latin-1', errors='replace').decode('latin-1')


def _generate_training_deck(project) -> tuple:
    """Generate a professional PowerPoint training deck (.pptx)"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    brand = _get_community_branding(project)
    config = project.configuration
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    PRIMARY = RGBColor(*brand["primary_color"])
    DARK = RGBColor(*brand["dark_color"])
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252)
    ACCENT = RGBColor(*brand["accent_color"])
    GRAY = RGBColor(100, 116, 139)

    def _set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_branded_slide(title_text, subtitle_text="", dark=True):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, DARK if dark else WHITE)
        # Accent bar at top
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY
        bar.line.fill.background()
        # Title
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE if dark else DARK
        if subtitle_text:
            sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(1))
            stf = sub.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = subtitle_text
            sp.font.size = Pt(20)
            sp.font.color.rgb = RGBColor(148, 163, 184) if dark else GRAY
        return slide

    def add_content_slide(title_text, bullets, two_col=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, WHITE)
        # Top accent bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.04))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY
        bar.line.fill.background()
        # Title
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK
        # Underline
        line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(2), Inches(0))
        line.line.color.rgb = PRIMARY
        line.line.width = Pt(3)

        if two_col and len(bullets) > 3:
            mid = len(bullets) // 2
            for col, items, x_pos in [(0, bullets[:mid], 0.8), (1, bullets[mid:], 7.0)]:
                box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.8), Inches(5.5), Inches(5))
                tf = box.text_frame
                tf.word_wrap = True
                for i, b in enumerate(items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(16)
                    p.font.color.rgb = DARK
                    p.space_before = Pt(8)
                    p.space_after = Pt(4)
        else:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
            tf = box.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
                p.font.color.rgb = DARK
                p.space_before = Pt(8)
                p.space_after = Pt(4)
        return slide

    # ---- SLIDE 1: Title ----
    add_branded_slide(
        f"{brand['name']}",
        f"OpenGov Implementation - Staff Training & Process Overview"
    )

    # ---- SLIDE 2: Agenda ----
    add_content_slide("Training Agenda", [
        "1. Project Overview & Goals",
        "2. Current vs. New Process Comparison",
        "3. Record Types & Applications",
        "4. Fee Schedules & Payment Processing",
        "5. Workflow & Approval Routing",
        "6. Departments & Staff Roles",
        "7. Required Documents & Submissions",
        "8. System Navigation & Key Features",
        "9. Q&A and Next Steps",
    ])

    # ---- SLIDE 3: Project Overview ----
    overview_bullets = [
        f"Community: {brand['name']}",
        f"Total Record Types: {len(config.record_types)}",
        f"Departments Configured: {len(config.departments)}",
        f"User Roles Defined: {len(config.user_roles)}",
    ]
    if project.community_url:
        overview_bullets.append(f"Community Website: {project.community_url}")
    overview_bullets.extend([
        "",
        "GOAL: Streamline permitting, licensing, and code enforcement",
        "through a centralized digital platform that improves efficiency,",
        "transparency, and the applicant experience."
    ])
    add_content_slide("Project Overview", overview_bullets)

    # ---- SLIDE 4: What's Changing ----
    add_content_slide("What's Changing", [
        "BEFORE: Paper-based or fragmented digital processes",
        "  - Manual routing, lost applications, inconsistent fee collection",
        "  - Limited visibility for applicants on application status",
        "  - Siloed departmental processes with no unified tracking",
        "",
        "AFTER: Unified OpenGov Digital Platform",
        "  - Online submission with automated workflow routing",
        "  - Real-time status tracking for staff and applicants",
        "  - Centralized fee calculation and payment processing",
        "  - Configurable business rules and conditional logic",
    ])

    # ---- SLIDES 5+: Each Record Type ----
    for rt in config.record_types[:8]:
        bullets = []
        if rt.description:
            bullets.append(f"Description: {rt.description[:200]}")
        if rt.category:
            bullets.append(f"Category: {rt.category}")
        if rt.form_fields:
            field_names = [f.name for f in rt.form_fields[:6]]
            bullets.append(f"Key Form Fields: {', '.join(field_names)}")
            required = [f.name for f in rt.form_fields if f.required][:4]
            if required:
                bullets.append(f"Required Fields: {', '.join(required)}")
        if rt.fees:
            for fee in rt.fees[:4]:
                bullets.append(f"Fee: {fee.name} - ${fee.amount:.2f} ({fee.fee_type})")
        if rt.workflow_steps:
            steps = sorted(rt.workflow_steps, key=lambda s: s.order)[:6]
            step_str = " -> ".join([s.name for s in steps])
            bullets.append(f"Workflow: {step_str}")
        if rt.required_documents:
            doc_names = [d.name for d in rt.required_documents[:4]]
            bullets.append(f"Required Documents: {', '.join(doc_names)}")
        add_content_slide(f"Record Type: {rt.name}", bullets)

    # ---- Fee Schedule Summary ----
    fee_bullets = []
    for rt in config.record_types:
        if rt.fees:
            fee_bullets.append(f"--- {rt.name} ---")
            for fee in rt.fees[:5]:
                fee_bullets.append(f"  {fee.name}: ${fee.amount:.2f} ({fee.fee_type})")
    if fee_bullets:
        add_content_slide("Complete Fee Schedule", fee_bullets[:16], two_col=True)

    # ---- Departments ----
    dept_bullets = []
    for d in config.departments[:10]:
        desc = f" - {d.description[:80]}" if d.description else ""
        dept_bullets.append(f"{d.name}{desc}")
    if dept_bullets:
        add_content_slide("Departments & Organization", dept_bullets, two_col=True)

    # ---- User Roles ----
    role_bullets = []
    for r in config.user_roles[:8]:
        desc = f" - {r.description[:80]}" if r.description else ""
        role_bullets.append(f"{r.name}{desc}")
        if r.permissions:
            role_bullets.append(f"  Permissions: {', '.join(r.permissions[:4])}")
    if role_bullets:
        add_content_slide("User Roles & Permissions", role_bullets)

    # ---- Key System Features ----
    add_content_slide("Key System Features", [
        "Online Application Portal - 24/7 submission for applicants",
        "Automated Routing - Applications go to the right department automatically",
        "Real-Time Status - Track applications from submission to approval",
        "Integrated Payments - Online fee calculation and payment processing",
        "Document Management - Upload, review, and store required documents",
        "Inspection Scheduling - Schedule and track field inspections",
        "Reporting & Analytics - Dashboards for workload and performance",
        "Notifications - Automated emails at each workflow stage",
    ])

    # ---- Next Steps ----
    add_branded_slide(
        "Next Steps",
        "1. Review this training material  |  2. Attend hands-on training sessions  |  3. Complete the Knowledge Quiz  |  4. Begin using the new system"
    )

    # ---- Thank You ----
    add_branded_slide(
        "Thank You",
        f"{brand['name']} - OpenGov Implementation Team"
    )

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return "training_deck.pptx", output.read()


def _generate_faq(project) -> tuple:
    """Generate comprehensive FAQ document (.pdf)"""
    from fpdf import FPDF

    brand = _get_community_branding(project)
    config = project.configuration
    pdf = FPDF()
    pdf.set_margins(15, 15, 15)
    pdf.set_auto_page_break(auto=True, margin=20)

    r, g, b = brand["primary_color"]
    dr, dg, db = brand["dark_color"]
    usable_w = pdf.w - pdf.l_margin - pdf.r_margin

    def add_header(text, size=16):
        pdf.set_x(pdf.l_margin)
        pdf.set_font("Arial", "B", size)
        pdf.set_text_color(r, g, b)
        pdf.cell(usable_w, 10, _safe_text(text), ln=True)
        pdf.set_text_color(dr, dg, db)
        pdf.ln(2)

    def add_qa(question, answer):
        pdf.set_x(pdf.l_margin)
        pdf.set_font("Arial", "B", 11)
        pdf.set_text_color(dr, dg, db)
        pdf.multi_cell(usable_w, 6, _safe_text(f"Q: {question}"))
        pdf.set_x(pdf.l_margin)
        pdf.set_font("Arial", "", 10)
        pdf.set_text_color(80, 80, 80)
        pdf.multi_cell(usable_w, 5, _safe_text(f"A: {answer}", 800))
        pdf.ln(4)

    # Cover page
    pdf.add_page()
    pdf.set_font("Arial", "B", 28)
    pdf.set_text_color(r, g, b)
    pdf.ln(30)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 15, _safe_text(brand["name"]), ln=True, align="C")
    pdf.set_font("Arial", "", 18)
    pdf.set_text_color(dr, dg, db)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 10, "Frequently Asked Questions", ln=True, align="C")
    pdf.set_font("Arial", "", 12)
    pdf.set_text_color(100, 100, 100)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 8, "OpenGov Implementation Guide", ln=True, align="C")
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 8, _safe_text(f"Generated: {datetime.utcnow().strftime('%B %d, %Y')}"), ln=True, align="C")

    # Section 1: General
    pdf.add_page()
    add_header("1. General Questions")

    add_qa("What is this new system?",
        f"{brand['name']} is transitioning to a unified OpenGov platform for managing permits, licenses, and code enforcement. "
        "This system replaces previous paper-based or fragmented digital processes with a centralized, online platform.")
    add_qa("Why are we making this change?",
        "The new system provides online submission, automated workflow routing, real-time status tracking, "
        "integrated payment processing, and comprehensive reporting. This improves efficiency for staff and transparency for applicants.")
    add_qa("When will the new system go live?",
        "The go-live date will be communicated by your project manager. Training sessions will be held before launch to ensure all staff are comfortable with the new processes.")
    add_qa("Will my existing applications be transferred?",
        "Your project team will determine the data migration strategy. Active applications may be migrated or completed under the current system depending on their status.")
    add_qa("Who do I contact for help?",
        f"Contact the {brand['name']} Help Desk or your designated OpenGov implementation specialist for technical support and process questions.")

    # Section 2: For Applicants
    add_header("2. For Applicants & Public Users")

    add_qa("How do I submit an application online?",
        f"Visit the {brand['name']} online portal, select the application type you need, complete the required form fields, upload supporting documents, pay applicable fees, and submit. You'll receive a confirmation email with a tracking number.")

    for rt in config.record_types[:5]:
        if rt.form_fields:
            required_fields = [f.name for f in rt.form_fields if f.required][:5]
            if required_fields:
                add_qa(f"What information do I need to apply for a {rt.name}?",
                    f"To apply for a {rt.name}, you'll need to provide: {', '.join(required_fields)}. "
                    + (f"Required documents include: {', '.join([d.name for d in rt.required_documents[:3]])}." if rt.required_documents else ""))
        if rt.fees:
            fee_list = ", ".join([f"{f.name}: ${f.amount:.2f}" for f in rt.fees[:3]])
            add_qa(f"What are the fees for a {rt.name}?",
                f"The fees for {rt.name} are: {fee_list}. Fees can be paid online by credit card or in person.")

    add_qa("How do I check the status of my application?",
        "Log into the portal with your account credentials. Your dashboard shows all submitted applications with their current status, pending actions, and estimated completion dates.")

    # Section 3: For Staff
    pdf.add_page()
    add_header("3. For Staff & Reviewers")

    add_qa("How are applications assigned to me?",
        "Applications are automatically routed to the appropriate department and reviewer based on the record type and configured workflow rules. You'll receive email notifications when new items require your attention.")

    for rt in config.record_types[:3]:
        if rt.workflow_steps:
            steps = sorted(rt.workflow_steps, key=lambda s: s.order)
            step_desc = " -> ".join([s.name for s in steps[:6]])
            add_qa(f"What is the review workflow for {rt.name}?",
                f"The workflow for {rt.name} follows these steps: {step_desc}. "
                f"Each step has designated reviewers and approval requirements.")

    add_qa("Can I send an application back for corrections?",
        "Yes. At any review step, you can request additional information or corrections from the applicant. "
        "The system will notify the applicant and pause the workflow until they respond.")
    add_qa("How do I approve or deny an application?",
        "Open the application from your task queue, review all submitted information and documents, "
        "add any review comments, then click Approve or Deny. The system will automatically route to the next step or notify the applicant of the decision.")

    # Section 4: Departments
    if config.departments:
        add_header("4. Departments & Responsibilities")
        for dept in config.departments[:6]:
            desc = dept.description if dept.description else "Responsible for reviewing and processing applications within their jurisdiction."
            add_qa(f"What does the {dept.name} department handle?", desc)

    # Section 5: Troubleshooting
    add_header("5. Troubleshooting & Technical Support")
    add_qa("I can't log into the system.",
        "Verify your username and password. If you've forgotten your password, use the 'Forgot Password' link on the login page. "
        "Contact the Help Desk if issues persist.")
    add_qa("My application appears stuck in a workflow step.",
        "Check the application's activity log for any pending actions or reviewer comments. "
        "If the application has been pending beyond the expected timeframe, contact your supervisor or the system administrator.")
    add_qa("A fee amount appears incorrect.",
        "Fee calculations are based on the configured fee schedule. If you believe there is an error, "
        "contact the system administrator with the application number and details of the discrepancy.")
    add_qa("I need to upload a large document but it won't go through.",
        "Check the file size limit (typically 25MB per file). Ensure the file format is accepted (PDF, JPG, PNG, DOC). "
        "Try compressing the file or splitting it into smaller documents if needed.")

    return "faq_guide.pdf", bytes(pdf.output())


def _generate_quiz(project) -> tuple:
    """Generate professional knowledge assessment quiz (.pdf)"""
    from fpdf import FPDF

    brand = _get_community_branding(project)
    config = project.configuration
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)

    r, g, b = brand["primary_color"]
    dr, dg, db = brand["dark_color"]

    questions = []
    answer_key = []

    # Generate questions from configuration data
    # Q: Record type count
    rt_count = len(config.record_types)
    questions.append({
        "q": f"How many record types are configured in the {brand['name']} system?",
        "options": [str(rt_count), str(rt_count + 2), str(max(1, rt_count - 1)), str(rt_count + 5)],
        "answer": "A"
    })

    # Q: Per record type - required documents
    for rt in config.record_types[:4]:
        if rt.required_documents and len(rt.required_documents) >= 2:
            correct_doc = rt.required_documents[0].name
            wrong_options = ["Notarized Affidavit", "Birth Certificate", "Vehicle Registration", "Insurance Binder", "Tax Clearance Letter"]
            wrong = [w for w in wrong_options if w.lower() != correct_doc.lower()][:3]
            questions.append({
                "q": f"Which of the following is a required document for a {rt.name} application?",
                "options": [correct_doc, wrong[0], wrong[1], wrong[2]],
                "answer": "A"
            })

        # Q: Fees
        if rt.fees:
            fee = rt.fees[0]
            correct = f"${fee.amount:.2f}"
            wrong_amounts = [f"${fee.amount * 1.5:.2f}", f"${fee.amount * 0.5:.2f}", f"${fee.amount * 2:.2f}"]
            questions.append({
                "q": f"What is the {fee.name} fee for a {rt.name}?",
                "options": [correct, wrong_amounts[0], wrong_amounts[1], wrong_amounts[2]],
                "answer": "A"
            })

        # Q: Workflow steps
        if rt.workflow_steps and len(rt.workflow_steps) >= 2:
            steps = sorted(rt.workflow_steps, key=lambda s: s.order)
            first_step = steps[0].name
            second_step = steps[1].name if len(steps) > 1 else "Review"
            questions.append({
                "q": f"What is the first step in the {rt.name} workflow?",
                "options": [first_step, second_step, "Final Approval", "Payment Collection"],
                "answer": "A"
            })

    # General knowledge questions
    general = [
        {"q": "What is the primary benefit of the new OpenGov system?",
         "options": ["Centralized digital workflow with automated routing", "Faster internet speeds", "Reduced staff headcount", "Elimination of all fees"],
         "answer": "A"},
        {"q": "How are applicants notified of status changes?",
         "options": ["Automated email notifications", "Phone calls from staff", "Physical mail only", "They must check in person"],
         "answer": "A"},
        {"q": "Who can request additional information from an applicant during review?",
         "options": ["The assigned reviewer at any workflow step", "Only the department head", "Only the system administrator", "No one - applications cannot be paused"],
         "answer": "A"},
        {"q": "What should you do if you identify an incorrect fee on an application?",
         "options": ["Contact the system administrator with the application number", "Approve it anyway", "Delete the application", "Ignore it"],
         "answer": "A"},
        {"q": "How are applications routed to the correct department?",
         "options": ["Automatically based on record type and workflow rules", "Manually by the front desk", "Randomly assigned", "Applicants choose their reviewer"],
         "answer": "A"},
    ]
    questions.extend(general)

    # Department questions
    if config.departments and len(config.departments) >= 2:
        dept = config.departments[0]
        wrong_depts = ["Human Resources", "Marketing", "Janitorial Services"]
        questions.append({
            "q": f"Which department is part of the {brand['name']} implementation?",
            "options": [dept.name, wrong_depts[0], wrong_depts[1], wrong_depts[2]],
            "answer": "A"
        })

    # Role questions
    if config.user_roles and len(config.user_roles) >= 2:
        role = config.user_roles[0]
        questions.append({
            "q": f"What is the '{role.name}' role responsible for in the system?",
            "options": [
                role.description[:80] if role.description else "Reviewing and processing applications",
                "Managing server hardware",
                "Designing marketing materials",
                "Maintaining the parking lot"
            ],
            "answer": "A"
        })

    # Shuffle options (but track answer)
    import random
    for q in questions:
        opts = q["options"][:]
        correct_text = opts[0]
        random.shuffle(opts)
        q["options"] = opts
        q["answer"] = chr(65 + opts.index(correct_text))

    # Limit to 20 questions
    questions = questions[:20]

    # --- Build PDF ---
    pdf.set_margins(15, 15, 15)
    usable_w = pdf.w - pdf.l_margin - pdf.r_margin

    # Cover page
    pdf.add_page()
    pdf.set_font("Arial", "B", 28)
    pdf.set_text_color(r, g, b)
    pdf.ln(25)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 15, _safe_text(brand["name"]), ln=True, align="C")
    pdf.set_font("Arial", "", 18)
    pdf.set_text_color(dr, dg, db)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 10, "Knowledge Assessment Quiz", ln=True, align="C")
    pdf.set_font("Arial", "", 12)
    pdf.set_text_color(100, 100, 100)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 8, "OpenGov Implementation Training", ln=True, align="C")
    pdf.ln(15)
    pdf.set_font("Arial", "", 11)
    pdf.set_text_color(80, 80, 80)
    pdf.set_x(pdf.l_margin)
    pdf.multi_cell(usable_w, 6, _safe_text(
        "Instructions: Select the best answer for each question. This quiz covers the key concepts "
        f"from the {brand['name']} OpenGov implementation training. There are {len(questions)} questions total. "
        "An answer key is provided at the end."
    ))
    pdf.ln(5)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 8, "Name: ________________________________________", ln=True)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 8, "Date:  ________________________________________", ln=True)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 8, _safe_text(f"Score: _______ / {len(questions)}"), ln=True)

    # Questions
    pdf.add_page()
    pdf.set_font("Arial", "B", 14)
    pdf.set_text_color(r, g, b)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 10, "Questions", ln=True)
    pdf.ln(2)

    for i, q in enumerate(questions):
        pdf.set_x(pdf.l_margin)
        pdf.set_font("Arial", "B", 11)
        pdf.set_text_color(dr, dg, db)
        pdf.multi_cell(usable_w, 6, _safe_text(f"{i+1}. {q['q']}"))
        pdf.set_font("Arial", "", 10)
        pdf.set_text_color(80, 80, 80)
        for j, opt in enumerate(q["options"]):
            letter = chr(65 + j)
            pdf.set_x(pdf.l_margin)
            pdf.cell(usable_w, 6, _safe_text(f"    {letter}) {opt}"), ln=True)
        pdf.ln(3)
        answer_key.append(f"{i+1}. {q['answer']}")

    # Answer Key
    pdf.add_page()
    pdf.set_font("Arial", "B", 14)
    pdf.set_text_color(r, g, b)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 10, "Answer Key", ln=True)
    pdf.ln(2)
    pdf.set_font("Arial", "", 11)
    pdf.set_text_color(dr, dg, db)
    for ans in answer_key:
        pdf.set_x(pdf.l_margin)
        pdf.cell(usable_w, 6, _safe_text(ans), ln=True)

    return "knowledge_quiz.pdf", bytes(pdf.output())


def _generate_handbook(project) -> tuple:
    """Generate comprehensive process handbook (.pdf)"""
    from fpdf import FPDF

    brand = _get_community_branding(project)
    config = project.configuration
    pdf = FPDF()
    pdf.set_margins(15, 15, 15)
    pdf.set_auto_page_break(auto=True, margin=20)

    r, g, b = brand["primary_color"]
    dr, dg, db = brand["dark_color"]
    usable_w = pdf.w - pdf.l_margin - pdf.r_margin

    def section_header(text, size=16):
        pdf.set_x(pdf.l_margin)
        pdf.set_font("Arial", "B", size)
        pdf.set_text_color(r, g, b)
        pdf.cell(usable_w, 10, _safe_text(text), ln=True)
        # Underline
        pdf.set_draw_color(r, g, b)
        pdf.line(pdf.l_margin, pdf.get_y(), 80, pdf.get_y())
        pdf.ln(4)
        pdf.set_text_color(dr, dg, db)

    def sub_header(text):
        pdf.set_x(pdf.l_margin)
        pdf.set_font("Arial", "B", 12)
        pdf.set_text_color(dr, dg, db)
        pdf.cell(usable_w, 8, _safe_text(text), ln=True)
        pdf.ln(1)

    def body_text(text):
        pdf.set_x(pdf.l_margin)
        pdf.set_font("Arial", "", 10)
        pdf.set_text_color(60, 60, 60)
        pdf.multi_cell(usable_w, 5, _safe_text(text, 1000))
        pdf.ln(2)

    def bullet(text):
        pdf.set_x(pdf.l_margin)
        pdf.set_font("Arial", "", 10)
        pdf.set_text_color(60, 60, 60)
        pdf.multi_cell(usable_w, 5, _safe_text(f"     - {text}", 500))

    # ---- COVER PAGE ----
    pdf.add_page()
    pdf.set_font("Arial", "B", 32)
    pdf.set_text_color(r, g, b)
    pdf.ln(40)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 15, _safe_text(brand["name"]), ln=True, align="C")
    pdf.set_font("Arial", "", 20)
    pdf.set_text_color(dr, dg, db)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 12, "Process Handbook", ln=True, align="C")
    pdf.set_font("Arial", "", 14)
    pdf.set_text_color(100, 100, 100)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 8, "OpenGov Implementation Guide", ln=True, align="C")
    pdf.ln(10)
    pdf.set_font("Arial", "", 11)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 6, _safe_text(f"Generated: {datetime.utcnow().strftime('%B %d, %Y')}"), ln=True, align="C")
    if project.community_url:
        pdf.set_x(pdf.l_margin)
        pdf.cell(usable_w, 6, _safe_text(f"Website: {project.community_url}"), ln=True, align="C")
    pdf.ln(20)
    pdf.set_draw_color(r, g, b)
    pdf.line(60, pdf.get_y(), 150, pdf.get_y())
    pdf.ln(5)
    pdf.set_font("Arial", "I", 10)
    pdf.set_text_color(120, 120, 120)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 6, "CONFIDENTIAL - For staff use only", ln=True, align="C")

    # ---- TABLE OF CONTENTS ----
    pdf.add_page()
    section_header("Table of Contents")
    pdf.set_font("Arial", "", 11)
    pdf.set_text_color(dr, dg, db)
    toc = [
        "1. Introduction & Project Overview",
        "2. System Overview & Key Features",
    ]
    for i, rt in enumerate(config.record_types, 1):
        toc.append(f"  2.{i}. {rt.name}")
    toc.extend([
        "3. Fee Schedules",
        "4. Departments & Organization",
        "5. User Roles & Permissions",
        "6. Workflow & Approval Processes",
        "7. Document Requirements",
        "8. Best Practices & Tips",
        "9. Support & Resources",
    ])
    for item in toc:
        pdf.set_x(pdf.l_margin)
        pdf.cell(usable_w, 6, _safe_text(item), ln=True)

    # ---- SECTION 1: INTRODUCTION ----
    pdf.add_page()
    section_header("1. Introduction & Project Overview")
    body_text(
        f"This handbook serves as a comprehensive reference guide for the {brand['name']} "
        "OpenGov implementation. It covers all configured record types, fee schedules, "
        "workflow processes, departmental responsibilities, and user roles. "
        "Staff should use this document alongside the Training Deck and FAQ Guide "
        "as part of the complete training program."
    )
    body_text(
        f"The {brand['name']} implementation includes {len(config.record_types)} record types, "
        f"{len(config.departments)} departments, and {len(config.user_roles)} user roles. "
        "The system is designed to streamline permitting, licensing, and code enforcement "
        "processes through automated workflow routing, online submissions, and integrated "
        "payment processing."
    )

    # ---- SECTION 2: RECORD TYPES ----
    pdf.add_page()
    section_header("2. Record Types & Applications")
    body_text(
        f"The following {len(config.record_types)} record types have been configured for {brand['name']}. "
        "Each record type defines the application form, required documents, fee schedule, "
        "and approval workflow for a specific permit, license, or enforcement action."
    )

    for idx, rt in enumerate(config.record_types, 1):
        pdf.add_page()
        sub_header(f"2.{idx}. {rt.name}")

        if rt.category:
            body_text(f"Category: {rt.category}")
        if rt.description:
            body_text(rt.description)

        # Form Fields
        if rt.form_fields:
            pdf.ln(2)
            sub_header("Application Form Fields")
            for ff in rt.form_fields[:12]:
                req_label = "Required" if ff.required else "Optional"
                bullet(f"{ff.name} ({ff.field_type}, {req_label})")
            if len(rt.form_fields) > 12:
                body_text(f"... and {len(rt.form_fields) - 12} additional fields")

        # Workflow
        if rt.workflow_steps:
            pdf.ln(2)
            sub_header("Workflow Steps")
            steps = sorted(rt.workflow_steps, key=lambda s: s.order)
            for ws in steps[:10]:
                role_info = f" [Assigned to: {ws.assigned_role}]" if ws.assigned_role else ""
                bullet(f"Step {ws.order}: {ws.name}{role_info}")

        # Fees
        if rt.fees:
            pdf.ln(2)
            sub_header("Fee Schedule")
            for fee in rt.fees[:8]:
                bullet(f"{fee.name}: ${fee.amount:.2f} ({fee.fee_type})")

        # Required Documents
        if rt.required_documents:
            pdf.ln(2)
            sub_header("Required Documents")
            for doc in rt.required_documents[:8]:
                desc = f" - {doc.description[:80]}" if doc.description else ""
                bullet(f"{doc.name}{desc}")

    # ---- SECTION 3: FEE SCHEDULES ----
    pdf.add_page()
    section_header("3. Complete Fee Schedule")
    body_text("The following is a consolidated view of all fees across record types.")
    for rt in config.record_types:
        if rt.fees:
            sub_header(rt.name)
            for fee in rt.fees[:10]:
                bullet(f"{fee.name}: ${fee.amount:.2f} ({fee.fee_type})")
            pdf.ln(2)

    # ---- SECTION 4: DEPARTMENTS ----
    pdf.add_page()
    section_header("4. Departments & Organization")
    if config.departments:
        body_text(f"{brand['name']} has {len(config.departments)} departments configured in the system:")
        for dept in config.departments:
            sub_header(dept.name)
            if dept.description:
                body_text(dept.description)
            else:
                body_text(f"Responsible for reviewing and processing applications assigned to {dept.name}.")
            pdf.ln(1)
    else:
        body_text("Departments have not been configured yet.")

    # ---- SECTION 5: USER ROLES ----
    pdf.add_page()
    section_header("5. User Roles & Permissions")
    if config.user_roles:
        body_text(f"The system defines {len(config.user_roles)} user roles with specific permissions:")
        for role in config.user_roles:
            sub_header(role.name)
            if role.description:
                body_text(role.description)
            if role.permissions:
                body_text(f"Permissions: {', '.join(role.permissions[:8])}")
            pdf.ln(1)
    else:
        body_text("User roles have not been configured yet.")

    # ---- SECTION 6: WORKFLOW PROCESSES ----
    pdf.add_page()
    section_header("6. Workflow & Approval Processes")
    body_text(
        "Each record type has a defined workflow that governs how applications move through "
        "the review and approval process. Below is a summary of each workflow."
    )
    for rt in config.record_types[:8]:
        if rt.workflow_steps:
            sub_header(f"{rt.name} Workflow")
            steps = sorted(rt.workflow_steps, key=lambda s: s.order)
            for ws in steps[:8]:
                role_info = f" ({ws.assigned_role})" if ws.assigned_role else ""
                bullet(f"{ws.order}. {ws.name}{role_info}")
            pdf.ln(2)

    # ---- SECTION 7: DOCUMENT REQUIREMENTS ----
    pdf.add_page()
    section_header("7. Document Requirements")
    body_text("Applicants must submit the following documents as part of their applications:")
    for rt in config.record_types[:8]:
        if rt.required_documents:
            sub_header(rt.name)
            for doc in rt.required_documents[:6]:
                desc = f": {doc.description[:100]}" if doc.description else ""
                bullet(f"{doc.name}{desc}")
            pdf.ln(2)

    # ---- SECTION 8: BEST PRACTICES ----
    pdf.add_page()
    section_header("8. Best Practices & Tips")
    tips = [
        "Always verify applicant information is complete before advancing the workflow",
        "Use the internal notes feature to communicate with other reviewers",
        "Set up email notification preferences to stay informed of new assignments",
        "Check the dashboard regularly for pending items and approaching deadlines",
        "Use the search function to quickly find applications by number, name, or address",
        "Document all communications with applicants in the application's activity log",
        "Follow the established fee schedule - do not manually adjust fees without authorization",
        "Report any system issues immediately to the system administrator",
        "Refer applicants to the online portal for 24/7 self-service options",
        "Complete the Knowledge Quiz to verify your understanding of the new system",
    ]
    for tip in tips:
        bullet(tip)
        pdf.ln(1)

    # ---- SECTION 9: SUPPORT ----
    pdf.add_page()
    section_header("9. Support & Resources")
    body_text(
        f"For technical support, contact the {brand['name']} Help Desk. "
        "Additional resources include:"
    )
    bullet("Training Deck (PowerPoint) - Comprehensive staff training presentation")
    bullet("FAQ Guide (PDF) - Answers to frequently asked questions")
    bullet("Knowledge Quiz (PDF) - Self-assessment to verify understanding")
    bullet(f"Community Website: {project.community_url or 'Contact your administrator'}")
    pdf.ln(5)
    body_text(
        "For system-specific questions, contact your OpenGov implementation specialist. "
        "For process-related questions, consult your department supervisor."
    )
    pdf.ln(10)
    pdf.set_font("Arial", "I", 9)
    pdf.set_text_color(150, 150, 150)
    pdf.set_x(pdf.l_margin)
    pdf.cell(usable_w, 5, _safe_text(f"This document was auto-generated for {brand['name']} by OpenGov Auto Implementation."), ln=True, align="C")

    return "process_handbook.pdf", bytes(pdf.output())


# ============================================================================
# DATA SOURCES - Multi-source context ingestion
# ============================================================================

PEER_CITY_TEMPLATES = [
    {
        "id": "small-town-basic",
        "name": "Small Town - Basic PLC",
        "description": "Starter config for towns under 15,000 pop. Covers building permits, business licenses, and code complaints.",
        "population": "< 15,000",
        "tags": ["small", "basic", "residential"],
        "record_types": [
            {"name": "Residential Building Permit", "category": "Building", "form_fields": [
                {"name": "Property Address", "field_type": "text", "required": True},
                {"name": "Parcel Number", "field_type": "text", "required": True},
                {"name": "Owner Name", "field_type": "text", "required": True},
                {"name": "Contractor Name", "field_type": "text", "required": True},
                {"name": "Contractor License #", "field_type": "text", "required": True},
                {"name": "Project Description", "field_type": "textarea", "required": True},
                {"name": "Estimated Value", "field_type": "number", "required": True},
                {"name": "Square Footage", "field_type": "number", "required": False},
            ], "fees": [
                {"name": "Building Permit Fee", "amount": 150.00, "fee_type": "flat"},
                {"name": "Plan Review Fee", "amount": 75.00, "fee_type": "flat"},
            ], "workflow_steps": [
                {"name": "Application Submitted", "order": 1, "assigned_role": "Clerk"},
                {"name": "Plan Review", "order": 2, "assigned_role": "Plan Reviewer"},
                {"name": "Permit Issued", "order": 3, "assigned_role": "Building Official"},
                {"name": "Inspection Scheduled", "order": 4, "assigned_role": "Inspector"},
                {"name": "Final Inspection", "order": 5, "assigned_role": "Inspector"},
            ], "required_documents": [
                {"name": "Site Plan", "required": True, "stage": "submission"},
                {"name": "Construction Plans", "required": True, "stage": "submission"},
                {"name": "Contractor License", "required": True, "stage": "submission"},
            ]},
            {"name": "Business License", "category": "Licensing", "form_fields": [
                {"name": "Business Name", "field_type": "text", "required": True},
                {"name": "Business Address", "field_type": "text", "required": True},
                {"name": "Owner Name", "field_type": "text", "required": True},
                {"name": "Business Type", "field_type": "select", "required": True, "options": ["Retail", "Food Service", "Professional Services", "Home Occupation", "Other"]},
                {"name": "EIN / Tax ID", "field_type": "text", "required": True},
                {"name": "Number of Employees", "field_type": "number", "required": True},
            ], "fees": [
                {"name": "Annual License Fee", "amount": 100.00, "fee_type": "flat"},
                {"name": "Late Renewal Penalty", "amount": 25.00, "fee_type": "flat"},
            ], "workflow_steps": [
                {"name": "Application Received", "order": 1, "assigned_role": "Clerk"},
                {"name": "Zoning Review", "order": 2, "assigned_role": "Zoning Officer"},
                {"name": "License Issued", "order": 3, "assigned_role": "Clerk"},
            ], "required_documents": [
                {"name": "State Registration", "required": True, "stage": "submission"},
                {"name": "Certificate of Insurance", "required": True, "stage": "submission"},
            ]},
            {"name": "Code Enforcement Complaint", "category": "Code Enforcement", "form_fields": [
                {"name": "Complaint Location", "field_type": "text", "required": True},
                {"name": "Violation Type", "field_type": "select", "required": True, "options": ["Property Maintenance", "Zoning Violation", "Noise", "Overgrown Vegetation", "Abandoned Vehicle", "Other"]},
                {"name": "Description", "field_type": "textarea", "required": True},
                {"name": "Complainant Name", "field_type": "text", "required": False},
                {"name": "Complainant Phone", "field_type": "text", "required": False},
            ], "fees": [], "workflow_steps": [
                {"name": "Complaint Filed", "order": 1, "assigned_role": "Clerk"},
                {"name": "Investigation", "order": 2, "assigned_role": "Code Officer"},
                {"name": "Notice Issued", "order": 3, "assigned_role": "Code Officer"},
                {"name": "Re-Inspection", "order": 4, "assigned_role": "Code Officer"},
                {"name": "Resolution", "order": 5, "assigned_role": "Code Officer"},
            ], "required_documents": [
                {"name": "Photos of Violation", "required": False, "stage": "submission"},
            ]},
        ],
        "departments": [
            {"name": "Building Department", "description": "Handles building permits, plan review, and inspections"},
            {"name": "Clerk's Office", "description": "Manages business licenses, applications intake, and public records"},
            {"name": "Code Enforcement", "description": "Investigates and resolves code violations and complaints"},
        ],
        "user_roles": [
            {"name": "Administrator", "description": "Full system access", "permissions": ["manage_users", "manage_config", "view_reports", "approve_all"]},
            {"name": "Clerk", "description": "Application intake and processing", "permissions": ["create_applications", "update_status", "collect_payments"]},
            {"name": "Plan Reviewer", "description": "Reviews construction plans and documents", "permissions": ["review_plans", "add_comments", "approve_plans"]},
            {"name": "Inspector", "description": "Conducts field inspections", "permissions": ["schedule_inspections", "record_results", "issue_violations"]},
            {"name": "Code Officer", "description": "Code enforcement investigations", "permissions": ["investigate_complaints", "issue_notices", "close_cases"]},
        ],
    },
    {
        "id": "mid-city-full",
        "name": "Mid-Size City - Full PLC Suite",
        "description": "Comprehensive config for cities 15,000-100,000 pop. Full building, planning, licensing, and code enforcement.",
        "population": "15,000 - 100,000",
        "tags": ["medium", "comprehensive", "commercial", "residential"],
        "record_types": [
            {"name": "Residential Building Permit", "category": "Building", "form_fields": [
                {"name": "Property Address", "field_type": "text", "required": True},
                {"name": "Parcel/APN", "field_type": "text", "required": True},
                {"name": "Owner Name", "field_type": "text", "required": True},
                {"name": "Owner Phone", "field_type": "text", "required": True},
                {"name": "Owner Email", "field_type": "email", "required": True},
                {"name": "Contractor Name", "field_type": "text", "required": True},
                {"name": "Contractor License #", "field_type": "text", "required": True},
                {"name": "Project Type", "field_type": "select", "required": True, "options": ["New Construction", "Addition", "Remodel", "Repair", "Demolition", "Accessory Structure"]},
                {"name": "Project Description", "field_type": "textarea", "required": True},
                {"name": "Estimated Valuation", "field_type": "number", "required": True},
                {"name": "Square Footage", "field_type": "number", "required": True},
                {"name": "Number of Stories", "field_type": "number", "required": False},
                {"name": "Flood Zone", "field_type": "select", "required": False, "options": ["Zone A", "Zone AE", "Zone X", "None"]},
            ], "fees": [
                {"name": "Building Permit Fee", "amount": 250.00, "fee_type": "flat"},
                {"name": "Plan Review Fee", "amount": 125.00, "fee_type": "flat"},
                {"name": "Technology Fee", "amount": 15.00, "fee_type": "flat"},
                {"name": "Valuation Surcharge", "amount": 0.00, "fee_type": "calculated", "formula": "valuation * 0.015"},
            ], "workflow_steps": [
                {"name": "Application Submitted", "order": 1, "assigned_role": "Permit Technician"},
                {"name": "Completeness Check", "order": 2, "assigned_role": "Permit Technician"},
                {"name": "Plan Review - Building", "order": 3, "assigned_role": "Plan Reviewer"},
                {"name": "Plan Review - Fire", "order": 4, "assigned_role": "Fire Marshal"},
                {"name": "Plan Review - Engineering", "order": 5, "assigned_role": "City Engineer"},
                {"name": "Corrections Required", "order": 6, "assigned_role": "Plan Reviewer"},
                {"name": "Approved for Issuance", "order": 7, "assigned_role": "Building Official"},
                {"name": "Permit Issued", "order": 8, "assigned_role": "Permit Technician"},
                {"name": "Inspections", "order": 9, "assigned_role": "Building Inspector"},
                {"name": "Final / Certificate of Occupancy", "order": 10, "assigned_role": "Building Official"},
            ], "required_documents": [
                {"name": "Site Plan", "required": True, "stage": "submission"},
                {"name": "Architectural Plans", "required": True, "stage": "submission"},
                {"name": "Structural Calculations", "required": True, "stage": "submission"},
                {"name": "Energy Compliance (Title 24)", "required": True, "stage": "submission"},
                {"name": "Soils Report", "required": False, "stage": "submission"},
                {"name": "Contractor License Copy", "required": True, "stage": "submission"},
            ]},
            {"name": "Commercial Building Permit", "category": "Building", "form_fields": [
                {"name": "Property Address", "field_type": "text", "required": True},
                {"name": "Parcel/APN", "field_type": "text", "required": True},
                {"name": "Owner Name", "field_type": "text", "required": True},
                {"name": "Applicant/Agent", "field_type": "text", "required": True},
                {"name": "Architect of Record", "field_type": "text", "required": True},
                {"name": "General Contractor", "field_type": "text", "required": True},
                {"name": "Project Type", "field_type": "select", "required": True, "options": ["New Commercial", "Tenant Improvement", "Addition", "Renovation", "Change of Use"]},
                {"name": "Occupancy Type", "field_type": "select", "required": True, "options": ["Assembly", "Business", "Educational", "Factory", "Hazardous", "Institutional", "Mercantile", "Residential", "Storage", "Utility"]},
                {"name": "Project Description", "field_type": "textarea", "required": True},
                {"name": "Estimated Valuation", "field_type": "number", "required": True},
                {"name": "Building Area (sq ft)", "field_type": "number", "required": True},
                {"name": "ADA Compliance", "field_type": "checkbox", "required": True},
            ], "fees": [
                {"name": "Commercial Permit Fee", "amount": 500.00, "fee_type": "flat"},
                {"name": "Plan Review Fee", "amount": 250.00, "fee_type": "flat"},
                {"name": "Fire Review Fee", "amount": 150.00, "fee_type": "flat"},
                {"name": "Technology Fee", "amount": 25.00, "fee_type": "flat"},
            ], "workflow_steps": [
                {"name": "Application Submitted", "order": 1, "assigned_role": "Permit Technician"},
                {"name": "Completeness Review", "order": 2, "assigned_role": "Permit Technician"},
                {"name": "Plan Review - Building", "order": 3, "assigned_role": "Plan Reviewer"},
                {"name": "Plan Review - Fire/Life Safety", "order": 4, "assigned_role": "Fire Marshal"},
                {"name": "Plan Review - Public Works", "order": 5, "assigned_role": "City Engineer"},
                {"name": "Plan Review - Planning/Zoning", "order": 6, "assigned_role": "Planner"},
                {"name": "Corrections Cycle", "order": 7, "assigned_role": "Plan Reviewer"},
                {"name": "Final Approval", "order": 8, "assigned_role": "Building Official"},
                {"name": "Permit Issued", "order": 9, "assigned_role": "Permit Technician"},
                {"name": "Inspections", "order": 10, "assigned_role": "Building Inspector"},
                {"name": "Certificate of Occupancy", "order": 11, "assigned_role": "Building Official"},
            ], "required_documents": [
                {"name": "Architectural Plans (stamped)", "required": True, "stage": "submission"},
                {"name": "Structural Plans & Calcs", "required": True, "stage": "submission"},
                {"name": "MEP Plans", "required": True, "stage": "submission"},
                {"name": "Fire Sprinkler Plans", "required": True, "stage": "submission"},
                {"name": "Energy Compliance", "required": True, "stage": "submission"},
                {"name": "ADA Compliance Report", "required": True, "stage": "submission"},
                {"name": "Geotechnical Report", "required": False, "stage": "submission"},
            ]},
            {"name": "Encroachment Permit", "category": "Public Works", "form_fields": [
                {"name": "Location/Address", "field_type": "text", "required": True},
                {"name": "Work Description", "field_type": "textarea", "required": True},
                {"name": "Encroachment Type", "field_type": "select", "required": True, "options": ["Sidewalk Cut", "Street Cut", "Driveway Approach", "Utility Trench", "Other"]},
                {"name": "Contractor Name", "field_type": "text", "required": True},
                {"name": "Start Date", "field_type": "date", "required": True},
                {"name": "End Date", "field_type": "date", "required": True},
                {"name": "Traffic Control Required", "field_type": "checkbox", "required": True},
            ], "fees": [
                {"name": "Encroachment Permit Fee", "amount": 200.00, "fee_type": "flat"},
                {"name": "Inspection Fee", "amount": 100.00, "fee_type": "flat"},
            ], "workflow_steps": [
                {"name": "Application Submitted", "order": 1, "assigned_role": "Permit Technician"},
                {"name": "Engineering Review", "order": 2, "assigned_role": "City Engineer"},
                {"name": "Permit Issued", "order": 3, "assigned_role": "Permit Technician"},
                {"name": "Pre-Construction Inspection", "order": 4, "assigned_role": "Inspector"},
                {"name": "Final Inspection", "order": 5, "assigned_role": "Inspector"},
            ], "required_documents": [
                {"name": "Traffic Control Plan", "required": True, "stage": "submission"},
                {"name": "Site Plan/Drawings", "required": True, "stage": "submission"},
                {"name": "Insurance Certificate", "required": True, "stage": "submission"},
            ]},
            {"name": "Business License", "category": "Licensing", "form_fields": [
                {"name": "Business Name (DBA)", "field_type": "text", "required": True},
                {"name": "Legal Entity Name", "field_type": "text", "required": True},
                {"name": "Business Address", "field_type": "text", "required": True},
                {"name": "Mailing Address", "field_type": "text", "required": True},
                {"name": "Owner/Manager Name", "field_type": "text", "required": True},
                {"name": "Phone", "field_type": "text", "required": True},
                {"name": "Email", "field_type": "email", "required": True},
                {"name": "NAICS Code", "field_type": "text", "required": False},
                {"name": "Business Type", "field_type": "select", "required": True, "options": ["Retail", "Restaurant/Food", "Professional Services", "Contractor", "Home Occupation", "Mobile/Vendor", "Manufacturing", "Other"]},
                {"name": "Number of Employees", "field_type": "number", "required": True},
                {"name": "Estimated Annual Revenue", "field_type": "number", "required": False},
            ], "fees": [
                {"name": "Business License Fee", "amount": 150.00, "fee_type": "flat"},
                {"name": "Home Occupation Surcharge", "amount": 50.00, "fee_type": "conditional"},
                {"name": "Late Penalty (per month)", "amount": 25.00, "fee_type": "flat"},
            ], "workflow_steps": [
                {"name": "Application Received", "order": 1, "assigned_role": "License Clerk"},
                {"name": "Zoning Verification", "order": 2, "assigned_role": "Planner"},
                {"name": "Fire Inspection Required", "order": 3, "assigned_role": "Fire Marshal"},
                {"name": "Approved", "order": 4, "assigned_role": "License Clerk"},
                {"name": "License Issued", "order": 5, "assigned_role": "License Clerk"},
            ], "required_documents": [
                {"name": "State Business Registration", "required": True, "stage": "submission"},
                {"name": "Federal EIN", "required": True, "stage": "submission"},
                {"name": "Certificate of Insurance", "required": True, "stage": "submission"},
                {"name": "Lease Agreement", "required": False, "stage": "submission"},
            ]},
            {"name": "Special Event Permit", "category": "Licensing", "form_fields": [
                {"name": "Event Name", "field_type": "text", "required": True},
                {"name": "Event Location", "field_type": "text", "required": True},
                {"name": "Event Date(s)", "field_type": "date", "required": True},
                {"name": "Setup Time", "field_type": "text", "required": True},
                {"name": "Event Hours", "field_type": "text", "required": True},
                {"name": "Expected Attendance", "field_type": "number", "required": True},
                {"name": "Alcohol Served", "field_type": "checkbox", "required": True},
                {"name": "Road Closures Needed", "field_type": "checkbox", "required": True},
                {"name": "Amplified Sound", "field_type": "checkbox", "required": True},
                {"name": "Food Vendors", "field_type": "number", "required": False},
            ], "fees": [
                {"name": "Event Permit Fee", "amount": 200.00, "fee_type": "flat"},
                {"name": "Road Closure Fee", "amount": 300.00, "fee_type": "conditional"},
            ], "workflow_steps": [
                {"name": "Application Submitted", "order": 1, "assigned_role": "Clerk"},
                {"name": "Police Review", "order": 2, "assigned_role": "Police Dept"},
                {"name": "Fire Review", "order": 3, "assigned_role": "Fire Marshal"},
                {"name": "Public Works Review", "order": 4, "assigned_role": "City Engineer"},
                {"name": "Final Approval", "order": 5, "assigned_role": "City Manager"},
                {"name": "Permit Issued", "order": 6, "assigned_role": "Clerk"},
            ], "required_documents": [
                {"name": "Event Site Map", "required": True, "stage": "submission"},
                {"name": "Certificate of Insurance", "required": True, "stage": "submission"},
                {"name": "ABC Permit (if alcohol)", "required": False, "stage": "submission"},
            ]},
            {"name": "Code Enforcement Case", "category": "Code Enforcement", "form_fields": [
                {"name": "Violation Location", "field_type": "text", "required": True},
                {"name": "Violation Category", "field_type": "select", "required": True, "options": ["Property Maintenance", "Zoning Violation", "Building w/o Permit", "Sign Violation", "Noise", "Overgrown Vegetation", "Abandoned Vehicle", "Illegal Dumping", "Other"]},
                {"name": "Description", "field_type": "textarea", "required": True},
                {"name": "Complainant Name", "field_type": "text", "required": False},
                {"name": "Complainant Phone", "field_type": "text", "required": False},
                {"name": "Anonymous Complaint", "field_type": "checkbox", "required": False},
            ], "fees": [], "workflow_steps": [
                {"name": "Complaint Received", "order": 1, "assigned_role": "Code Enforcement Clerk"},
                {"name": "Case Assigned", "order": 2, "assigned_role": "Code Enforcement Supervisor"},
                {"name": "Initial Investigation", "order": 3, "assigned_role": "Code Officer"},
                {"name": "Notice of Violation", "order": 4, "assigned_role": "Code Officer"},
                {"name": "Compliance Period", "order": 5, "assigned_role": "Code Officer"},
                {"name": "Re-Inspection", "order": 6, "assigned_role": "Code Officer"},
                {"name": "Citation / Abatement", "order": 7, "assigned_role": "Code Officer"},
                {"name": "Case Closed", "order": 8, "assigned_role": "Code Enforcement Supervisor"},
            ], "required_documents": [
                {"name": "Violation Photos", "required": False, "stage": "investigation"},
            ]},
        ],
        "departments": [
            {"name": "Building & Safety", "description": "Building permits, plan review, inspections, and code compliance"},
            {"name": "Planning & Zoning", "description": "Land use planning, zoning verification, and entitlements"},
            {"name": "Public Works/Engineering", "description": "Encroachment permits, infrastructure, and capital projects"},
            {"name": "Fire Prevention", "description": "Fire safety review, inspections, and special hazard permits"},
            {"name": "Business License Division", "description": "Business licensing, renewals, and revenue collection"},
            {"name": "Code Enforcement", "description": "Municipal code compliance, investigations, and abatement"},
        ],
        "user_roles": [
            {"name": "System Administrator", "description": "Full system configuration and management", "permissions": ["manage_all"]},
            {"name": "Department Director", "description": "Department oversight and final approvals", "permissions": ["approve_all", "view_reports", "manage_staff"]},
            {"name": "Permit Technician", "description": "Intake, routing, and issuance of permits", "permissions": ["create_applications", "update_status", "collect_payments", "issue_permits"]},
            {"name": "Plan Reviewer", "description": "Technical review of construction plans", "permissions": ["review_plans", "add_comments", "request_corrections", "approve_plans"]},
            {"name": "Building Inspector", "description": "Field inspections for building permits", "permissions": ["schedule_inspections", "record_results", "approve_inspections"]},
            {"name": "Fire Marshal / Fire Reviewer", "description": "Fire and life safety reviews and inspections", "permissions": ["review_fire", "approve_fire", "schedule_inspections"]},
            {"name": "City Engineer", "description": "Engineering review for public works permits", "permissions": ["review_engineering", "approve_engineering"]},
            {"name": "Planner / Zoning Officer", "description": "Zoning verification and planning review", "permissions": ["review_zoning", "verify_land_use"]},
            {"name": "Code Enforcement Officer", "description": "Investigates complaints and enforces codes", "permissions": ["investigate_complaints", "issue_notices", "issue_citations"]},
            {"name": "License Clerk", "description": "Business license processing and renewals", "permissions": ["process_licenses", "collect_payments", "issue_licenses"]},
            {"name": "Public Portal User", "description": "External applicant self-service access", "permissions": ["submit_applications", "upload_documents", "make_payments", "view_status"]},
        ],
    },
    {
        "id": "county-planning",
        "name": "County - Planning & Land Use",
        "description": "County-level configuration focused on planning, land use, subdivisions, and environmental review.",
        "population": "County-level",
        "tags": ["county", "planning", "land_use", "environmental"],
        "record_types": [
            {"name": "Conditional Use Permit (CUP)", "category": "Planning", "form_fields": [
                {"name": "Property Address/APN", "field_type": "text", "required": True},
                {"name": "Property Owner", "field_type": "text", "required": True},
                {"name": "Applicant Name", "field_type": "text", "required": True},
                {"name": "Proposed Use", "field_type": "textarea", "required": True},
                {"name": "Current Zoning", "field_type": "text", "required": True},
                {"name": "Lot Size (acres)", "field_type": "number", "required": True},
                {"name": "Adjacent Land Uses", "field_type": "textarea", "required": True},
            ], "fees": [
                {"name": "CUP Application Fee", "amount": 2500.00, "fee_type": "flat"},
                {"name": "Environmental Review Fee", "amount": 1500.00, "fee_type": "flat"},
                {"name": "Public Notice Fee", "amount": 300.00, "fee_type": "flat"},
            ], "workflow_steps": [
                {"name": "Application Submitted", "order": 1, "assigned_role": "Planning Technician"},
                {"name": "Completeness Review", "order": 2, "assigned_role": "Planner"},
                {"name": "Environmental Review (CEQA)", "order": 3, "assigned_role": "Environmental Planner"},
                {"name": "Staff Report Prepared", "order": 4, "assigned_role": "Senior Planner"},
                {"name": "Public Hearing Notice", "order": 5, "assigned_role": "Planning Technician"},
                {"name": "Planning Commission Hearing", "order": 6, "assigned_role": "Planning Commission"},
                {"name": "Decision / Conditions", "order": 7, "assigned_role": "Planning Director"},
                {"name": "Appeal Period", "order": 8, "assigned_role": "Planning Technician"},
            ], "required_documents": [
                {"name": "Site Plan (to scale)", "required": True, "stage": "submission"},
                {"name": "Project Description Narrative", "required": True, "stage": "submission"},
                {"name": "Environmental Assessment", "required": True, "stage": "review"},
                {"name": "Traffic Impact Study", "required": False, "stage": "review"},
            ]},
            {"name": "Subdivision / Tentative Map", "category": "Planning", "form_fields": [
                {"name": "Property Address/APN", "field_type": "text", "required": True},
                {"name": "Owner/Developer", "field_type": "text", "required": True},
                {"name": "Engineer of Record", "field_type": "text", "required": True},
                {"name": "Number of Lots", "field_type": "number", "required": True},
                {"name": "Total Acreage", "field_type": "number", "required": True},
                {"name": "Subdivision Type", "field_type": "select", "required": True, "options": ["Residential", "Commercial", "Mixed Use", "Industrial"]},
                {"name": "Proposed Infrastructure", "field_type": "textarea", "required": True},
            ], "fees": [
                {"name": "Tentative Map Fee", "amount": 5000.00, "fee_type": "flat"},
                {"name": "Per Lot Fee", "amount": 200.00, "fee_type": "per_unit"},
                {"name": "Environmental Review", "amount": 3000.00, "fee_type": "flat"},
            ], "workflow_steps": [
                {"name": "Application Filed", "order": 1, "assigned_role": "Planning Technician"},
                {"name": "Completeness Review", "order": 2, "assigned_role": "Senior Planner"},
                {"name": "Agency Circulation", "order": 3, "assigned_role": "Planner"},
                {"name": "Environmental Review", "order": 4, "assigned_role": "Environmental Planner"},
                {"name": "Conditions Development", "order": 5, "assigned_role": "Senior Planner"},
                {"name": "Planning Commission", "order": 6, "assigned_role": "Planning Commission"},
                {"name": "Board of Supervisors", "order": 7, "assigned_role": "Board"},
                {"name": "Final Map Recordation", "order": 8, "assigned_role": "Planning Technician"},
            ], "required_documents": [
                {"name": "Tentative Map (engineered)", "required": True, "stage": "submission"},
                {"name": "Preliminary Grading Plan", "required": True, "stage": "submission"},
                {"name": "Utility Plan", "required": True, "stage": "submission"},
                {"name": "Environmental Impact Report", "required": True, "stage": "review"},
                {"name": "Traffic Study", "required": True, "stage": "review"},
                {"name": "Water/Sewer Will-Serve Letters", "required": True, "stage": "review"},
            ]},
            {"name": "Grading Permit", "category": "Building", "form_fields": [
                {"name": "Site Address", "field_type": "text", "required": True},
                {"name": "Cut Volume (cubic yards)", "field_type": "number", "required": True},
                {"name": "Fill Volume (cubic yards)", "field_type": "number", "required": True},
                {"name": "Disturbed Area (sq ft)", "field_type": "number", "required": True},
                {"name": "SWPPP Required", "field_type": "checkbox", "required": True},
            ], "fees": [
                {"name": "Grading Permit Fee", "amount": 750.00, "fee_type": "flat"},
                {"name": "Erosion Control Fee", "amount": 250.00, "fee_type": "flat"},
            ], "workflow_steps": [
                {"name": "Application", "order": 1, "assigned_role": "Permit Technician"},
                {"name": "Engineering Review", "order": 2, "assigned_role": "County Engineer"},
                {"name": "Environmental Check", "order": 3, "assigned_role": "Environmental Planner"},
                {"name": "Permit Issued", "order": 4, "assigned_role": "Permit Technician"},
                {"name": "Inspections", "order": 5, "assigned_role": "Inspector"},
            ], "required_documents": [
                {"name": "Grading Plan (engineered)", "required": True, "stage": "submission"},
                {"name": "Soils/Geotechnical Report", "required": True, "stage": "submission"},
                {"name": "SWPPP", "required": True, "stage": "submission"},
            ]},
        ],
        "departments": [
            {"name": "Planning & Development", "description": "Land use planning, zoning, and entitlements"},
            {"name": "Public Works / Engineering", "description": "Grading, infrastructure, and public improvement review"},
            {"name": "Environmental Resources", "description": "CEQA review, environmental compliance, and sustainability"},
            {"name": "Building & Safety", "description": "Structural review, building permits, and inspections"},
        ],
        "user_roles": [
            {"name": "Planning Director", "description": "Oversees planning department, final authority", "permissions": ["manage_all", "approve_all"]},
            {"name": "Senior Planner", "description": "Complex project review and staff reports", "permissions": ["review_plans", "write_reports", "present_hearings"]},
            {"name": "Associate Planner", "description": "Project review and processing", "permissions": ["review_plans", "add_conditions"]},
            {"name": "Planning Technician", "description": "Intake, routing, and records", "permissions": ["create_applications", "route_applications"]},
            {"name": "Environmental Planner", "description": "CEQA and environmental review", "permissions": ["review_environmental", "approve_environmental"]},
            {"name": "County Engineer", "description": "Engineering and infrastructure review", "permissions": ["review_engineering", "approve_engineering"]},
        ],
    },
]


def _extract_with_ai(prompt_text, project_context="", operation_type="extraction"):
    """Use Claude to extract structured data from text"""
    if not claude_service.is_available():
        print(f"[AI] Skipping AI extraction ({operation_type}): service not available")
        return None
    try:
        print(f"[AI] Running AI extraction: {operation_type} ({len(prompt_text)} chars)")
        response = claude_service.client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4000,
            messages=[{"role": "user", "content": prompt_text}]
        )

        try:
            tokens_used = (response.usage.input_tokens or 0) + (response.usage.output_tokens or 0)
            ai_usage_tracker.record_call(operation_type, tokens_used, True)
            print(f"[AI] Extraction complete: {operation_type} - {tokens_used} tokens")
        except Exception:
            ai_usage_tracker.record_call(operation_type, 0, True)

        return response.content[0].text
    except Exception as e:
        print(f"[AI] ERROR in extraction ({operation_type}): {e}")
        ai_usage_tracker.record_call(operation_type, 0, False)
        return None


def _scrape_url_text(url):
    """Scrape text content from a URL"""
    import urllib.request
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=15) as resp:
            html = resp.read().decode("utf-8", errors="replace")
        import re
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        return text[:15000]
    except Exception as e:
        return f"Error fetching URL: {str(e)}"


# --- 1. MUNICIPAL CODE / ORDINANCE PARSER ---
@app.post("/api/projects/{project_id}/sources/municipal-code")
async def parse_municipal_code(project_id: str, data: dict):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data

    url = data.get("url", "")
    text = data.get("text", "")
    name = data.get("name", "")

    if not url and not text:
        raise HTTPException(status_code=400, detail="Provide a URL or paste text content")

    source_name = name or (f"Municipal Code: {url[:50]}" if url else "Pasted Municipal Code")
    source = {
        "id": str(uuid.uuid4())[:8],
        "source_type": "municipal_code",
        "name": source_name,
        "status": "processing",
        "url": url,
        "created_at": datetime.utcnow().isoformat(),
    }

    try:
        raw_text = text
        if url and not text:
            raw_text = _scrape_url_text(url)
            if raw_text.startswith("Error") or len(raw_text.strip()) < 100:
                # Many municipal code sites (Municode, etc) use JS rendering
                # or require downloads. Try alternate approaches.
                raw_text = ""
                # Try common Municode API patterns
                if "municode.com" in url.lower():
                    # Try the print/export version of Municode URLs
                    import urllib.request
                    alt_urls = []
                    if "/codes/" in url:
                        alt_urls.append(url.replace("/codes/", "/print/"))
                    for alt in alt_urls:
                        try:
                            req = urllib.request.Request(alt, headers={"User-Agent": "Mozilla/5.0"})
                            with urllib.request.urlopen(req, timeout=15) as resp:
                                html = resp.read().decode("utf-8", errors="replace")
                            import re
                            text_cleaned = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
                            text_cleaned = re.sub(r'<style[^>]*>.*?</style>', '', text_cleaned, flags=re.DOTALL)
                            text_cleaned = re.sub(r'<[^>]+>', ' ', text_cleaned)
                            text_cleaned = re.sub(r'\s+', ' ', text_cleaned).strip()
                            if len(text_cleaned) > 200:
                                raw_text = text_cleaned[:15000]
                                break
                        except Exception:
                            continue

                if not raw_text:
                    source["status"] = "error"
                    source["error_message"] = (
                        "Could not scrape this URL directly. Many municipal code sites use JavaScript rendering "
                        "or require file downloads. Please try: (1) Copy and paste the code text directly, "
                        "(2) Download the PDF from the site and paste its text content, or "
                        "(3) Try a direct link to a specific chapter/section."
                    )
                    sources = list(project.data_sources) if project.data_sources else []
                    sources.append(source)
                    store.update_project(project_id, data_sources=sources)
                    return source

        if not raw_text or len(raw_text.strip()) < 50:
            source["status"] = "error"
            source["error_message"] = "Not enough text content to analyze. Please paste the municipal code text directly."
            sources = list(project.data_sources) if project.data_sources else []
            sources.append(source)
            store.update_project(project_id, data_sources=sources)
            return source

        source["raw_text"] = raw_text[:5000]

        ai_prompt = f"""Analyze this municipal code / ordinance text and extract ALL permit types, license types, and code enforcement processes mentioned.

For each item found, provide:
- type: "permit", "license", or "enforcement"
- name: the official name (e.g., "Building Permit", "Business License", "Code Enforcement Complaint")
- description: brief description of what it covers and when it's needed
- requirements: list of any mentioned requirements (documents, fees, conditions)
- department: which department typically handles this
- triggers: what triggers the need for this permit/license

Be thorough - extract every permit, license, and enforcement action you can find.
Respond in JSON format as a list of objects.

Municipal Code Text:
{raw_text[:10000]}

Community context: {project.customer_name} - {project.community_url}"""

        ai_result = _extract_with_ai(ai_prompt, "", "municipal_code_analysis")
        if ai_result:
            try:
                import re
                json_match = re.search(r'\[.*\]', ai_result, re.DOTALL)
                if json_match:
                    extracted = json.loads(json_match.group())
                else:
                    extracted = [{"raw_analysis": ai_result}]
            except json.JSONDecodeError:
                extracted = [{"raw_analysis": ai_result}]
            source["extracted_data"] = {"requirements": extracted, "url": url, "text_length": len(raw_text)}
        else:
            # Fallback: keyword extraction
            keywords = {
                "permits": ["building permit", "grading permit", "demolition permit", "electrical permit",
                            "plumbing permit", "mechanical permit", "sign permit", "encroachment permit",
                            "excavation permit", "fire permit", "special event permit", "conditional use",
                            "variance", "site plan", "subdivision", "zoning permit"],
                "licenses": ["business license", "contractor license", "liquor license", "vendor permit",
                             "home occupation", "peddler license", "taxi license", "alarm permit",
                             "dog license", "solicitor permit", "rental license"],
                "enforcement": ["code enforcement", "violation", "nuisance", "abatement", "citation",
                                "property maintenance", "zoning violation", "abandoned vehicle",
                                "overgrown vegetation", "illegal dumping"]
            }
            found = []
            text_lower = raw_text.lower()
            for category, terms in keywords.items():
                for term in terms:
                    if term in text_lower:
                        found.append({
                            "type": category.rstrip("s"),
                            "name": term.title(),
                            "description": f"Found reference to '{term}' in municipal code",
                            "source": "keyword_match"
                        })
            source["extracted_data"] = {"requirements": found, "url": url, "text_length": len(raw_text)}

        source["status"] = "completed"
    except Exception as e:
        source["status"] = "error"
        source["error_message"] = str(e)[:300]

    sources = list(project.data_sources) if project.data_sources else []
    sources.append(source)
    store.update_project(project_id, data_sources=sources)
    return source


# --- 2. EXISTING FORM INGESTION ---
@app.post("/api/projects/{project_id}/sources/existing-form")
async def ingest_existing_form(project_id: str, data: dict):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data

    form_text = data.get("text", "")
    form_name = data.get("name", "Existing Form")
    form_url = data.get("url", "")

    if form_url and not form_text:
        form_text = _scrape_url_text(form_url)

    if not form_text:
        raise HTTPException(status_code=400, detail="Form text or URL is required")

    source = {
        "id": str(uuid.uuid4())[:8],
        "source_type": "existing_form",
        "name": form_name,
        "status": "processing",
        "url": form_url,
        "created_at": datetime.utcnow().isoformat(),
    }

    try:
        ai_prompt = f"""Analyze this existing government application form and extract all form fields, their types, and whether they are required.

For each field, provide:
- name: field label as shown on the form
- field_type: "text", "number", "date", "email", "select", "checkbox", "textarea"
- required: true or false
- options: if it's a select/dropdown, list the options
- section: which section of the form this belongs to

Also identify:
- form_name: what type of application/permit this form is for
- department: which department likely uses this form
- documents_mentioned: any required documents or attachments mentioned

Respond in JSON format.

Form Content:
{form_text[:8000]}"""

        ai_result = _extract_with_ai(ai_prompt, "", "form_field_extraction")
        if ai_result:
            try:
                import re
                json_match = re.search(r'\{.*\}', ai_result, re.DOTALL)
                if json_match:
                    extracted = json.loads(json_match.group())
                else:
                    extracted = {"raw_analysis": ai_result}
            except json.JSONDecodeError:
                extracted = {"raw_analysis": ai_result}
            source["extracted_data"] = extracted
        else:
            # Fallback: basic field pattern detection
            import re
            field_patterns = re.findall(r'([A-Z][A-Za-z\s/]+)[:_]{1,}', form_text[:5000])
            fields = []
            for f in field_patterns[:30]:
                name = f.strip()
                if len(name) > 2 and len(name) < 60:
                    ftype = "text"
                    if any(w in name.lower() for w in ["date", "when"]):
                        ftype = "date"
                    elif any(w in name.lower() for w in ["email", "e-mail"]):
                        ftype = "email"
                    elif any(w in name.lower() for w in ["number", "amount", "qty", "quantity", "#"]):
                        ftype = "number"
                    elif any(w in name.lower() for w in ["description", "explain", "comments", "notes"]):
                        ftype = "textarea"
                    fields.append({"name": name, "field_type": ftype, "required": True})
            source["extracted_data"] = {"form_name": form_name, "fields": fields}

        source["status"] = "completed"
    except Exception as e:
        source["status"] = "error"
        source["error_message"] = str(e)[:300]

    sources = list(project.data_sources) if project.data_sources else []
    sources.append(source)
    store.update_project(project_id, data_sources=sources)
    return source


# --- 3. FEE SCHEDULE PARSER ---
@app.post("/api/projects/{project_id}/sources/fee-schedule")
async def parse_fee_schedule(project_id: str, data: dict):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data

    fee_text = data.get("text", "")
    fee_url = data.get("url", "")
    fee_name = data.get("name", "Fee Schedule")

    if fee_url and not fee_text:
        fee_text = _scrape_url_text(fee_url)

    if not fee_text:
        raise HTTPException(status_code=400, detail="Fee schedule text or URL is required")

    source = {
        "id": str(uuid.uuid4())[:8],
        "source_type": "fee_schedule",
        "name": fee_name,
        "status": "processing",
        "url": fee_url,
        "created_at": datetime.utcnow().isoformat(),
    }

    try:
        ai_prompt = f"""Analyze this government fee schedule and extract all fees.

For each fee, provide:
- name: fee name
- amount: dollar amount (number only)
- fee_type: "flat", "calculated", "per_unit", "percentage", "deposit"
- applies_to: which permit/license/service this fee applies to
- conditions: any conditions or notes about when this fee applies
- formula: if calculated, the formula (e.g., "valuation * 0.01")

Respond in JSON format as a list of fee objects.

Fee Schedule:
{fee_text[:8000]}"""

        ai_result = _extract_with_ai(ai_prompt, "", "fee_schedule_extraction")
        if ai_result:
            try:
                import re
                json_match = re.search(r'\[.*\]', ai_result, re.DOTALL)
                if json_match:
                    extracted = json.loads(json_match.group())
                else:
                    extracted = [{"raw_analysis": ai_result}]
            except json.JSONDecodeError:
                extracted = [{"raw_analysis": ai_result}]
            source["extracted_data"] = {"fees": extracted}
        else:
            # Fallback: regex for dollar amounts
            import re
            fee_matches = re.findall(r'([A-Za-z][A-Za-z\s/()-]+?)\s*[\$:]?\s*\$\s*(\d+(?:,\d{3})*(?:\.\d{2})?)', fee_text[:5000])
            fees = []
            for name, amount in fee_matches[:30]:
                name = name.strip()
                if len(name) > 2 and len(name) < 80:
                    fees.append({
                        "name": name,
                        "amount": float(amount.replace(",", "")),
                        "fee_type": "flat",
                        "applies_to": "Unknown"
                    })
            source["extracted_data"] = {"fees": fees}

        source["status"] = "completed"
    except Exception as e:
        source["status"] = "error"
        source["error_message"] = str(e)[:300]

    sources = list(project.data_sources) if project.data_sources else []
    sources.append(source)
    store.update_project(project_id, data_sources=sources)
    return source


# --- 4. CROSS-SOURCE RECONCILIATION ---
@app.post("/api/projects/{project_id}/sources/reconcile")
async def reconcile_sources(project_id: str):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data

    if not project.configuration:
        raise HTTPException(status_code=400, detail="No configuration exists yet. Run analysis first.")

    config = project.configuration if isinstance(project.configuration, Configuration) else Configuration(**project.configuration)
    sources = project.data_sources or []
    completed_sources = [s for s in sources if (s.get("status") if isinstance(s, dict) else s.status) == "completed"]

    if not completed_sources:
        raise HTTPException(status_code=400, detail="No completed data sources to reconcile")

    # Gather all extracted data
    all_municipal_reqs = []
    all_form_fields = []
    all_fees = []
    source_summaries = []

    for src in completed_sources:
        s = src if isinstance(src, dict) else src.dict()
        ed = s.get("extracted_data") or {}
        stype = s.get("source_type", "")

        if stype == "municipal_code":
            reqs = ed.get("requirements", [])
            all_municipal_reqs.extend(reqs)
            source_summaries.append(f"Municipal Code ({s.get('url','')}): found {len(reqs)} requirements")
        elif stype == "existing_form":
            fields = ed.get("fields", [])
            all_form_fields.extend(fields)
            source_summaries.append(f"Form '{s.get('name','')}': found {len(fields)} fields")
        elif stype == "fee_schedule":
            fees = ed.get("fees", [])
            all_fees.extend(fees)
            source_summaries.append(f"Fee Schedule '{s.get('name','')}': found {len(fees)} fees")

    # Build current config summary
    config_rt_names = [rt.name.lower() for rt in config.record_types]
    config_summary = {
        "record_types": [{"name": rt.name, "fields": len(rt.form_fields), "fees": len(rt.fees),
                          "steps": len(rt.workflow_steps), "docs": len(rt.required_documents)}
                         for rt in config.record_types],
        "departments": [d.name for d in config.departments],
        "roles": [r.name for r in config.user_roles],
    }

    items = []

    # AI reconciliation
    ai_prompt = f"""You are an expert at configuring government PLC (Permitting, Licensing & Code Enforcement) systems.

Compare these data sources against the current configuration and identify gaps, conflicts, and enrichment opportunities.

CURRENT CONFIGURATION:
{json.dumps(config_summary, indent=2)}

DATA SOURCES COLLECTED:
{chr(10).join(source_summaries)}

MUNICIPAL CODE REQUIREMENTS:
{json.dumps(all_municipal_reqs[:20], indent=2)}

FORM FIELDS FROM EXISTING FORMS:
{json.dumps(all_form_fields[:30], indent=2)}

FEE SCHEDULE DATA:
{json.dumps(all_fees[:20], indent=2)}

For each finding, provide:
- action: "add" (missing from config), "update" (exists but incomplete), or "flag" (potential conflict)
- target: "record_type", "fee", "form_field", "document", "workflow_step", or "department"
- record_type_name: which record type this relates to (if applicable)
- confidence: 0.0 to 1.0
- title: short title
- description: detailed explanation
- suggested_data: specific data to add/update (as JSON object)

Respond as a JSON array of findings. Focus on the most impactful items first."""

    ai_result = _extract_with_ai(ai_prompt, "", "reconciliation_analysis")
    if ai_result:
        try:
            import re
            json_match = re.search(r'\[.*\]', ai_result, re.DOTALL)
            if json_match:
                ai_items = json.loads(json_match.group())
                for item in ai_items[:25]:
                    items.append({
                        "id": str(uuid.uuid4())[:8],
                        "action": item.get("action", "flag"),
                        "target": item.get("target", "record_type"),
                        "target_id": "",
                        "record_type_name": item.get("record_type_name", ""),
                        "confidence": item.get("confidence", 0.5),
                        "source_ids": [s.get("id", "") for s in completed_sources],
                        "title": item.get("title", "Finding"),
                        "description": item.get("description", ""),
                        "suggested_data": item.get("suggested_data"),
                        "status": "pending"
                    })
        except (json.JSONDecodeError, Exception):
            pass

    # Fallback / supplemental: rule-based reconciliation
    # Check municipal code requirements against existing record types
    for req in all_municipal_reqs:
        if isinstance(req, dict) and req.get("name"):
            req_name_lower = req["name"].lower()
            if not any(req_name_lower in rn or rn in req_name_lower for rn in config_rt_names):
                items.append({
                    "id": str(uuid.uuid4())[:8],
                    "action": "add",
                    "target": "record_type",
                    "target_id": "",
                    "record_type_name": req.get("name", ""),
                    "confidence": 0.7,
                    "source_ids": [],
                    "title": f"Missing Record Type: {req.get('name', '')}",
                    "description": f"Municipal code references '{req.get('name','')}' but no matching record type exists. {req.get('description','')}",
                    "suggested_data": {"name": req.get("name", ""), "description": req.get("description", ""), "category": req.get("type", "permit").title()},
                    "status": "pending"
                })

    # Check fee schedule against existing fees
    for fee in all_fees:
        if isinstance(fee, dict) and fee.get("name"):
            applies_to = fee.get("applies_to", "").lower()
            for rt in config.record_types:
                if applies_to and (applies_to in rt.name.lower() or rt.name.lower() in applies_to):
                    existing_fee_names = [f.name.lower() for f in rt.fees]
                    if fee["name"].lower() not in existing_fee_names:
                        items.append({
                            "id": str(uuid.uuid4())[:8],
                            "action": "add",
                            "target": "fee",
                            "target_id": rt.id,
                            "record_type_name": rt.name,
                            "confidence": 0.8,
                            "source_ids": [],
                            "title": f"Missing Fee: {fee['name']} on {rt.name}",
                            "description": f"Fee schedule lists '{fee['name']}' (${fee.get('amount', 0):.2f}) for {rt.name} but it's not in the configuration.",
                            "suggested_data": {"name": fee["name"], "amount": fee.get("amount", 0), "fee_type": fee.get("fee_type", "flat")},
                            "status": "pending"
                        })

    # Deduplicate by title
    seen_titles = set()
    unique_items = []
    for item in items:
        if item["title"] not in seen_titles:
            seen_titles.add(item["title"])
            unique_items.append(item)

    store.update_project(project_id, reconciliation_items=unique_items[:30])
    return {"items": unique_items[:30], "source_count": len(completed_sources)}


# --- 5. PEER CITY TEMPLATES ---
@app.get("/api/templates/peer-cities")
async def list_peer_city_templates(search: str = ""):
    results = []
    for t in PEER_CITY_TEMPLATES:
        if search:
            search_lower = search.lower()
            if (search_lower in t["name"].lower() or
                search_lower in t["description"].lower() or
                any(search_lower in tag for tag in t["tags"])):
                results.append({k: v for k, v in t.items() if k != "record_types" and k != "departments" and k != "user_roles"})
        else:
            results.append({k: v for k, v in t.items() if k != "record_types" and k != "departments" and k != "user_roles"})
    return {"templates": results}


@app.get("/api/templates/peer-cities/{template_id}")
async def get_peer_city_template(template_id: str):
    for t in PEER_CITY_TEMPLATES:
        if t["id"] == template_id:
            return t
    raise HTTPException(status_code=404, detail="Template not found")


@app.post("/api/projects/{project_id}/sources/apply-template")
async def apply_peer_template(project_id: str, data: dict):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data

    template_id = data.get("template_id", "")
    merge_mode = data.get("mode", "merge")  # "merge" or "replace"

    template = None
    for t in PEER_CITY_TEMPLATES:
        if t["id"] == template_id:
            template = t
            break

    if not template:
        raise HTTPException(status_code=404, detail="Template not found")

    # Build configuration from template
    record_types = []
    for rt_data in template["record_types"]:
        rt = RecordType(
            name=rt_data["name"],
            category=rt_data.get("category", ""),
            description=rt_data.get("description", ""),
            form_fields=[FormField(
                name=f["name"], field_type=f.get("field_type", "text"),
                required=f.get("required", True),
                options=f.get("options"),
            ) for f in rt_data.get("form_fields", [])],
            fees=[Fee(
                name=f["name"], amount=f.get("amount", 0),
                fee_type=f.get("fee_type", "flat"),
                when_applied="submission", formula=f.get("formula", ""),
            ) for f in rt_data.get("fees", [])],
            workflow_steps=[WorkflowStep(
                name=s["name"], order=s["order"],
                assigned_role=s.get("assigned_role", ""),
                status_to=s["name"].lower().replace(" ", "_"),
            ) for s in rt_data.get("workflow_steps", [])],
            required_documents=[RequiredDocument(
                name=d["name"], required=d.get("required", True),
                stage=d.get("stage", "submission"),
                description=d.get("description", ""),
            ) for d in rt_data.get("required_documents", [])],
        )
        record_types.append(rt)

    departments = [Department(
        name=d["name"], description=d.get("description", "")
    ) for d in template.get("departments", [])]

    user_roles = [UserRole(
        name=r["name"], description=r.get("description", ""),
        permissions=r.get("permissions", []),
    ) for r in template.get("user_roles", [])]

    if merge_mode == "replace" or not project.configuration:
        config = Configuration(
            record_types=record_types,
            departments=departments,
            user_roles=user_roles,
            generated_at=datetime.utcnow(),
            summary=f"Configuration from template: {template['name']}"
        )
    else:
        existing = project.configuration if isinstance(project.configuration, Configuration) else Configuration(**project.configuration)
        existing_rt_names = {rt.name.lower() for rt in existing.record_types}
        for rt in record_types:
            if rt.name.lower() not in existing_rt_names:
                existing.record_types.append(rt)
        existing_dept_names = {d.name.lower() for d in existing.departments}
        for d in departments:
            if d.name.lower() not in existing_dept_names:
                existing.departments.append(d)
        existing_role_names = {r.name.lower() for r in existing.user_roles}
        for r in user_roles:
            if r.name.lower() not in existing_role_names:
                existing.user_roles.append(r)
        existing.summary += f" | Merged with template: {template['name']}"
        config = existing

    store.save_configuration(project_id, config)

    # Track as data source
    source = {
        "id": str(uuid.uuid4())[:8],
        "source_type": "peer_template",
        "name": f"Template: {template['name']}",
        "status": "completed",
        "created_at": datetime.utcnow().isoformat(),
        "extracted_data": {"template_id": template_id, "mode": merge_mode,
                           "record_types_added": len(record_types),
                           "departments_added": len(departments),
                           "roles_added": len(user_roles)},
    }
    sources = list(project.data_sources) if project.data_sources else []
    sources.append(source)
    store.update_project(project_id, data_sources=sources)

    return {"message": f"Template '{template['name']}' applied ({merge_mode})",
            "record_types": len(config.record_types),
            "departments": len(config.departments),
            "user_roles": len(config.user_roles)}


# --- 6. VALIDATION AGENT ---
@app.post("/api/projects/{project_id}/validate")
async def validate_configuration(project_id: str):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data

    if not project.configuration:
        raise HTTPException(status_code=400, detail="No configuration to validate")

    config = project.configuration if isinstance(project.configuration, Configuration) else Configuration(**project.configuration)
    findings = []

    # --- RULE-BASED VALIDATION ---

    # 1. Record type completeness
    for rt in config.record_types:
        if not rt.form_fields:
            findings.append({
                "id": str(uuid.uuid4())[:8], "severity": "critical", "category": "completeness",
                "title": f"{rt.name}: No form fields defined",
                "description": f"Record type '{rt.name}' has no form fields. Applicants won't be able to submit any data.",
                "record_type_id": rt.id,
                "recommendation": "Add form fields that capture the information needed for this application type.",
                "auto_fixable": False
            })
        elif len(rt.form_fields) < 3:
            findings.append({
                "id": str(uuid.uuid4())[:8], "severity": "warning", "category": "completeness",
                "title": f"{rt.name}: Very few form fields ({len(rt.form_fields)})",
                "description": f"Record type '{rt.name}' only has {len(rt.form_fields)} form fields. Most application types need at least 5-8 fields.",
                "record_type_id": rt.id,
                "recommendation": "Consider adding fields for applicant contact info, project details, and property information.",
                "auto_fixable": False
            })

        if not rt.workflow_steps:
            findings.append({
                "id": str(uuid.uuid4())[:8], "severity": "critical", "category": "workflow",
                "title": f"{rt.name}: No workflow steps defined",
                "description": f"Record type '{rt.name}' has no workflow. Applications will have no review or approval process.",
                "record_type_id": rt.id,
                "recommendation": "Add at least: Submission -> Review -> Approval workflow steps.",
                "auto_fixable": True,
                "fix_data": {"workflow_steps": [
                    {"name": "Application Submitted", "order": 1, "assigned_role": "Clerk", "status_to": "submitted"},
                    {"name": "Under Review", "order": 2, "assigned_role": "Reviewer", "status_to": "under_review"},
                    {"name": "Approved", "order": 3, "assigned_role": "Supervisor", "status_to": "approved"},
                ]}
            })
        else:
            # Check for unassigned workflow steps
            unassigned = [s for s in rt.workflow_steps if not s.assigned_role]
            if unassigned:
                findings.append({
                    "id": str(uuid.uuid4())[:8], "severity": "warning", "category": "workflow",
                    "title": f"{rt.name}: {len(unassigned)} workflow steps have no assigned role",
                    "description": f"Steps without assigned roles: {', '.join([s.name for s in unassigned])}. These won't route to anyone.",
                    "record_type_id": rt.id,
                    "recommendation": "Assign a user role to each workflow step so applications get routed correctly.",
                    "auto_fixable": False
                })

            # Check workflow has proper start and end
            steps_sorted = sorted(rt.workflow_steps, key=lambda s: s.order)
            first_step = steps_sorted[0].name.lower()
            last_step = steps_sorted[-1].name.lower()
            if not any(w in first_step for w in ["submit", "receive", "intake", "filed", "application"]):
                findings.append({
                    "id": str(uuid.uuid4())[:8], "severity": "info", "category": "best_practice",
                    "title": f"{rt.name}: First workflow step may not be intake",
                    "description": f"First step is '{steps_sorted[0].name}'. Best practice is to start with an intake/submission step.",
                    "record_type_id": rt.id,
                    "recommendation": "Consider renaming or reordering so the first step clearly represents application intake.",
                    "auto_fixable": False
                })

        # Fee validation
        if not rt.fees and rt.category and rt.category.lower() not in ["code enforcement", "enforcement", "complaint"]:
            findings.append({
                "id": str(uuid.uuid4())[:8], "severity": "warning", "category": "fees",
                "title": f"{rt.name}: No fees configured",
                "description": f"Record type '{rt.name}' has no fees. Most permit and license types require at least an application fee.",
                "record_type_id": rt.id,
                "recommendation": "Add applicable fees (application fee, review fee, inspection fee, etc.).",
                "auto_fixable": False
            })

        # Zero-amount fees
        zero_fees = [f for f in rt.fees if f.amount == 0 and f.fee_type == "flat"]
        if zero_fees:
            findings.append({
                "id": str(uuid.uuid4())[:8], "severity": "info", "category": "fees",
                "title": f"{rt.name}: {len(zero_fees)} fees have $0.00 amount",
                "description": f"Fees with zero amount: {', '.join([f.name for f in zero_fees])}",
                "record_type_id": rt.id,
                "recommendation": "Verify these fees are intentionally $0.00 or update with correct amounts.",
                "auto_fixable": False
            })

        # Document requirements
        if not rt.required_documents and rt.category and rt.category.lower() in ["building", "planning"]:
            findings.append({
                "id": str(uuid.uuid4())[:8], "severity": "warning", "category": "documents",
                "title": f"{rt.name}: No required documents",
                "description": f"Building/planning record types typically require supporting documents (plans, reports, etc.).",
                "record_type_id": rt.id,
                "recommendation": "Add required documents such as site plans, construction plans, or engineering reports.",
                "auto_fixable": False
            })

        # Required fields check - common fields every app should have
        if rt.form_fields:
            field_names_lower = [f.name.lower() for f in rt.form_fields]
            essential_fields = {
                "address": ["address", "property address", "location", "site address"],
                "applicant": ["applicant", "owner", "owner name", "applicant name", "contact name"],
            }
            for field_group, variations in essential_fields.items():
                if not any(any(v in fn for v in variations) for fn in field_names_lower):
                    findings.append({
                        "id": str(uuid.uuid4())[:8], "severity": "warning", "category": "completeness",
                        "title": f"{rt.name}: Missing {field_group} field",
                        "description": f"No field for '{field_group}' was found. Most applications need this information.",
                        "record_type_id": rt.id,
                        "recommendation": f"Add a field for {field_group} information.",
                        "auto_fixable": True,
                        "fix_data": {"add_field": {"name": field_group.title(), "field_type": "text", "required": True}}
                    })

    # 2. Department validation
    if not config.departments:
        findings.append({
            "id": str(uuid.uuid4())[:8], "severity": "warning", "category": "completeness",
            "title": "No departments configured",
            "description": "No departments have been set up. Departments help organize workflow routing and reporting.",
            "record_type_id": "",
            "recommendation": "Add departments that reflect your organization's structure (e.g., Building, Planning, Code Enforcement).",
            "auto_fixable": False
        })

    # 3. Role validation
    if not config.user_roles:
        findings.append({
            "id": str(uuid.uuid4())[:8], "severity": "critical", "category": "completeness",
            "title": "No user roles configured",
            "description": "No user roles have been defined. Without roles, workflow steps cannot be assigned to staff.",
            "record_type_id": "",
            "recommendation": "Add user roles (e.g., Administrator, Reviewer, Inspector, Clerk, Public User).",
            "auto_fixable": True,
            "fix_data": {"roles": [
                {"name": "Administrator", "description": "Full system access", "permissions": ["manage_all"]},
                {"name": "Reviewer", "description": "Reviews and approves applications", "permissions": ["review", "approve"]},
                {"name": "Clerk", "description": "Application intake and processing", "permissions": ["create", "update"]},
            ]}
        })
    else:
        # Check for admin role
        has_admin = any("admin" in r.name.lower() for r in config.user_roles)
        if not has_admin:
            findings.append({
                "id": str(uuid.uuid4())[:8], "severity": "warning", "category": "best_practice",
                "title": "No administrator role found",
                "description": "No role with 'admin' or 'administrator' in the name was found.",
                "record_type_id": "",
                "recommendation": "Add an Administrator role with full system access for configuration management.",
                "auto_fixable": False
            })

        # Check workflow roles exist in user_roles
        role_names_lower = {r.name.lower() for r in config.user_roles}
        for rt in config.record_types:
            for ws in rt.workflow_steps:
                if ws.assigned_role and ws.assigned_role.lower() not in role_names_lower:
                    findings.append({
                        "id": str(uuid.uuid4())[:8], "severity": "warning", "category": "workflow",
                        "title": f"Workflow role '{ws.assigned_role}' not in user roles",
                        "description": f"Step '{ws.name}' in {rt.name} is assigned to '{ws.assigned_role}' but this role doesn't exist in the configured user roles.",
                        "record_type_id": rt.id,
                        "recommendation": f"Either add '{ws.assigned_role}' as a user role or update the workflow step assignment.",
                        "auto_fixable": False
                    })

    # 4. Overall health score
    critical_count = sum(1 for f in findings if f.get("severity") == "critical")
    warning_count = sum(1 for f in findings if f.get("severity") == "warning")
    info_count = sum(1 for f in findings if f.get("severity") == "info")

    score = max(0, 100 - (critical_count * 15) - (warning_count * 5) - (info_count * 1))

    # Add success findings
    if config.record_types:
        findings.insert(0, {
            "id": str(uuid.uuid4())[:8], "severity": "success", "category": "completeness",
            "title": f"{len(config.record_types)} record types configured",
            "description": f"Your configuration includes {len(config.record_types)} record types covering your PLC processes.",
            "record_type_id": "", "recommendation": "", "auto_fixable": False
        })
    if config.departments:
        findings.insert(0, {
            "id": str(uuid.uuid4())[:8], "severity": "success", "category": "completeness",
            "title": f"{len(config.departments)} departments set up",
            "description": "Departments are configured for organizational structure.",
            "record_type_id": "", "recommendation": "", "auto_fixable": False
        })

    # AI-powered additional insights
    if ANTHROPIC_AVAILABLE and config.record_types:
        config_summary = json.dumps({
            "record_types": [{"name": rt.name, "category": rt.category, "fields": len(rt.form_fields),
                              "fees": len(rt.fees), "workflow_steps": len(rt.workflow_steps)}
                             for rt in config.record_types],
            "departments": [d.name for d in config.departments],
            "roles": [r.name for r in config.user_roles],
        })
        ai_prompt = f"""Review this government PLC configuration and provide 3-5 best practice recommendations.
Focus on common gaps that implementation consultants typically find.

Configuration:
{config_summary}

Provide each recommendation as JSON with: severity ("info" or "warning"), category ("best_practice"), title, description, recommendation.
Return as a JSON array."""

        ai_result = _extract_with_ai(ai_prompt, "", "validation_recommendations")
        if ai_result:
            try:
                import re
                json_match = re.search(r'\[.*\]', ai_result, re.DOTALL)
                if json_match:
                    ai_findings = json.loads(json_match.group())
                    for af in ai_findings[:5]:
                        findings.append({
                            "id": str(uuid.uuid4())[:8],
                            "severity": af.get("severity", "info"),
                            "category": "best_practice",
                            "title": af.get("title", "AI Recommendation"),
                            "description": af.get("description", ""),
                            "record_type_id": "",
                            "recommendation": af.get("recommendation", ""),
                            "auto_fixable": False
                        })
            except (json.JSONDecodeError, Exception):
                pass

    store.update_project(project_id, validation_findings=findings)
    return {"findings": findings, "score": score,
            "summary": {"critical": critical_count, "warning": warning_count, "info": info_count,
                         "total": len(findings), "score": score}}


# --- DATA SOURCES LIST ---
@app.get("/api/projects/{project_id}/sources")
async def list_data_sources(project_id: str):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data
    return {
        "sources": project.data_sources or [],
        "reconciliation_items": project.reconciliation_items or [],
        "validation_findings": project.validation_findings or [],
    }


@app.delete("/api/projects/{project_id}/sources/{source_id}")
async def delete_data_source(project_id: str, source_id: str):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data
    sources = [s for s in (project.data_sources or []) if (s.get("id") if isinstance(s, dict) else s.id) != source_id]
    store.update_project(project_id, data_sources=sources)
    return {"message": "Source deleted"}


@app.post("/api/projects/{project_id}/reconciliation/{item_id}/accept")
async def accept_reconciliation(project_id: str, item_id: str):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data

    items = project.reconciliation_items or []
    for item in items:
        i = item if isinstance(item, dict) else item.dict()
        if i.get("id") == item_id:
            i["status"] = "accepted"

            # Apply the suggested data
            if i.get("suggested_data") and project.configuration:
                config = project.configuration if isinstance(project.configuration, Configuration) else Configuration(**project.configuration)
                if i["target"] == "record_type" and i["action"] == "add":
                    sd = i["suggested_data"]
                    new_rt = RecordType(name=sd.get("name", "New Record Type"),
                                        description=sd.get("description", ""),
                                        category=sd.get("category", ""))
                    config.record_types.append(new_rt)
                    store.save_configuration(project_id, config)
                elif i["target"] == "fee" and i.get("target_id"):
                    sd = i["suggested_data"]
                    for rt in config.record_types:
                        if rt.id == i["target_id"]:
                            new_fee = Fee(name=sd.get("name", ""), amount=sd.get("amount", 0),
                                          fee_type=sd.get("fee_type", "flat"), when_applied="submission")
                            rt.fees.append(new_fee)
                            store.save_configuration(project_id, config)
                            break

    store.update_project(project_id, reconciliation_items=items)
    return {"message": "Recommendation accepted and applied"}


@app.post("/api/projects/{project_id}/reconciliation/{item_id}/reject")
async def reject_reconciliation(project_id: str, item_id: str):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data
    items = project.reconciliation_items or []
    for item in items:
        i = item if isinstance(item, dict) else item.dict()
        if i.get("id") == item_id:
            i["status"] = "rejected"
    store.update_project(project_id, reconciliation_items=items)
    return {"message": "Recommendation rejected"}


@app.post("/api/projects/{project_id}/validate/auto-fix/{finding_id}")
async def auto_fix_finding(project_id: str, finding_id: str):
    project_data = store.get_project(project_id)
    if not project_data:
        raise HTTPException(status_code=404, detail="Project not found")
    project = Project(**project_data) if isinstance(project_data, dict) else project_data

    if not project.configuration:
        raise HTTPException(status_code=400, detail="No configuration")

    config = project.configuration if isinstance(project.configuration, Configuration) else Configuration(**project.configuration)
    findings = project.validation_findings or []

    for f in findings:
        fi = f if isinstance(f, dict) else f.dict()
        if fi.get("id") == finding_id and fi.get("auto_fixable") and fi.get("fix_data"):
            fd = fi["fix_data"]

            if "workflow_steps" in fd and fi.get("record_type_id"):
                for rt in config.record_types:
                    if rt.id == fi["record_type_id"]:
                        rt.workflow_steps = [WorkflowStep(
                            name=s["name"], order=s["order"],
                            assigned_role=s.get("assigned_role", ""),
                            status_to=s.get("status_to", ""),
                        ) for s in fd["workflow_steps"]]
                        break

            if "roles" in fd:
                for r in fd["roles"]:
                    config.user_roles.append(UserRole(
                        name=r["name"], description=r.get("description", ""),
                        permissions=r.get("permissions", []),
                    ))

            if "add_field" in fd and fi.get("record_type_id"):
                for rt in config.record_types:
                    if rt.id == fi["record_type_id"]:
                        af = fd["add_field"]
                        rt.form_fields.append(FormField(
                            name=af["name"], field_type=af.get("field_type", "text"),
                            required=af.get("required", True),
                        ))
                        break

            fi["severity"] = "success"
            fi["title"] = f"[FIXED] {fi['title']}"
            store.save_configuration(project_id, config)
            break

    store.update_project(project_id, validation_findings=findings)
    return {"message": "Auto-fix applied"}


# ============================================================================
# CITY PREVIEW - Fetch metadata from community website
# ============================================================================

@app.post("/api/city-preview")
async def get_city_preview(data: dict):
    """Fetch metadata from a city/community website URL for personalization."""
    import re
    from urllib.request import urlopen, Request
    from urllib.error import URLError

    url = data.get("url", "").strip()
    if not url:
        raise HTTPException(status_code=400, detail="URL is required")

    # Ensure URL has scheme
    if not url.startswith("http"):
        url = "https://" + url

    result = {
        "url": url,
        "city_name": "",
        "title": "",
        "description": "",
        "og_image": "",
        "favicon": "",
        "theme_color": "",
    }

    try:
        req = Request(url, headers={
            "User-Agent": "Mozilla/5.0 (compatible; PLCAutoConfig/1.0)"
        })
        with urlopen(req, timeout=8) as resp:
            html = resp.read().decode("utf-8", errors="ignore")[:50000]  # limit

        # Extract title
        title_match = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
        if title_match:
            result["title"] = title_match.group(1).strip()[:200]

        # Extract meta description
        desc_match = re.search(r'<meta[^>]+name=["\']description["\'][^>]+content=["\']([^"\']*)["\']', html, re.IGNORECASE)
        if not desc_match:
            desc_match = re.search(r'<meta[^>]+content=["\']([^"\']*)["\'][^>]+name=["\']description["\']', html, re.IGNORECASE)
        if desc_match:
            result["description"] = desc_match.group(1).strip()[:500]

        # Extract og:image
        og_match = re.search(r'<meta[^>]+property=["\']og:image["\'][^>]+content=["\']([^"\']*)["\']', html, re.IGNORECASE)
        if not og_match:
            og_match = re.search(r'<meta[^>]+content=["\']([^"\']*)["\'][^>]+property=["\']og:image["\']', html, re.IGNORECASE)
        if og_match:
            img_url = og_match.group(1).strip()
            if img_url.startswith("/"):
                from urllib.parse import urlparse
                parsed = urlparse(url)
                img_url = f"{parsed.scheme}://{parsed.netloc}{img_url}"
            result["og_image"] = img_url

        # Extract favicon
        fav_match = re.search(r'<link[^>]+rel=["\'](?:shortcut )?icon["\'][^>]+href=["\']([^"\']*)["\']', html, re.IGNORECASE)
        if not fav_match:
            fav_match = re.search(r'<link[^>]+href=["\']([^"\']*)["\'][^>]+rel=["\'](?:shortcut )?icon["\']', html, re.IGNORECASE)
        if fav_match:
            fav_url = fav_match.group(1).strip()
            if fav_url.startswith("/"):
                from urllib.parse import urlparse
                parsed = urlparse(url)
                fav_url = f"{parsed.scheme}://{parsed.netloc}{fav_url}"
            result["favicon"] = fav_url
        else:
            from urllib.parse import urlparse
            parsed = urlparse(url)
            result["favicon"] = f"{parsed.scheme}://{parsed.netloc}/favicon.ico"

        # Extract theme-color
        theme_match = re.search(r'<meta[^>]+name=["\']theme-color["\'][^>]+content=["\']([^"\']*)["\']', html, re.IGNORECASE)
        if theme_match:
            result["theme_color"] = theme_match.group(1).strip()

        # Try to derive city name from title
        title = result["title"]
        city_name = title
        # Common patterns: "City of X", "X, State", "Welcome to X"
        for pattern in [
            r"(?:City of |Town of |Village of )([A-Za-z\s]+)",
            r"Welcome to (?:the )?(?:City of |Town of )?([A-Za-z\s]+)",
            r"^([A-Za-z\s]+?)(?:\s*[-|,])",
        ]:
            m = re.search(pattern, title, re.IGNORECASE)
            if m:
                city_name = m.group(1).strip()
                break
        result["city_name"] = city_name[:100]

    except (URLError, Exception) as e:
        # If we can't fetch, try to derive city name from URL
        from urllib.parse import urlparse
        parsed = urlparse(url)
        domain = parsed.netloc.replace("www.", "")
        domain_parts = domain.split(".")
        if domain_parts:
            city_guess = domain_parts[0].replace("-", " ").title()
            result["city_name"] = city_guess
        result["description"] = f"Could not fetch site details: {str(e)[:100]}"

    return result


# ============================================================================
# AI CONSULTANT - Multi-agent Q&A for project questions
# ============================================================================

@app.post("/api/projects/{project_id}/consultant/ask")
async def consultant_ask(project_id: str, data: dict):
    """
    AI Consultant: orchestrates multiple agents to answer questions about
    the project's configuration, uploaded data, community, and best practices.
    """
    question = data.get("question", "").strip()
    history = data.get("history", [])

    if not question:
        raise HTTPException(status_code=400, detail="Question is required")

    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    # Gather context from all agents
    agents_consulted = []
    sources = []
    context_parts = []

    # Agent 1: Configuration Agent - search project configuration
    if project.configuration:
        agents_consulted.append("configuration")
        config = project.configuration
        config_context = []

        if config.record_types:
            rt_summaries = []
            for rt in config.record_types:
                rt_info = f"- **{rt.name}**"
                if rt.category:
                    rt_info += f" (Category: {rt.category})"
                if rt.description:
                    rt_info += f": {rt.description[:150]}"

                fields = [f.name for f in rt.form_fields[:5]] if rt.form_fields else []
                if fields:
                    rt_info += f"\n  Form fields: {', '.join(fields)}"

                fees = [f"{f.name}: ${f.amount:.2f}" for f in rt.fees[:3]] if rt.fees else []
                if fees:
                    rt_info += f"\n  Fees: {', '.join(fees)}"

                steps = sorted(rt.workflow_steps, key=lambda s: s.order)[:5] if rt.workflow_steps else []
                if steps:
                    step_names = [s.name for s in steps]
                    rt_info += f"\n  Workflow: {' → '.join(step_names)}"

                docs = [d.name for d in rt.required_documents[:3]] if rt.required_documents else []
                if docs:
                    rt_info += f"\n  Required docs: {', '.join(docs)}"

                rt_summaries.append(rt_info)

            config_context.append(f"Record Types ({len(config.record_types)} total):\n" + "\n".join(rt_summaries))
            sources.append(f"Project configuration: {len(config.record_types)} record types")

        if config.departments:
            dept_info = ", ".join([f"{d.name}" + (f" ({d.description[:50]})" if d.description else "") for d in config.departments])
            config_context.append(f"Departments: {dept_info}")
            sources.append(f"Departments: {len(config.departments)} configured")

        if config.user_roles:
            role_info = ", ".join([f"{r.name}" + (f" - {r.description[:50]}" if r.description else "") for r in config.user_roles])
            config_context.append(f"User Roles: {role_info}")
            sources.append(f"User roles: {len(config.user_roles)} configured")

        # Note: conditional_rules not yet in Configuration model

        if config_context:
            context_parts.append("=== CONFIGURATION DATA ===\n" + "\n\n".join(config_context))

    # Agent 2: Document Agent - search uploaded file data
    if project.uploaded_files:
        agents_consulted.append("documents")
        file_info = []
        for uf in project.uploaded_files:
            file_info.append(f"- {uf.filename}: {uf.rows_count} rows, columns: {', '.join(uf.columns[:8])}")
        context_parts.append("=== UPLOADED DATA ===\n" + "\n".join(file_info))
        sources.append(f"Uploaded files: {len(project.uploaded_files)} files analyzed")

    # Agent 3: Community Agent - community research data
    if project.community_research:
        agents_consulted.append("community")
        research = project.community_research
        if isinstance(research, str) and len(research) > 10:
            context_parts.append(f"=== COMMUNITY RESEARCH ===\n{research[:2000]}")
            sources.append(f"Community research: {project.community_url or 'website analysis'}")
        elif isinstance(research, dict):
            research_text = json.dumps(research, indent=2)[:2000]
            context_parts.append(f"=== COMMUNITY RESEARCH ===\n{research_text}")
            sources.append(f"Community research: {project.community_url or 'website analysis'}")

    # Agent 4: Best Practices Agent - always available
    agents_consulted.append("best_practices")
    best_practices = """=== OPENGOV PLC BEST PRACTICES ===
- Record types should have clear naming conventions matching the community's terminology
- Each record type should have at minimum: application form fields, fee schedule, workflow steps, and required documents
- Workflows should follow the principle of least-touch: auto-route to the right department based on record type
- Fee schedules should include both flat fees and calculated fees where applicable
- Conditional logic can automate routing, required fields, and fee calculations
- Departments should map to the community's actual organizational structure
- User roles should follow principle of least privilege with clear separation of duties
- Required documents should specify accepted formats and file size limits
- Community research helps ensure the configuration matches existing local ordinances and processes
- Staff training (LMS materials) should be generated after configuration is finalized
- Plan reviews and inspections should be separate workflow steps with assigned roles
- Public-facing portals should show real-time status updates for applicants
- Code enforcement workflows should include violation types, inspection scheduling, and citation management"""
    context_parts.append(best_practices)
    sources.append("OpenGov PLC best practices knowledge base")

    # Build the full context
    full_context = "\n\n".join(context_parts)

    # Generate answer
    answer = ""

    if ANTHROPIC_AVAILABLE:
        try:
            client = Anthropic()
            system_prompt = f"""You are an expert OpenGov PLC consultant helping configure Permitting, Licensing & Code Enforcement systems for {project.customer_name or 'a community'}.

You have access to the following context about this project:

{full_context}

Instructions:
- Answer the user's question thoroughly using the provided context
- Reference specific configuration details (record types, fees, workflows) when relevant
- Provide actionable recommendations based on best practices
- If information is missing from the context, say so and suggest what data might help
- Use **bold** for key terms and bullet points for clarity
- Keep responses focused and practical
- If asked about processes, describe the current configuration AND suggest improvements"""

            api_messages = []
            for h in history[-6:]:
                api_messages.append({"role": h["role"], "content": h["content"]})
            api_messages.append({"role": "user", "content": question})

            response = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=1500,
                system=system_prompt,
                messages=api_messages,
            )
            answer = response.content[0].text
        except Exception as e:
            answer = _generate_fallback_answer(question, project, full_context)
    else:
        answer = _generate_fallback_answer(question, project, full_context)

    return {
        "answer": answer,
        "sources": sources,
        "agents_consulted": agents_consulted,
    }


def _generate_fallback_answer(question: str, project, context: str) -> str:
    """Generate a structured answer without AI when Claude API is not available."""
    q_lower = question.lower()
    parts = []
    config = project.configuration

    if not config:
        return "This project doesn't have a configuration yet. Please upload data and run the analysis first to generate a configuration that I can help answer questions about."

    # Check what the question is about and build relevant answer
    if any(w in q_lower for w in ["permit", "record type", "license", "enforcement", "type"]):
        if config.record_types:
            parts.append(f"**{project.customer_name} has {len(config.record_types)} record types configured:**\n")
            for rt in config.record_types:
                parts.append(f"**{rt.name}**" + (f" ({rt.category})" if rt.category else ""))
                if rt.description:
                    parts.append(f"  {rt.description[:200]}")
                if rt.fees:
                    fee_list = ", ".join([f"{f.name}: ${f.amount:.2f}" for f in rt.fees[:3]])
                    parts.append(f"  Fees: {fee_list}")
                if rt.workflow_steps:
                    steps = sorted(rt.workflow_steps, key=lambda s: s.order)
                    step_names = [s.name for s in steps[:5]]
                    parts.append(f"  Workflow: {' → '.join(step_names)}")
                parts.append("")

    elif any(w in q_lower for w in ["fee", "cost", "price", "charge"]):
        if config.record_types:
            parts.append("**Fee Schedule:**\n")
            for rt in config.record_types:
                if rt.fees:
                    parts.append(f"**{rt.name}:**")
                    for fee in rt.fees:
                        parts.append(f"  - {fee.name}: ${fee.amount:.2f} ({fee.fee_type})")
                    parts.append("")

    elif any(w in q_lower for w in ["workflow", "process", "step", "approval"]):
        if config.record_types:
            parts.append("**Workflow Processes:**\n")
            for rt in config.record_types[:5]:
                if rt.workflow_steps:
                    steps = sorted(rt.workflow_steps, key=lambda s: s.order)
                    parts.append(f"**{rt.name}:**")
                    for s in steps:
                        assigned = f" (Assigned to: {s.assigned_role})" if s.assigned_role else ""
                        parts.append(f"  {s.order}. {s.name}{assigned}")
                    parts.append("")

    elif any(w in q_lower for w in ["department", "org", "team"]):
        if config.departments:
            parts.append("**Departments:**\n")
            for d in config.departments:
                parts.append(f"- **{d.name}**" + (f": {d.description}" if d.description else ""))

    elif any(w in q_lower for w in ["role", "user", "permission", "access"]):
        if config.user_roles:
            parts.append("**User Roles:**\n")
            for r in config.user_roles:
                parts.append(f"- **{r.name}**" + (f": {r.description}" if r.description else ""))
                if r.permissions:
                    parts.append(f"  Permissions: {', '.join(r.permissions[:5])}")

    elif any(w in q_lower for w in ["document", "require", "submit", "upload"]):
        if config.record_types:
            parts.append("**Required Documents:**\n")
            for rt in config.record_types[:5]:
                if rt.required_documents:
                    parts.append(f"**{rt.name}:**")
                    for doc in rt.required_documents:
                        parts.append(f"  - {doc.name}" + (f": {doc.description}" if doc.description else ""))
                    parts.append("")

    elif any(w in q_lower for w in ["best practice", "recommend", "suggestion", "improve", "conditional"]):
        parts.append("**OpenGov PLC Best Practices:**\n")
        parts.append("- Ensure each record type has complete fee schedules, workflows, and required documents")
        parts.append("- Use conditional logic to automate routing based on application type and value")
        parts.append("- Map departments to the community's actual organizational structure")
        parts.append("- Configure user roles with principle of least privilege")
        parts.append("- Generate LMS training materials after finalizing configuration")
        parts.append("- Review community research to align with existing local ordinances")
        if config.record_types:
            incomplete = [rt.name for rt in config.record_types if not rt.fees or not rt.workflow_steps]
            if incomplete:
                parts.append(f"\n**Note:** These record types may need attention: {', '.join(incomplete[:3])}")

    else:
        # General overview
        parts.append(f"**{project.customer_name} Project Overview:**\n")
        if config.record_types:
            parts.append(f"- **{len(config.record_types)} record types** configured")
        if config.departments:
            parts.append(f"- **{len(config.departments)} departments** set up")
        if config.user_roles:
            parts.append(f"- **{len(config.user_roles)} user roles** defined")
        if project.uploaded_files:
            parts.append(f"- **{len(project.uploaded_files)} data files** uploaded and analyzed")
        parts.append("\nAsk me about specific record types, fees, workflows, departments, roles, documents, or best practices for more details.")

    return "\n".join(parts) if parts else "I don't have enough information to answer that question yet. Try uploading data and running the analysis first."


# ============================================================================
# STATIC FRONTEND SERVING
# Serve the built React SPA from api/static/ (copied during deploy)
# ============================================================================

STATIC_DIR = os.path.join(os.path.dirname(__file__), "static")


@app.get("/assets/{file_path:path}")
async def serve_assets(file_path: str):
    """Serve static JS/CSS assets"""
    from fastapi.responses import FileResponse
    full_path = os.path.join(STATIC_DIR, "assets", file_path)
    if os.path.isfile(full_path):
        media_type = None
        if file_path.endswith(".js"):
            media_type = "application/javascript"
        elif file_path.endswith(".css"):
            media_type = "text/css"
        return FileResponse(full_path, media_type=media_type)
    raise HTTPException(status_code=404, detail="Asset not found")


@app.get("/{full_path:path}")
async def serve_spa(full_path: str):
    """SPA catch-all: serve index.html for all non-API routes"""
    from fastapi.responses import HTMLResponse
    index_path = os.path.join(STATIC_DIR, "index.html")
    if os.path.isfile(index_path):
        with open(index_path, "r") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>PLC AutoConfig</h1><p>Frontend not built. Run deploy.sh to build and deploy.</p>",
        status_code=200,
    )
